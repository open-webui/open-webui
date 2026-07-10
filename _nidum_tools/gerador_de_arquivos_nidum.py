"""
title: Gerador de Arquivos Nidum
author: Nidum
version: 2.2.0
description: Gera PPTX, XLSX, DOCX, PDF, HTML e APRESENTACAO HTML navegavel no servidor com alto padrao de acabamento (UX/UI) e a identidade do brandbook Nidum: paleta, fonte Maxima Nouva embutida, logos, contraste correto, layouts variados, tabelas refinadas, rodapes e numeracao. Devolve link de download nativo.
requirements: python-pptx, openpyxl, python-docx, reportlab
changelog:
  2.2.0:
    - GERACAO NAO-BLOQUEANTE: a montagem dos arquivos (pptx/xlsx/docx/pdf/html),
      que e trabalho sincrono pesado de CPU, agora roda em thread separada
      (asyncio.to_thread). Antes, um PDF grande travava o event loop inteiro do
      Open WebUI e congelava TODOS os usuarios ate terminar. O upload ao Storage
      tambem foi para thread (importante quando o Storage virar rede, ex.: R2).
    - UPLOAD TOLERANTE A R2: o Cloudflare R2 nao implementa o header
      x-amz-tagging no PutObject e derruba uploads com tags (NotImplemented).
      O upload agora tenta com tags e, se falhar, repete SEM tags - o arquivo
      salva de qualquer forma. Prepara a migracao STORAGE_PROVIDER=s3 -> R2.
    - TRACEBACK SO NO LOG: erros internos nao despejam mais o traceback
      completo no chat do usuario (vazava caminhos e estrutura interna).
      Agora: log completo no servidor (logger "gerador_nidum") + mensagem
      limpa com CODIGO DE ERRO de 8 caracteres para correlacionar no log.
    - LIMPEZA: removidas as valves OPENWEBUI_BASE_URL e OPENWEBUI_API_KEY,
      que nunca eram usadas (campo de secret morto e passivo de seguranca).
"""

# ATENCAO: este arquivo usa apenas caracteres ASCII no codigo, de proposito.
# Nao insira bullets unicode, travessoes ou emojis aqui.
# Os textos com acento vem do usuario em tempo de execucao, o que nao e problema.

import asyncio
import io
import os
import json
import logging
import uuid
import inspect

from pydantic import BaseModel

log = logging.getLogger("gerador_nidum")

# ----------------------------------------------------------------------------
# Identidade visual da Nidum (brandbook MKT) - cores em hex, sem o '#'
# ----------------------------------------------------------------------------
NIDUM_TERRACOTA = "9A4A2E"   # cor de assinatura / destaque
NIDUM_VERDE = "647260"       # verde oliva - titulos e blocos
NIDUM_AZUL = "4F7187"        # azul aco - blocos
NIDUM_CINZA = "8A8880"       # cinza quente / taupe - antetitulos
NIDUM_PRETO = "1F1E1B"       # quase preto
NIDUM_CREME = "EAE6DC"       # fundo principal dos slides de conteudo
NIDUM_CREME_ALT = "E5E0D5"   # creme alternativo
NIDUM_BRANCO = "FFFFFF"
# Tipografia da marca
NIDUM_FONT = "Maxima Nouva"  # titulos, subtitulos e corpo
NIDUM_FONT_LOGO = "Ibrand"   # exclusiva do logotipo
# Compatibilidade com versoes anteriores do gerador (nomes antigos -> paleta nova)
NIDUM_SALVIA = NIDUM_AZUL
NIDUM_OCRE = NIDUM_TERRACOTA
NIDUM_INK = NIDUM_PRETO
NIDUM_PAPER = NIDUM_CREME


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _brand_dir():
    # Resolve a pasta de assets de marca (fontes + logos), embutida na imagem.
    candidates = []
    try:
        import open_webui

        candidates.append(
            os.path.join(os.path.dirname(open_webui.__file__), "static", "brand")
        )
    except Exception:
        pass
    candidates.append("/app/backend/open_webui/static/brand")
    for d in candidates:
        if d and os.path.isdir(d):
            return d
    return None


def _logo_path(cor):
    # cor: areia | azul | escuro | terracota | verde
    d = _brand_dir()
    if not d:
        return None
    p = os.path.join(d, "logos", "nidum-" + str(cor) + ".png")
    return p if os.path.isfile(p) else None


def _font_path(fname):
    d = _brand_dir()
    if not d:
        return None
    p = os.path.join(d, "fonts", fname)
    return p if os.path.isfile(p) else None


def _font_face(fname, weight, style="normal", family="Maxima Nouva"):
    # Gera uma regra @font-face com a fonte embutida em base64 (HTML autocontido).
    import base64

    p = _font_path(fname)
    if not p:
        return ""
    try:
        b64 = base64.b64encode(open(p, "rb").read()).decode("ascii")
    except Exception:
        log.exception("gerador_nidum: falha ao ler fonte %s", fname)
        return ""
    return (
        "@font-face{font-family:'%s';font-weight:%s;font-style:%s;"
        "src:url(data:font/ttf;base64,%s) format('truetype');}"
        % (family, weight, style, b64)
    )


def _brand_css():
    # Folha de estilo da marca Nidum para HTML/PDF (fontes embutidas + paleta).
    faces = "".join(
        [
            _font_face("MaximaNouva-Regular.ttf", 400),
            _font_face("MaximaNouva-SemiBold.ttf", 600),
            _font_face("MaximaNouva-Bold.ttf", 700),
            _font_face("MaximaNouva-Italic.ttf", 400, "italic"),
            _font_face("Ibrand.ttf", "100 900", "normal", "Ibrand"),
        ]
    )
    rules = (
        ":root{--terracota:#9A4A2E;--verde:#647260;--azul:#4F7187;"
        "--cinza:#8A8880;--preto:#1F1E1B;--creme:#EAE6DC;--cremealt:#E5E0D5;}"
        "*{box-sizing:border-box}"
        "html{background:#EAE6DC}"
        "body{background:#EAE6DC;color:#1F1E1B;"
        "font-family:'Maxima Nouva',-apple-system,Segoe UI,Roboto,Arial,sans-serif;"
        "line-height:1.72;margin:0 auto;max-width:880px;padding:64px 44px 72px;"
        "font-size:18px;-webkit-font-smoothing:antialiased;}"
        "h1{font-family:'Ibrand','Maxima Nouva',sans-serif;font-size:2.7em;"
        "color:#647260;font-weight:400;line-height:1.08;"
        "letter-spacing:-.005em;margin:0 0 .5em;}"
        "h2{font-family:'Ibrand','Maxima Nouva',sans-serif;font-size:1.8em;"
        "color:#647260;font-weight:400;line-height:1.18;"
        "margin:1.9em 0 .5em;padding-bottom:.22em;"
        "border-bottom:2px solid rgba(154,74,46,.22);}"
        "h3{font-size:1.28em;color:#9A4A2E;font-weight:600;margin:1.5em 0 .4em;}"
        "h4{font-size:1em;color:#647260;text-transform:uppercase;"
        "letter-spacing:.14em;margin:1.5em 0 .4em;}"
        "p{margin:0 0 1.1em;}"
        "a{color:#9A4A2E;text-decoration:none;"
        "border-bottom:1px solid rgba(154,74,46,.4);}"
        "strong,b{color:#9A4A2E;font-weight:700;}"
        "ul,ol{margin:0 0 1.1em 1.3em;} li{margin:.45em 0;}"
        "blockquote{margin:1.7em 0;padding:.7em 1.5em;border-left:4px solid #9A4A2E;"
        "background:#E5E0D5;border-radius:0 12px 12px 0;color:#5b574f;font-style:italic;}"
        "hr{border:none;border-top:2px solid rgba(138,136,128,.32);margin:2.4em 0;}"
        "img{max-width:100%;height:auto;border-radius:14px;"
        "box-shadow:0 12px 32px rgba(31,30,27,.15);}"
        "table{border-collapse:separate;border-spacing:0;width:100%;margin:1.7em 0;"
        "border-radius:12px;overflow:hidden;box-shadow:0 6px 20px rgba(31,30,27,.08);}"
        "th{background:#647260;color:#EAE6DC;text-align:left;padding:12px 14px;"
        "font-weight:600;}"
        "td{padding:11px 14px;border-bottom:1px solid #DEDAD0;}"
        "tr:nth-child(even) td{background:rgba(255,255,255,.45);}"
        "tr:last-child td{border-bottom:none;}"
        "code{font-family:Consolas,Menlo,monospace;background:#E5E0D5;"
        "border-radius:6px;padding:.12em .4em;font-size:.92em;}"
        "pre{font-family:Consolas,Menlo,monospace;background:#E5E0D5;"
        "border-radius:10px;padding:1em 1.2em;overflow:auto;}"
        ".nidum-footer{margin-top:64px;padding-top:18px;"
        "border-top:2px solid rgba(138,136,128,.3);display:flex;align-items:center;"
        "gap:10px;color:#8A8880;font-size:.9em;letter-spacing:.02em;}"
        ".nidum-footer img{height:20px;box-shadow:none;border-radius:0;margin:0;}"
        "@media(max-width:900px){body{padding:40px 22px 60px;font-size:17px;}}"
        "@media print{body{max-width:none;padding:0;background:#fff;}"
        ".nidum-footer{display:none;}}"
    )
    return "<style>/*NIDUM_BRAND*/\n" + faces + rules + "</style>"


def _injetar_marca_html(conteudo):
    # Insere a folha de marca no <head> e um rodape sobrio antes de </body>.
    if "NIDUM_BRAND" in conteudo:
        return conteudo
    css = _brand_css()
    low = conteudo.lower()
    idx = low.find("</head>")
    if idx != -1:
        conteudo = conteudo[:idx] + css + conteudo[idx:]
    else:
        bidx = low.find("<body")
        if bidx != -1:
            end = conteudo.find(">", bidx)
            conteudo = (
                conteudo[: end + 1] + css + conteudo[end + 1 :]
                if end != -1
                else css + conteudo
            )
        else:
            conteudo = css + conteudo
    logo = _logo_b64("terracota")
    footer = (
        "<footer class='nidum-footer'>"
        + (("<img src='" + logo + "'>") if logo else "")
        + "<span>nidum &middot; fazer da casa um ninho.</span></footer>"
    )
    bidx = conteudo.lower().rfind("</body>")
    if bidx != -1:
        conteudo = conteudo[:bidx] + footer + conteudo[bidx:]
    else:
        conteudo = conteudo + footer
    return conteudo


def _esc(t):
    return (
        str(t if t is not None else "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def _logo_b64(cor):
    import base64

    p = _logo_path(cor)
    if not p:
        return ""
    try:
        return "data:image/png;base64," + base64.b64encode(
            open(p, "rb").read()
        ).decode("ascii")
    except Exception:
        log.exception("gerador_nidum: falha ao ler logo %s", cor)
        return ""


DECK_CSS = (
    "*{box-sizing:border-box;margin:0;padding:0}"
    ":root{--terracota:#9A4A2E;--verde:#647260;--azul:#4F7187;--cinza:#8A8880;"
    "--preto:#1F1E1B;--creme:#EAE6DC;--cremealt:#E5E0D5;--rn:26px;"
    "--sc:0 26px 70px rgba(31,30,27,.34)}"
    "html,body{height:100%;background:#15140f;"
    "font-family:'Maxima Nouva',-apple-system,Segoe UI,Roboto,Arial,sans-serif}"
    ".deck{position:fixed;inset:0;display:flex;align-items:center;"
    "justify-content:center;padding:3vh 3vw}"
    ".slide{position:absolute;width:min(94vw,1180px);aspect-ratio:16/9;"
    "border-radius:var(--rn);overflow:hidden;box-shadow:var(--sc);opacity:0;"
    "transform:translateY(16px) scale(.985);transition:opacity .5s ease,"
    "transform .55s cubic-bezier(.2,.7,.2,1);pointer-events:none;display:flex;"
    "flex-direction:column;justify-content:center;padding:6% 7%;background:#EAE6DC;"
    "color:#1F1E1B}"
    ".slide.active{opacity:1;transform:none;pointer-events:auto}"
    ".slide .ante{font-size:.92rem;letter-spacing:.2em;text-transform:uppercase;"
    "color:var(--cinza);margin-bottom:1rem}"
    ".slide h1{font-family:'Ibrand','Maxima Nouva',sans-serif;"
    "font-size:clamp(2.1rem,5.4vw,3.9rem);color:var(--verde);"
    "font-weight:400;line-height:1.05}"
    ".slide h2{font-family:'Ibrand','Maxima Nouva',sans-serif;"
    "font-size:clamp(1.6rem,3.8vw,2.6rem);color:var(--verde);"
    "font-weight:400;line-height:1.1}"
    ".slide p{font-size:clamp(1rem,1.7vw,1.28rem);line-height:1.62;margin-top:1rem;"
    "max-width:60ch;color:var(--preto)}"
    ".slide ul{margin:1.1rem 0 0 1.2rem} .slide li{font-size:clamp(1rem,1.6vw,1.2rem);"
    "line-height:1.5;margin:.55rem 0}"
    ".slide .logo{height:30px;position:absolute;right:6%;bottom:6.4%;opacity:.92}"
    ".slide.cover{align-items:center;text-align:center}"
    ".slide.cover .logo-c{height:clamp(54px,8vw,108px);margin-bottom:1.4rem}"
    ".slide.cover .sub{color:var(--terracota);font-size:clamp(1.1rem,2vw,1.5rem);"
    "margin-top:.7rem}"
    ".slide.fill{justify-content:center;color:var(--creme)}"
    ".slide.fill .ante{color:rgba(234,230,220,.72)}"
    ".slide.fill h1,.slide.fill h2{color:var(--creme)} .slide.fill p{color:var(--creme)}"
    ".slide.split{flex-direction:row;padding:0}"
    ".split .left{flex:0 0 42%;display:flex;align-items:center;padding:6.5%}"
    ".split .left h2{color:var(--creme)}"
    ".split .right{flex:1;background:var(--creme);display:flex;flex-direction:column;"
    "justify-content:center;padding:6.5%}"
    ".num{display:flex;gap:1.1rem;align-items:baseline;margin:.85rem 0}"
    ".num .n{font-size:clamp(1.8rem,3vw,2.7rem);color:var(--terracota);"
    "font-weight:700;line-height:1;min-width:1.5em}"
    ".num .t{font-size:clamp(1rem,1.6vw,1.22rem);line-height:1.45}"
    ".num .t b{color:var(--verde)}"
    ".cards{display:grid;grid-template-columns:1fr 1fr;gap:1.1rem;margin-top:1.3rem}"
    ".card{border-radius:18px;padding:1.3rem 1.45rem;color:var(--creme)}"
    ".card h3{font-size:1.22rem;margin-bottom:.35rem} "
    ".card p{color:rgba(234,230,220,.92);margin-top:.1rem;font-size:1.02rem;max-width:none}"
    ".nav{position:fixed;bottom:20px;left:50%;transform:translateX(-50%);display:flex;"
    "align-items:center;gap:14px;background:rgba(31,30,27,.55);"
    "-webkit-backdrop-filter:blur(7px);backdrop-filter:blur(7px);padding:8px 16px;"
    "border-radius:999px;z-index:10}"
    ".nav button{border:none;background:transparent;color:#EAE6DC;font-size:1.3rem;"
    "cursor:pointer;width:34px;height:34px;border-radius:50%;line-height:1;"
    "transition:background .15s}"
    ".nav button:hover{background:rgba(234,230,220,.18)}"
    ".dots{display:flex;gap:7px}"
    ".dot{width:8px;height:8px;border-radius:50%;background:rgba(234,230,220,.42);"
    "cursor:pointer;transition:background .2s,width .2s}"
    ".dot.on{background:#EAE6DC;width:22px;border-radius:999px}"
    ".count{color:#EAE6DC;font-size:.84rem;min-width:52px;text-align:center;"
    "letter-spacing:.04em}"
    "@media print{.nav{display:none}.slide{position:relative;opacity:1!important;"
    "transform:none!important;margin:0 auto 24px;page-break-after:always}}"
)

DECK_JS = (
    "<script>var S=[].slice.call(document.querySelectorAll('.slide')),"
    "D=[].slice.call(document.querySelectorAll('.dot')),i=0;"
    "function go(n){i=Math.max(0,Math.min(S.length-1,n));"
    "S.forEach(function(s,k){s.classList.toggle('active',k===i);});"
    "D.forEach(function(d,k){d.classList.toggle('on',k===i);});"
    "var c=document.getElementById('cnt');if(c){c.textContent=(i+1)+' / '+S.length;}}"
    "document.addEventListener('keydown',function(e){"
    "if(e.key==='ArrowRight'||e.key==='PageDown'||e.key===' '){e.preventDefault();go(i+1);}"
    "if(e.key==='ArrowLeft'||e.key==='PageUp'){e.preventDefault();go(i-1);}});"
    "go(0);</script>"
)


def _slide_html(s, tipo, mapa, cores_secao, cores_cartao, sec, logo_t, logo_a):
    titulo = _esc(s.get("titulo"))
    sub = _esc(s.get("subtitulo"))
    texto = _esc(s.get("texto"))
    logo_creme = ("<img class='logo' src='" + logo_t + "'>") if logo_t else ""
    logo_color = ("<img class='logo' src='" + logo_a + "'>") if logo_a else ""

    def bl():
        items = _lista(s.get("bullets"))
        if not items:
            return ""
        return "<ul>" + "".join("<li>" + _esc(b) + "</li>" for b in items) + "</ul>"

    ante = ("<div class='ante'>" + sub + "</div>") if sub else ""
    h1 = ("<h1>" + titulo + "</h1>") if titulo else ""
    h2 = ("<h2>" + titulo + "</h2>") if titulo else ""
    par = ("<p>" + texto + "</p>") if texto else ""

    if tipo == "capa":
        lg = ("<img class='logo-c' src='" + logo_t + "'>") if logo_t else ""
        subc = ("<div class='sub'>" + sub + "</div>") if sub else ""
        return "<section class='slide cover'>" + lg + h1 + subc + "</section>"
    if tipo == "secao":
        c = mapa.get(str(s.get("cor", "")).lower(), cores_secao[sec[0] % len(cores_secao)])
        sec[0] += 1
        return ("<section class='slide fill' style='background:" + c + "'>"
                + ante + h1 + logo_color + "</section>")
    if tipo == "destaque":
        c = mapa.get(str(s.get("cor", "terracota")).lower(), "#9A4A2E")
        return ("<section class='slide fill' style='background:" + c + "'>"
                + ante + h1 + par + logo_color + "</section>")
    if tipo == "divisao":
        c = mapa.get(str(s.get("cor", "verde")).lower(), "#647260")
        right = ante + par + bl()
        return ("<section class='slide split'><div class='left' style='background:"
                + c + "'>" + h2 + "</div><div class='right'>" + right + "</div></section>")
    if tipo == "numerada":
        itens = _lista_de_dicts(s.get("itens"))
        norm = []
        for it in itens:
            norm.append(
                (_esc(it.get("titulo")),
                 _esc(it.get("texto") or " ".join(_lista(it.get("bullets")))))
            )
        if not norm:
            norm = [("", _esc(b)) for b in _lista(s.get("bullets"))]
        rows = ""
        for k, (t_i, x_i) in enumerate(norm[:6]):
            inner = ("<b>" + t_i + "</b> " + x_i) if t_i else x_i
            rows += ("<div class='num'><div class='n'>" + str(k + 1)
                     + "</div><div class='t'>" + inner + "</div></div>")
        return "<section class='slide'>" + ante + h2 + rows + logo_creme + "</section>"
    if tipo == "cartoes":
        itens = _lista_de_dicts(s.get("itens"))
        if not itens:
            itens = [{"texto": b} for b in _lista(s.get("bullets"))]
        cards = ""
        for k, it in enumerate(itens[:4]):
            c = cores_cartao[k % len(cores_cartao)]
            ch = ("<h3>" + _esc(it.get("titulo")) + "</h3>") if it.get("titulo") else ""
            cp = ("<p>" + _esc(it.get("texto")) + "</p>") if it.get("texto") else ""
            cards += ("<div class='card' style='background:" + c + "'>" + ch + cp + "</div>")
        return ("<section class='slide'>" + ante + h2 + "<div class='cards'>"
                + cards + "</div>" + logo_creme + "</section>")
    if tipo == "encerramento":
        lg = ("<img class='logo-c' src='" + logo_t + "'>") if logo_t else ""
        subc = ("<div class='sub'>" + sub + "</div>") if sub else ""
        return "<section class='slide cover'>" + lg + h1 + subc + "</section>"
    return ("<section class='slide'>" + ante + h2 + par + bl()
            + logo_creme + "</section>")


def _get_user_id(__user__):
    if __user__ is None:
        return None
    if isinstance(__user__, dict):
        return __user__.get("id")
    return getattr(__user__, "id", None)


def _coerce(value):
    # O modelo as vezes envia listas/objetos como string JSON. Converte de volta.
    if isinstance(value, str):
        try:
            return json.loads(value)
        except Exception:
            return value
    return value


def _lista_de_dicts(value):
    # Garante uma lista de dicionarios, aceitando string JSON ou dict unico.
    value = _coerce(value)
    if isinstance(value, dict):
        value = [value]
    out = []
    for item in value or []:
        item = _coerce(item)
        if isinstance(item, dict):
            out.append(item)
    return out


def _lista(value):
    # Garante uma lista simples, aceitando string JSON.
    value = _coerce(value)
    if value is None:
        return []
    if isinstance(value, (list, tuple)):
        return list(value)
    return [value]


def _texto_para_slide(txt):
    # Converte uma string "Titulo\n- bullet\n- bullet\ntexto" em slide dict.
    linhas = [l.strip() for l in str(txt).split("\n") if l.strip()]
    titulo = ""
    bullets = []
    textos = []
    for idx, l in enumerate(linhas):
        if idx == 0 and l[:1] not in ("-", "*"):
            titulo = l
        elif l[:1] in ("-", "*"):
            bullets.append(l[1:].strip())
        else:
            textos.append(l)
    slide = {"tipo": "conteudo", "titulo": titulo}
    if bullets:
        slide["bullets"] = bullets
    if textos:
        slide["texto"] = " ".join(textos)
    return slide


def _texto_para_secao(txt):
    # Converte uma string "Heading\n- bullet\nparagrafo" em secao dict.
    linhas = [l.strip() for l in str(txt).split("\n") if l.strip()]
    heading = ""
    paragrafos = []
    bullets = []
    for idx, l in enumerate(linhas):
        if idx == 0 and l[:1] not in ("-", "*"):
            heading = l
        elif l[:1] in ("-", "*"):
            bullets.append(l[1:].strip())
        else:
            paragrafos.append(l)
    sec = {"heading": heading}
    if paragrafos:
        sec["paragrafos"] = paragrafos
    if bullets:
        sec["bullets"] = bullets
    return sec


def _itens_loose(value, conversor_string):
    # Aceita lista de dicts OU lista de strings OU dict/string unico.
    # Strings sao convertidas em dict pelo conversor informado. Isso torna a
    # tool tolerante ao jeito que o modelo costuma enviar (lista de textos).
    value = _coerce(value)
    if isinstance(value, (dict, str)):
        value = [value]
    out = []
    for item in value or []:
        item = _coerce(item)
        if isinstance(item, dict):
            out.append(item)
        elif isinstance(item, str) and item.strip():
            out.append(conversor_string(item))
    return out


def _diag_entrada_vazia(metodo, nome_param, valor_raw):
    # Em vez de gerar um arquivo vazio em silencio, explica o que chegou.
    # (Mantido no chat: o conteudo e o proprio dado enviado, nada interno vaza.)
    log.warning(
        "gerador_nidum: %s recebeu 0 itens validos em '%s' (tipo=%s)",
        metodo, nome_param, type(valor_raw).__name__,
    )
    return (
        "DIAGNOSTICO " + metodo + ": recebi 0 itens validos em '" + nome_param
        + "' apos o parse, entao o arquivo ficaria vazio. Isso costuma significar "
        "que o conteudo nao foi enviado como lista de objetos. tipo_recebido="
        + type(valor_raw).__name__ + " | inicio=" + repr(valor_raw)[:400]
    )


def _erro_limpo(metodo):
    # v2.2.0: em vez de despejar o traceback no chat (vazava caminhos internos),
    # loga o traceback completo no servidor e devolve uma mensagem limpa com um
    # codigo curto para correlacionar no log.
    codigo = uuid.uuid4().hex[:8].upper()
    log.exception("gerador_nidum: %s falhou (codigo %s)", metodo, codigo)
    return (
        "Nao consegui gerar o arquivo desta vez (erro interno em " + metodo
        + "). Codigo do erro: " + codigo
        + " - informe este codigo a Tecnologia para localizar o detalhe nos logs."
    )


async def _salvar_e_linkar(data_bytes, filename, content_type, user_id):
    """Grava os bytes pelos modulos internos do Open WebUI (sem HTTP) e
    registra o arquivo, devolvendo link de download nativo.

    v2.2.0:
    - O upload roda em thread (asyncio.to_thread): Storage.upload_file e
      sincrono e, quando o Storage for rede (ex.: Cloudflare R2), bloquearia
      o event loop durante a transferencia.
    - TOLERANCIA A TAGS: provedores S3 sem suporte a x-amz-tagging no PutObject
      (caso documentado do Cloudflare R2) derrubam o upload com NotImplemented.
      A sequencia de tentativas agora inclui variantes SEM tags, entao o
      arquivo salva de qualquer forma (a tag de user_id e conveniencia, nao
      requisito).

    Nota de compatibilidade mantida da v2.1.0: Storage.upload_file e sincrono,
    mas Files.insert_new_file pode ser assincrono - usamos inspect.isawaitable
    para funcionar nas duas versoes do Open WebUI."""
    from open_webui.storage.provider import Storage
    from open_webui.models.files import Files, FileForm

    file_id = str(uuid.uuid4())
    stored_name = file_id + "_" + filename
    tags = {"OpenWebUI-User-Id": user_id} if user_id else {}

    def _upload():
        tentativas = []
        if tags:
            tentativas.append(
                ("bytesio+tags",
                 lambda: Storage.upload_file(io.BytesIO(data_bytes), stored_name, tags))
            )
        tentativas.append(
            ("bytesio",
             lambda: Storage.upload_file(io.BytesIO(data_bytes), stored_name, {}))
        )
        if tags:
            tentativas.append(
                ("raw+tags",
                 lambda: Storage.upload_file(data_bytes, stored_name, tags))
            )
        tentativas.append(
            ("raw",
             lambda: Storage.upload_file(data_bytes, stored_name, {}))
        )
        tentativas.append(
            ("raw-sem-param-tags",
             lambda: Storage.upload_file(data_bytes, stored_name))
        )
        last_err = None
        for rotulo, tentar in tentativas:
            try:
                return tentar()
            except Exception as e:
                last_err = e
                log.warning(
                    "gerador_nidum: upload variante '%s' falhou: %s", rotulo, e
                )
        raise RuntimeError("Falha ao salvar no Storage: " + str(last_err))

    result = await asyncio.to_thread(_upload)
    if isinstance(result, tuple) and len(result) >= 2:
        file_path = result[1]
    else:
        file_path = result
    if file_path is None:
        raise RuntimeError("Storage.upload_file devolveu caminho vazio.")

    meta = {"name": filename, "content_type": content_type, "size": len(data_bytes)}
    form = FileForm(id=file_id, filename=filename, path=file_path, meta=meta, data={})

    # Registro no banco. PRECISA de await na versao async.
    inserted = Files.insert_new_file(user_id, form)
    if inspect.isawaitable(inserted):
        inserted = await inserted
    if inserted is None:
        raise RuntimeError("Falha ao registrar o arquivo no banco (insert_new_file).")

    download_path = "/api/v1/files/" + file_id + "/content"
    return "[Clique aqui para baixar " + filename + "](" + download_path + ")"


class Tools:
    class Valves(BaseModel):
        # v2.2.0: OPENWEBUI_BASE_URL e OPENWEBUI_API_KEY removidas - nunca eram
        # usadas no codigo (campo de secret morto).
        pass

    def __init__(self):
        self.valves = self.Valves()

    # ------------------------------------------------------------------
    # Metodos publicos (expostos ao modelo). Erros internos: traceback vai
    # para o LOG do servidor; o chat recebe mensagem limpa com codigo de erro.
    # ------------------------------------------------------------------
    async def gerar_pptx(
        self, titulo: str, slides: list, marca: bool = True, __user__: dict = None
    ) -> str:
        """Gera uma apresentacao PowerPoint (.pptx) e devolve um link de download.

        :param titulo: titulo geral da apresentacao.
        :param slides: lista de slides. Cada slide e um dicionario com:
            tipo ('capa'/'secao'/'conteudo'/'encerramento'), titulo,
            subtitulo (opcional), bullets (opcional, lista), texto (opcional).
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            return await self._pptx(titulo, slides, marca, __user__)
        except Exception:
            return _erro_limpo("gerar_pptx")

    async def gerar_xlsx(
        self, titulo: str, planilhas: list, marca: bool = True, __user__: dict = None
    ) -> str:
        """Gera uma planilha Excel (.xlsx) e devolve um link de download.

        :param titulo: titulo geral.
        :param planilhas: lista de abas. Cada aba e um dicionario com:
            nome, cabecalhos (lista), linhas (lista de listas).
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            return await self._xlsx(titulo, planilhas, marca, __user__)
        except Exception:
            return _erro_limpo("gerar_xlsx")

    async def gerar_docx(
        self, titulo: str, secoes: list, marca: bool = True, __user__: dict = None
    ) -> str:
        """Gera um documento Word (.docx) e devolve um link de download.

        :param titulo: titulo do documento.
        :param secoes: lista de secoes. Cada secao e um dicionario com:
            heading, paragrafos (opcional, lista), bullets (opcional, lista).
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            return await self._docx(titulo, secoes, marca, __user__)
        except Exception:
            return _erro_limpo("gerar_docx")

    async def gerar_pdf(
        self, titulo: str, secoes: list, marca: bool = True, __user__: dict = None
    ) -> str:
        """Gera um documento PDF e devolve um link de download.

        :param titulo: titulo do documento.
        :param secoes: lista de secoes. Cada secao e um dicionario com:
            heading, paragrafos (opcional, lista), bullets (opcional, lista),
            tabela (opcional, lista de listas).
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            return await self._pdf(titulo, secoes, marca, __user__)
        except Exception:
            return _erro_limpo("gerar_pdf")

    async def gerar_html(
        self, titulo: str, html: str, __user__: dict = None
    ) -> str:
        """Gera um arquivo HTML (.html) e devolve um link de download.

        :param titulo: titulo/nome do arquivo.
        :param html: documento HTML completo (string), pronto para abrir no navegador.
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            conteudo = _coerce(html)
            conteudo = conteudo if isinstance(conteudo, str) else str(conteudo or "")
            conteudo = conteudo.strip()
            if not conteudo:
                return _diag_entrada_vazia("gerar_html", "html", html)

            def _montar():
                c = conteudo
                # Se vier so um fragmento, embrulha num documento HTML minimo.
                if "<html" not in c.lower():
                    c = (
                        '<!DOCTYPE html>\n<html lang="pt-br"><head><meta charset="utf-8">'
                        '<meta name="viewport" content="width=device-width, initial-scale=1">'
                        "<title>" + (titulo or "Documento") + "</title></head><body>\n"
                        + c + "\n</body></html>"
                    )
                # _injetar_marca_html le fontes/logo do disco (IO) - por isso
                # roda aqui dentro da thread.
                return _injetar_marca_html(c).encode("utf-8")

            data = await asyncio.to_thread(_montar)
            nome = (titulo or "pagina").strip().replace(" ", "_") + ".html"
            link = await _salvar_e_linkar(
                data, nome, "text/html", _get_user_id(__user__)
            )
            return "Arquivo gerado com sucesso. Link para download: " + link
        except Exception:
            return _erro_limpo("gerar_html")

    async def gerar_apresentacao_html(
        self, titulo: str, slides: list, __user__: dict = None
    ) -> str:
        """Gera uma APRESENTACAO em HTML navegavel (deck) com a identidade Nidum.

        Deck autocontido: 1 slide por vez, navegacao por setas/teclado/dots,
        cantos arredondados, transicoes, contraste correto e fonte embutida.
        :param slides: mesma estrutura do gerar_pptx (lista de slides com tipo,
            titulo, subtitulo, texto, bullets, cor, itens).
        :return: link /api/v1/files/{id}/content para baixar o arquivo.
        """
        try:
            return await self._apresentacao_html(titulo, slides, __user__)
        except Exception:
            return _erro_limpo("gerar_apresentacao_html")

    # ------------------------------------------------------------------
    # Implementacoes (privadas - nao expostas ao modelo).
    # v2.2.0: a montagem pesada roda em asyncio.to_thread para NAO travar o
    # event loop do Open WebUI (um render grande congelava todos os usuarios).
    # ------------------------------------------------------------------
    async def _apresentacao_html(self, titulo, slides, __user__):
        raw = slides
        slides = _itens_loose(slides, _texto_para_slide)
        if not slides:
            return _diag_entrada_vazia("gerar_apresentacao_html", "slides", raw)

        def _montar():
            faces = "".join(
                [
                    _font_face("MaximaNouva-Thin.ttf", 300),
                    _font_face("MaximaNouva-Regular.ttf", 400),
                    _font_face("MaximaNouva-SemiBold.ttf", 600),
                    _font_face("MaximaNouva-Bold.ttf", 700),
                    _font_face("MaximaNouva-Italic.ttf", 400, "italic"),
                    _font_face("Ibrand.ttf", "100 900", "normal", "Ibrand"),
                ]
            )
            logo_t = _logo_b64("terracota")
            logo_a = _logo_b64("areia")
            mapa = {
                "verde": "#647260", "azul": "#4F7187", "terracota": "#9A4A2E",
                "preto": "#1F1E1B", "creme": "#EAE6DC",
            }
            cores_secao = ["#647260", "#4F7187", "#9A4A2E"]
            cores_cartao = ["#4F7187", "#647260", "#9A4A2E", "#1F1E1B"]
            sec = [0]

            partes = []
            for s in slides:
                tipo = (s.get("tipo") or "conteudo").lower()
                partes.append(
                    _slide_html(
                        s, tipo, mapa, cores_secao, cores_cartao, sec, logo_t, logo_a
                    )
                )
            deck = "".join(partes)
            dots = "".join(
                "<span class='dot' onclick='go(" + str(k) + ")'></span>"
                for k in range(len(slides))
            )
            nav = (
                "<div class='nav'><button onclick='go(i-1)'>&#8249;</button>"
                "<div class='dots'>" + dots + "</div>"
                "<button onclick='go(i+1)'>&#8250;</button>"
                "<span class='count' id='cnt'></span></div>"
            )
            html = (
                "<!DOCTYPE html><html lang=\"pt-br\"><head><meta charset=\"utf-8\">"
                "<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">"
                "<title>" + _esc(titulo or "Apresentacao Nidum") + "</title>"
                "<style>" + faces + DECK_CSS + "</style></head><body>"
                "<div class='deck'>" + deck + "</div>" + nav + DECK_JS + "</body></html>"
            )
            return html.encode("utf-8")

        data = await asyncio.to_thread(_montar)
        nome = (titulo or "apresentacao").strip().replace(" ", "_") + ".html"
        link = await _salvar_e_linkar(
            data, nome, "text/html", _get_user_id(__user__)
        )
        return "Arquivo gerado com sucesso. Link para download: " + link

    async def _pptx(self, titulo, slides, marca, __user__):
        raw_slides = slides
        slides = _itens_loose(slides, _texto_para_slide)
        if not slides:
            return _diag_entrada_vazia("gerar_pptx", "slides", raw_slides)

        def _montar():
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            def col(h):
                return RGBColor(*_hex_to_rgb(h))

            terracota = col(NIDUM_TERRACOTA)
            verde = col(NIDUM_VERDE)
            azul = col(NIDUM_AZUL)
            cinza = col(NIDUM_CINZA)
            preto = col(NIDUM_PRETO)
            creme = col(NIDUM_CREME)
            branco = col(NIDUM_BRANCO)
            cores_secao = [verde, azul, terracota]
            cores_cartao = [azul, verde, terracota, preto]
            mapa_cor = {"verde": verde, "azul": azul, "terracota": terracota,
                        "preto": preto, "creme": creme}
            sec_idx = 0

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]
            SW = prs.slide_width
            SH = prs.slide_height

            def add_fundo(slide, cor):
                shp = slide.shapes.add_shape(1, 0, 0, SW, SH)
                shp.fill.solid()
                shp.fill.fore_color.rgb = cor
                shp.line.fill.background()
                shp.shadow.inherit = False
                slide.shapes._spTree.remove(shp._element)
                slide.shapes._spTree.insert(2, shp._element)
                return shp

            def add_caixa(slide, left, top, width, height, anchor=None):
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.word_wrap = True
                if anchor is not None:
                    tf.vertical_anchor = anchor
                return tf

            def estilo(p, size, cor, bold=False, upper=False, italic=False, align=None):
                if upper and p.text:
                    p.text = p.text.upper()
                # Titulos (bold) usam a Ibrand (com peso proprio); corpo usa Maxima.
                fonte = NIDUM_FONT_LOGO if bold else NIDUM_FONT
                negr = bold and fonte == NIDUM_FONT
                p.font.name = fonte
                p.font.size = Pt(size)
                p.font.bold = negr
                p.font.italic = italic
                p.font.color.rgb = cor
                if align is not None:
                    p.alignment = align
                for r in p.runs:
                    r.font.name = fonte
                    r.font.size = Pt(size)
                    r.font.bold = negr
                    r.font.italic = italic
                    r.font.color.rgb = cor

            def add_logo(slide, cor_logo, left, top, width):
                p = _logo_path(cor_logo)
                if not p:
                    return
                try:
                    slide.shapes.add_picture(p, left, top, width=width)
                except Exception:
                    log.exception("gerador_nidum: falha ao inserir logo no pptx")

            def add_bloco(slide, left, top, w, h, cor):
                sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
                sh.fill.solid()
                sh.fill.fore_color.rgb = cor
                sh.line.fill.background()
                sh.shadow.inherit = False
                return sh

            def add_card(slide, left, top, w, h, cor, titulo_c, texto_c):
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
                )
                card.fill.solid()
                card.fill.fore_color.rgb = cor
                card.line.fill.background()
                card.shadow.inherit = False
                tf = card.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.22)
                tf.margin_right = Inches(0.22)
                tf.margin_top = Inches(0.18)
                p = tf.paragraphs[0]
                p.text = titulo_c or ""
                estilo(p, 15, creme, bold=True)
                if texto_c:
                    p2 = tf.add_paragraph()
                    p2.text = texto_c
                    estilo(p2, 11, creme)
                return card

            def cor_de(nome, padrao):
                return mapa_cor.get(str(nome).lower(), padrao) if nome else padrao

            for s in slides:
                tipo = (s.get("tipo") or "conteudo").lower()
                slide = prs.slides.add_slide(blank)
                cor_titulo = verde if marca else preto
                cor_corpo = preto

                if tipo == "capa":
                    if marca:
                        add_fundo(slide, creme)
                        add_logo(slide, "terracota", Inches(3.67), Inches(1.3), Inches(6.0))
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(4.4), Inches(11.33), Inches(2.0),
                        MSO_ANCHOR.TOP,
                    )
                    p = tf.paragraphs[0]
                    p.text = s.get("titulo") or titulo or ""
                    estilo(p, 40, cor_titulo, bold=True, align=PP_ALIGN.CENTER)
                    if s.get("subtitulo"):
                        p2 = tf.add_paragraph()
                        p2.text = s["subtitulo"]
                        estilo(p2, 18, terracota if marca else preto, align=PP_ALIGN.CENTER)

                elif tipo == "secao":
                    cor_sec = cores_secao[sec_idx % len(cores_secao)]
                    sec_idx += 1
                    cor_sec = cor_de(s.get("cor"), cor_sec)
                    if marca:
                        add_fundo(slide, cor_sec)
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.3),
                        MSO_ANCHOR.MIDDLE,
                    )
                    if s.get("subtitulo"):
                        pa = tf.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 16, creme if marca else preto, upper=True)
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = s.get("titulo", "")
                    estilo(p, 40, creme if marca else preto, bold=True)
                    if marca:
                        add_logo(slide, "areia", Inches(0.7), Inches(6.6), Inches(1.8))

                elif tipo == "destaque":
                    cor_d = cor_de(s.get("cor"), terracota)
                    if marca:
                        add_fundo(slide, cor_d)
                    tf = add_caixa(
                        slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(3.1),
                        MSO_ANCHOR.MIDDLE,
                    )
                    if s.get("subtitulo"):
                        pa = tf.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 16, creme if marca else preto, upper=True)
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = s.get("titulo", "")
                    estilo(p, 34, creme if marca else preto, bold=True)
                    if s.get("texto"):
                        p2 = tf.add_paragraph()
                        p2.text = s["texto"]
                        estilo(p2, 16, creme if marca else preto)
                    if marca:
                        add_logo(slide, "areia", Inches(11.0), Inches(6.8), Inches(1.6))

                elif tipo == "divisao":
                    cor_b = cor_de(s.get("cor"), verde)
                    if marca:
                        add_fundo(slide, creme)
                        add_bloco(slide, 0, 0, Inches(5.4), SH, cor_b)
                    tf_l = add_caixa(
                        slide, Inches(0.7), Inches(0.9), Inches(4.2), Inches(5.6),
                        MSO_ANCHOR.MIDDLE,
                    )
                    pl = tf_l.paragraphs[0]
                    pl.text = s.get("titulo", "")
                    estilo(pl, 28, creme if marca else preto, bold=True)
                    tf_r = add_caixa(slide, Inches(5.9), Inches(0.9), Inches(6.7), Inches(5.6))
                    first = True
                    if s.get("subtitulo"):
                        pa = tf_r.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 13, cinza, upper=True)
                        first = False
                    if s.get("texto"):
                        p = tf_r.paragraphs[0] if first else tf_r.add_paragraph()
                        first = False
                        p.text = s["texto"]
                        estilo(p, 15, preto)
                    for b in _lista(s.get("bullets")):
                        p = tf_r.paragraphs[0] if first else tf_r.add_paragraph()
                        first = False
                        p.text = str(b)
                        estilo(p, 15, preto)

                elif tipo == "numerada":
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.3))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, cinza, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 28, verde if marca else preto, bold=True)
                    itens = _lista_de_dicts(s.get("itens"))
                    norm = []
                    for it in itens:
                        norm.append(
                            (it.get("titulo") or "",
                             it.get("texto") or " ".join(_lista(it.get("bullets"))))
                        )
                    if not norm:
                        norm = [("", str(b)) for b in _lista(s.get("bullets"))]
                    for idx, (t_i, txt_i) in enumerate(norm[:6]):
                        y = Inches(2.0 + idx * 0.82)
                        tfn = add_caixa(slide, Inches(0.9), y, Inches(0.9), Inches(0.8))
                        pn = tfn.paragraphs[0]
                        pn.text = str(idx + 1)
                        estilo(pn, 30, terracota, bold=True)
                        tfx = add_caixa(slide, Inches(1.9), y, Inches(10.5), Inches(0.8),
                                        MSO_ANCHOR.MIDDLE)
                        px = tfx.paragraphs[0]
                        px.text = (t_i + " - " + txt_i).strip(" -") if t_i else txt_i
                        estilo(px, 15, preto)

                elif tipo == "cartoes":
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.1))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, cinza, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 26, verde if marca else preto, bold=True)
                    itens = _lista_de_dicts(s.get("itens"))
                    if not itens:
                        itens = [{"titulo": "", "texto": str(b)} for b in _lista(s.get("bullets"))]
                    itens = itens[:4]
                    xs = [Inches(0.9), Inches(6.85)]
                    ys = [Inches(1.9), Inches(4.35)]
                    for idx, it in enumerate(itens):
                        add_card(
                            slide, xs[idx % 2], ys[idx // 2], Inches(5.6), Inches(2.2),
                            cores_cartao[idx % len(cores_cartao)],
                            it.get("titulo") or "", it.get("texto") or "",
                        )

                elif tipo == "encerramento":
                    if marca:
                        add_fundo(slide, creme)
                        add_logo(slide, "terracota", Inches(4.67), Inches(2.4), Inches(4.0))
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(4.6), Inches(11.33), Inches(1.2),
                        MSO_ANCHOR.TOP,
                    )
                    p = tf.paragraphs[0]
                    p.text = s.get("titulo") or "Fazer da casa um ninho."
                    estilo(p, 22, cor_titulo, align=PP_ALIGN.CENTER)
                    if s.get("texto"):
                        p2 = tf.add_paragraph()
                        p2.text = s["texto"]
                        estilo(p2, 14, cinza if marca else preto, align=PP_ALIGN.CENTER)

                else:
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.4))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, cinza if marca else preto, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 30, cor_titulo, bold=True)

                    tf_b = add_caixa(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.5))
                    first = True
                    if s.get("texto"):
                        p = tf_b.paragraphs[0]
                        p.text = s["texto"]
                        estilo(p, 16, cor_corpo)
                        first = False
                    for b in _lista(s.get("bullets")):
                        p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
                        first = False
                        p.text = str(b)
                        estilo(p, 16, cor_corpo)
                    if marca:
                        add_logo(slide, "terracota", Inches(11.1), Inches(6.8), Inches(1.6))

            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = (titulo or "apresentacao").strip().replace(" ", "_") + ".pptx"
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. Link para download: " + link

    async def _xlsx(self, titulo, planilhas, marca, __user__):
        raw_planilhas = planilhas
        planilhas = _lista_de_dicts(planilhas)
        if not planilhas:
            return _diag_entrada_vazia("gerar_xlsx", "planilhas", raw_planilhas)

        def _montar():
            from openpyxl import Workbook
            from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            wb.remove(wb.active)

            f_titulo = Font(name=NIDUM_FONT_LOGO, size=18, bold=False, color=NIDUM_VERDE)
            f_header = Font(name=NIDUM_FONT, size=11, bold=True, color=NIDUM_CREME)
            f_body = Font(name=NIDUM_FONT, size=11, color=NIDUM_PRETO)
            fill_header = PatternFill("solid", fgColor=NIDUM_VERDE)
            fill_alt = PatternFill("solid", fgColor=NIDUM_CREME_ALT)
            linha_side = Side(style="thin", color="DEDAD0")
            borda_b = Border(bottom=linha_side)
            al_h = Alignment(horizontal="left", vertical="center")
            al_b = Alignment(horizontal="left", vertical="center", wrap_text=True)

            for pl in planilhas:
                ws = wb.create_sheet(title=(pl.get("nome") or "Planilha")[:31])
                if marca:
                    ws.sheet_properties.tabColor = NIDUM_VERDE
                cabec = _lista(pl.get("cabecalhos"))
                linhas = [_lista(r) for r in _lista(pl.get("linhas"))]
                ncols = max([len(cabec)] + [len(r) for r in linhas] + [1])

                r = 1
                if marca:
                    ws.merge_cells(
                        start_row=1, start_column=1, end_row=1, end_column=ncols
                    )
                    ct = ws.cell(row=1, column=1, value=(pl.get("nome") or titulo or "Planilha"))
                    ct.font = f_titulo
                    ct.alignment = al_h
                    ws.row_dimensions[1].height = 28
                    r = 2

                header_row = None
                if cabec:
                    header_row = r
                    for j, val in enumerate(cabec, start=1):
                        cell = ws.cell(row=r, column=j, value=val)
                        cell.alignment = al_h
                        if marca:
                            cell.fill = fill_header
                            cell.font = f_header
                    ws.row_dimensions[r].height = 22
                    r += 1

                data_ini = r
                for ld in linhas:
                    for j in range(1, ncols + 1):
                        val = ld[j - 1] if j - 1 < len(ld) else None
                        cell = ws.cell(row=r, column=j, value=val)
                        cell.font = f_body
                        cell.alignment = al_b
                        cell.border = borda_b
                        if marca and ((r - data_ini) % 2 == 1):
                            cell.fill = fill_alt
                    r += 1

                for j in range(1, ncols + 1):
                    maxlen = len(str(cabec[j - 1])) if cabec and j - 1 < len(cabec) else 0
                    for ld in linhas:
                        if j - 1 < len(ld) and ld[j - 1] is not None:
                            maxlen = max(maxlen, len(str(ld[j - 1])))
                    ws.column_dimensions[get_column_letter(j)].width = min(
                        max(maxlen + 3, 12), 44
                    )

                if header_row:
                    ws.freeze_panes = ws.cell(row=header_row + 1, column=1).coordinate

            buf = io.BytesIO()
            wb.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = (titulo or "planilha").strip().replace(" ", "_") + ".xlsx"
        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. Link para download: " + link

    async def _docx(self, titulo, secoes, marca, __user__):
        raw_secoes = secoes
        secoes = _itens_loose(secoes, _texto_para_secao)
        if not secoes:
            return _diag_entrada_vazia("gerar_docx", "secoes", raw_secoes)

        def _montar():
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            verde = RGBColor(*_hex_to_rgb(NIDUM_VERDE))
            terracota = RGBColor(*_hex_to_rgb(NIDUM_TERRACOTA))
            cinza = RGBColor(*_hex_to_rgb(NIDUM_CINZA))
            ink = RGBColor(*_hex_to_rgb(NIDUM_PRETO))

            doc = Document()
            if marca:
                try:
                    normal = doc.styles["Normal"]
                    normal.font.name = NIDUM_FONT
                    normal.font.size = Pt(11)
                    normal.font.color.rgb = ink
                    pf = normal.paragraph_format
                    pf.line_spacing = 1.32
                    pf.space_after = Pt(8)
                    for lvl, cor, sz in (
                        ("Title", verde, 30),
                        ("Heading 1", verde, 17),
                        ("Heading 2", terracota, 13),
                    ):
                        try:
                            st = doc.styles[lvl]
                            st.font.name = NIDUM_FONT_LOGO
                            st.font.color.rgb = cor
                            st.font.size = Pt(sz)
                        except Exception:
                            pass
                except Exception:
                    log.exception("gerador_nidum: falha ao aplicar estilos docx")

            h = doc.add_heading(titulo or "Documento", level=0)
            for run in h.runs:
                run.font.name = NIDUM_FONT_LOGO
                if marca:
                    run.font.color.rgb = verde

            for sec in secoes:
                if sec.get("heading"):
                    hs = doc.add_heading(sec["heading"], level=1)
                    for run in hs.runs:
                        run.font.name = NIDUM_FONT_LOGO
                        if marca:
                            run.font.color.rgb = verde
                for par in _lista(sec.get("paragrafos")):
                    p = doc.add_paragraph(str(par))
                    for run in p.runs:
                        run.font.name = NIDUM_FONT
                        run.font.size = Pt(11)
                        run.font.color.rgb = ink
                for b in _lista(sec.get("bullets")):
                    pb = doc.add_paragraph(str(b), style="List Bullet")
                    for run in pb.runs:
                        run.font.name = NIDUM_FONT
                        run.font.color.rgb = ink

            if marca:
                try:
                    fp = doc.sections[0].footer.paragraphs[0]
                    fr = fp.add_run("nidum  -  fazer da casa um ninho.")
                    fr.font.name = NIDUM_FONT
                    fr.font.size = Pt(8)
                    fr.font.color.rgb = cinza
                    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception:
                    log.exception("gerador_nidum: falha ao aplicar rodape docx")

            buf = io.BytesIO()
            doc.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = (titulo or "documento").strip().replace(" ", "_") + ".docx"
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. Link para download: " + link

    async def _pdf(self, titulo, secoes, marca, __user__):
        raw_secoes = secoes
        secoes = _itens_loose(secoes, _texto_para_secao)
        if not secoes:
            return _diag_entrada_vazia("gerar_pdf", "secoes", raw_secoes)

        def _montar():
            from reportlab.lib.pagesizes import A4
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.platypus import (
                SimpleDocTemplate,
                Paragraph,
                Spacer,
                Table,
                TableStyle,
                ListFlowable,
                ListItem,
                Image,
                HRFlowable,
            )
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            verde = colors.HexColor("#" + NIDUM_VERDE)
            terracota = colors.HexColor("#" + NIDUM_TERRACOTA)
            preto = colors.HexColor("#" + NIDUM_PRETO)
            creme = colors.HexColor("#" + NIDUM_CREME)
            cremealt = colors.HexColor("#" + NIDUM_CREME_ALT)
            cinza = colors.HexColor("#" + NIDUM_CINZA)
            branco = colors.HexColor("#" + NIDUM_BRANCO)
            ink = preto
            linha_cor = colors.HexColor("#DEDAD0")

            # Registra a fonte da marca (fallback Helvetica se nao encontrar)
            FONT = "Helvetica"
            FONT_B = "Helvetica-Bold"
            FONT_TITULO = "Helvetica-Bold"
            try:
                reg = _font_path("MaximaNouva-Regular.ttf")
                bold = _font_path("MaximaNouva-Bold.ttf")
                ib = _font_path("Ibrand.ttf")
                if reg:
                    pdfmetrics.registerFont(TTFont("MaximaNouva", reg))
                    FONT = "MaximaNouva"
                if bold:
                    pdfmetrics.registerFont(TTFont("MaximaNouva-Bold", bold))
                    FONT_B = "MaximaNouva-Bold"
                if reg and bold:
                    pdfmetrics.registerFontFamily(
                        "MaximaNouva", normal="MaximaNouva", bold="MaximaNouva-Bold"
                    )
                if ib:
                    pdfmetrics.registerFont(TTFont("Ibrand", ib))
                    FONT_TITULO = "Ibrand"
                else:
                    FONT_TITULO = FONT_B
            except Exception:
                log.exception("gerador_nidum: falha ao registrar fontes no pdf")
                FONT, FONT_B, FONT_TITULO = "Helvetica", "Helvetica-Bold", "Helvetica-Bold"

            styles = getSampleStyleSheet()
            st_titulo = ParagraphStyle(
                "NidumTitulo", parent=styles["Title"],
                textColor=verde if marca else ink, fontName=FONT_TITULO,
                fontSize=25, leading=29,
            )
            st_head = ParagraphStyle(
                "NidumHead", parent=styles["Heading1"],
                textColor=verde if marca else ink, fontName=FONT_TITULO,
                fontSize=15, leading=19,
            )
            st_body = ParagraphStyle(
                "NidumBody",
                parent=styles["BodyText"],
                textColor=ink,
                fontName=FONT,
                fontSize=11,
                leading=16,
            )

            buf = io.BytesIO()
            docp = SimpleDocTemplate(
                buf,
                pagesize=A4,
                leftMargin=22 * mm,
                rightMargin=22 * mm,
                topMargin=20 * mm,
                bottomMargin=18 * mm,
            )

            def _fundo(canvas, doc):
                if not marca:
                    return
                canvas.saveState()
                canvas.setFillColor(creme)
                canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=1, stroke=0)
                # rodape: assinatura + numero de pagina
                canvas.setFillColor(cinza)
                canvas.setFont(FONT, 8)
                canvas.drawString(
                    22 * mm, 11 * mm, "nidum  -  fazer da casa um ninho."
                )
                canvas.drawRightString(
                    doc.pagesize[0] - 22 * mm, 11 * mm, str(doc.page)
                )
                canvas.restoreState()

            flow = []
            if marca:
                lp = _logo_path("terracota")
                if lp:
                    try:
                        flow.append(
                            Image(lp, width=40 * mm, height=22.5 * mm, hAlign="LEFT")
                        )
                        flow.append(Spacer(1, 2 * mm))
                    except Exception:
                        log.exception("gerador_nidum: falha ao inserir logo no pdf")
            flow.append(Paragraph(titulo or "Documento", st_titulo))
            flow.append(Spacer(1, 2 * mm))
            if marca:
                flow.append(
                    HRFlowable(
                        width="100%", thickness=1.5, color=terracota,
                        spaceBefore=1, spaceAfter=7 * mm, lineCap="round",
                    )
                )
            else:
                flow.append(Spacer(1, 6 * mm))

            for sec in secoes:
                if sec.get("heading"):
                    flow.append(Paragraph(sec["heading"], st_head))
                    flow.append(Spacer(1, 2 * mm))
                for par in _lista(sec.get("paragrafos")):
                    flow.append(Paragraph(str(par), st_body))
                    flow.append(Spacer(1, 2 * mm))
                bullets = _lista(sec.get("bullets"))
                if bullets:
                    items = [ListItem(Paragraph(str(b), st_body)) for b in bullets]
                    flow.append(ListFlowable(items, bulletType="bullet"))
                    flow.append(Spacer(1, 2 * mm))
                tabela = sec.get("tabela")
                if tabela:
                    tabela = [_lista(r) for r in _lista(tabela)]
                    t = Table(tabela, hAlign="LEFT")
                    estilo_t = [
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("FONTSIZE", (0, 0), (-1, -1), 10),
                        ("TEXTCOLOR", (0, 1), (-1, -1), ink),
                        ("LEFTPADDING", (0, 0), (-1, -1), 9),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 9),
                        ("TOPPADDING", (0, 0), (-1, -1), 7),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                        ("LINEBELOW", (0, 1), (-1, -1), 0.4, linha_cor),
                    ]
                    if marca:
                        estilo_t.append(("BACKGROUND", (0, 0), (-1, 0), verde))
                        estilo_t.append(("TEXTCOLOR", (0, 0), (-1, 0), creme))
                        estilo_t.append(("FONTNAME", (0, 0), (-1, 0), FONT_B))
                        estilo_t.append(("FONTNAME", (0, 1), (-1, -1), FONT))
                        estilo_t.append(("TOPPADDING", (0, 0), (-1, 0), 9))
                        estilo_t.append(("BOTTOMPADDING", (0, 0), (-1, 0), 9))
                        estilo_t.append(
                            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [branco, cremealt])
                        )
                    t.setStyle(TableStyle(estilo_t))
                    flow.append(t)
                    flow.append(Spacer(1, 4 * mm))

            docp.build(flow, onFirstPage=_fundo, onLaterPages=_fundo)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = (titulo or "documento").strip().replace(" ", "_") + ".pdf"
        link = await _salvar_e_linkar(
            data, nome, "application/pdf", _get_user_id(__user__)
        )
        return "Arquivo gerado com sucesso. Link para download: " + link
