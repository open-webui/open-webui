import asyncio
import json
import logging
import re
import sys
import time

import ftfy
import requests
from azure.identity import DefaultAzureCredential
from langchain_community.document_loaders import (
    AzureAIDocumentIntelligenceLoader,
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    YoutubeLoader,
)
from langchain_core.documents import Document
from open_webui.env import (
    AIOHTTP_CLIENT_SESSION_SSL,
    GLOBAL_LOG_LEVEL,
    MINERU_MAX_MARKDOWN_BYTES,
    REQUESTS_VERIFY,
)
from open_webui.retrieval.loaders.datalab_marker import DatalabMarkerLoader
from open_webui.retrieval.loaders.external_document import ExternalDocumentLoader
from open_webui.retrieval.loaders.mineru import MinerULoader
from open_webui.retrieval.loaders.mistral import MistralLoader
from open_webui.retrieval.loaders.paddleocr_vl import PaddleOCRVLLoader

logging.basicConfig(stream=sys.stdout, level=GLOBAL_LOG_LEVEL)
log = logging.getLogger(__name__)

known_source_ext = [
    'go',
    'py',
    'java',
    'sh',
    'bat',
    'ps1',
    'cmd',
    'js',
    'ts',
    'css',
    'cpp',
    'hpp',
    'h',
    'c',
    'cs',
    'sql',
    'log',
    'ini',
    'pl',
    'pm',
    'r',
    'dart',
    'dockerfile',
    'env',
    'php',
    'hs',
    'hsc',
    'lua',
    'nginxconf',
    'conf',
    'm',
    'mm',
    'plsql',
    'perl',
    'rb',
    'rs',
    'db2',
    'scala',
    'bash',
    'swift',
    'vue',
    'svelte',
    'ex',
    'exs',
    'erl',
    'tsx',
    'jsx',
    'hs',
    'lhs',
    'json',
    'yaml',
    'yml',
    'toml',
]


class ExcelLoader:
    """Fallback Excel loader using pandas when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import pandas as pd

        text_parts = []
        xls = pd.ExcelFile(self.file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text_parts.append(f'Sheet: {sheet_name}\n{df.to_string(index=False)}')
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class PptxLoader:
    """Fallback PowerPoint loader using python-pptx when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                text_parts.append(f'Slide {i}:\n' + '\n'.join(slide_texts))
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class TikaLoader:
    def __init__(self, url, file_path, mime_type=None, extract_images=None):
        self.url = url
        self.file_path = file_path
        self.mime_type = mime_type

        self.extract_images = extract_images

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            data = f.read()

        if self.mime_type is not None:
            headers = {'Content-Type': self.mime_type}
        else:
            headers = {}

        if self.extract_images == True:
            headers['X-Tika-PDFextractInlineImages'] = 'true'

        endpoint = self.url
        if not endpoint.endswith('/'):
            endpoint += '/'
        endpoint += 'tika/text'

        r = requests.put(endpoint, data=data, headers=headers, verify=REQUESTS_VERIFY)

        if r.ok:
            raw_metadata = r.json()
            text = raw_metadata.get('X-TIKA:content', '<No text content found>').strip()

            if 'Content-Type' in raw_metadata:
                headers['Content-Type'] = raw_metadata['Content-Type']

            log.debug('Tika extracted text: %s', text)

            return [Document(page_content=text, metadata=headers)]
        else:
            raise Exception(f'Error calling Tika: {r.reason}')


class DoclingLoader:
    def __init__(self, url, api_key=None, file_path=None, mime_type=None, params=None, timeout=None, status_callback=None):
        self.url = url.rstrip('/')
        self.api_key = api_key
        self.file_path = file_path
        self.mime_type = mime_type
        self.params = params or {}
        self.timeout = timeout  # total seconds to wait; None = infinite
        self.status_callback = status_callback  # optional callable(dict) to persist queue state

    def _build_headers(self) -> dict:
        headers = {}
        if self.api_key:
            headers["X-Api-Key"] = f"{self.api_key}"
        return headers

    def _submit_file(self) -> tuple[str, int | None]:
        headers = self._build_headers()
        page_break_marker = '\f'
        with open(self.file_path, "rb") as f:
            r = requests.post(
                f"{self.url}/v1/convert/file/async",
                files={
                    'files': (
                        self.file_path,
                        f,
                        self.mime_type or 'application/octet-stream',
                    )
                },
                data={
                    'image_export_mode': 'placeholder',
                    'md_page_break_placeholder': page_break_marker,
                    **self.params,
                },
                headers=headers,
                verify=AIOHTTP_CLIENT_SESSION_SSL,
                timeout=30,
            )

        if not r.ok:
            error_msg = f"Error calling Docling API: {r.reason}"
            if r.text:
                try:
                    error_data = r.json()
                    if 'detail' in error_data:
                        error_msg += f' - {error_data["detail"]}'
                except Exception:
                    error_msg += f' - {r.text}'
            raise Exception(f'Error calling Docling: {error_msg}')

        submit_data = r.json()
        task_id = submit_data.get("task_id")
        task_position = submit_data.get("task_position")
        if not task_id:
            raise Exception("Docling async submit did not return a task_id")
        log.info(
            "Docling task submitted: %s, queue position: %s",
            task_id,
            task_position,
        )
        if self.status_callback:
            self.status_callback({"task_id": task_id, "task_position": task_position})

        return task_id, task_position

    def _poll_task_until_done(self, task_id: str) -> dict:
        headers = self._build_headers()
        deadline = time.monotonic() + self.timeout if self.timeout is not None else None
        poll_wait = 30  # long-poll window per request (seconds)

        while True:
            if deadline is not None:
                remaining = deadline - time.monotonic()
                if remaining <= 0:
                    raise Exception(
                        f"Docling conversion timed out after {self.timeout}s (task_id={task_id})"
                    )
                poll_wait = min(poll_wait, int(remaining) + 1)

            poll_start = time.monotonic()
            try:
                status_r = requests.get(
                    f"{self.url}/v1/status/poll/{task_id}",
                    params={"wait": poll_wait},
                    headers=headers,
                    timeout=poll_wait + 10,
                )
            except requests.Timeout:
                log.warning("Docling status poll timed out for task %s, retrying", task_id)
                elapsed = time.monotonic() - poll_start
                if elapsed < poll_wait:
                    time.sleep(poll_wait - elapsed)
                continue

            if not status_r.ok:
                raise Exception(f"Error polling Docling task status: {status_r.reason}")

            status_data = status_r.json()
            task_status = status_data.get("task_status", "")
            log.debug(
                "Docling task %s: status=%s, queue_position=%s",
                task_id,
                task_status,
                status_data.get("task_position"),
            )
            if self.status_callback and status_data.get("task_position") is not None:
                self.status_callback({"task_position": status_data["task_position"]})

            if task_status == "success":
                return status_data
            elif task_status == "failure":
                error_msg = status_data.get("error_message") or "Unknown error"
                raise Exception(f"Docling conversion failed: {error_msg}")
            # else "pending" or "started" – keep polling

            elapsed = time.monotonic() - poll_start
            if elapsed < poll_wait:
                time.sleep(poll_wait - elapsed)

    def _retrieve_result(self, task_id: str) -> dict:
        headers = self._build_headers()
        result_r = requests.get(f"{self.url}/v1/result/{task_id}", headers=headers, timeout=30)
        if not result_r.ok:
            raise Exception(f"Error retrieving Docling result: {result_r.reason}")
        return result_r.json()

    def format_result(self, result_json: dict) -> list[Document]:
        document_data = result_json.get("document", {})
        text = document_data.get("md_content", "<No text content found>")
        metadata = {"Content-Type": self.mime_type} if self.mime_type else {}
        log.debug("Docling extracted text: %s", text)
        return [Document(page_content=text, metadata=metadata)]
    
    def load(self) -> list[Document]:
        task_id, _ = self._submit_file()
        self._poll_task_until_done(task_id)
        result_json = self._retrieve_result(task_id)
        return self.format_result(result_json)

    def load_from_task_id(self, task_id: str) -> list[Document]:
        """Resume processing from an already-submitted docling task_id.

        Used on server restart when docling-serve is still running the task:
        skips re-submission and goes straight to polling → retrieve → format.
        """
        self._poll_task_until_done(task_id)
        result_json = self._retrieve_result(task_id)
        return self.format_result(result_json)

class DoclingLoaderJson(DoclingLoader):
    """Return the Docling JSON document payload split into vector-DB-friendly Documents.

    chunk_mode controls how the structured content is split:
      "item"  - one Document per Docling content item (finest granularity)
      "page"  - one Document per page
      "chunk" - sliding-window chunks with overlap (default)

    For "chunk" mode chunk_size and chunk_overlap are measured in the unit set by
    chunk_content_type: "character" (default) counts Unicode characters; "token" counts
    tiktoken subword tokens.  Both are absolute values, not fractions.
    The page number of every Document is taken from the first content item that
    falls into that chunk / page / item.
    """

    def __init__(
        self,
        *args,
        chunk_mode: str = "chunk",
        chunk_size: int = 1000,
        chunk_overlap: int = 100,
        chunk_content_type: str = "character",
        tiktoken_encoding_name: str = "cl100k_base",
        **kwargs,
    ):
        super().__init__(*args, **kwargs)
        self.chunk_mode = chunk_mode
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.chunk_content_type = chunk_content_type
        # Initialise tiktoken encoder if the unit is "token".
        self._enc = None
        if self.chunk_content_type == "token":
            try:
                import tiktoken
                self._enc = tiktoken.get_encoding(tiktoken_encoding_name)
            except Exception as exc:
                log.warning(
                    "tiktoken unavailable (%s); falling back to character-based chunking", exc
                )
                self.chunk_content_type = "character"
        to_formats = list(self.params.get("to_formats") or ["md"])
        if "json" not in to_formats:
            to_formats.append("json")
        self.params = {**self.params, "to_formats": to_formats}

    # ------------------------------------------------------------------ helpers

    def _word_sizes(self, words: list[str]) -> list[int]:
        """Return the size of each word-token in the configured unit.

        "character" (default): number of Unicode characters.
        "token"               : tiktoken subword token count (approximate per word).
        """
        if self.chunk_content_type == "token" and self._enc is not None:
            return [len(self._enc.encode(w)) for w in words]
        return [len(w) for w in words]

    def _extract_ordered_items(self, document_data: dict) -> list[dict]:
        """Return all content items sorted by page then spatial position (top→bottom, left→right).

        Each entry: {"text": str, "label": str, "page_no": int, "bbox": tuple|None}
        Items without any prov entry containing a page_no are omitted.
        """
        json_content = document_data.get("json_content") or document_data.get("json") or {}

        items_out = []
        for name in ("texts", "tables", "pictures", "key_value_items"):
            cand = document_data.get(name) or (
                json_content.get(name) if isinstance(json_content, dict) else None
            )
            if not isinstance(cand, list):
                continue
            for it in cand:
                text = it.get("text") or it.get("orig") or ""
                label = it.get("label") or ""
                for prov_entry in it.get("prov") or []:
                    if not isinstance(prov_entry, dict):
                        continue
                    pno = prov_entry.get("page_no")
                    if not isinstance(pno, int):
                        continue
                    bbox = None
                    b = prov_entry.get("bbox")
                    if isinstance(b, dict):
                        try:
                            bbox = (
                                float(b.get("l", 0)),
                                float(b.get("t", 0)),
                                float(b.get("r", 0)),
                                float(b.get("b", 0)),
                            )
                        except Exception:
                            bbox = None
                    items_out.append(
                        {
                            "text": text,
                            "label": label,
                            "page_no": pno,
                            "bbox": bbox,
                            "level": it.get("level"),          # section_header depth
                            "marker": it.get("marker"),        # list_item bullet/number
                            "enumerated": it.get("enumerated"),# list_item ordered flag
                            "hyperlink": it.get("hyperlink"),  # hyperlink target if any
                        }
                    )
                    break  # first valid prov entry is sufficient for sorting/metadata

        def _sort_key(entry):
            b = entry["bbox"]
            if b:
                cy = (b[1] + b[3]) / 2.0
                cx = (b[0] + b[2]) / 2.0
                return (entry["page_no"], -cy, cx)  # -cy: PDF y-axis goes upward
            return (entry["page_no"], 0.0, 0.0)

        items_out.sort(key=_sort_key)
        return items_out

    def _extract_doc_metadata(self, document_data: dict) -> dict:
        """Extract document-level fields shared by every point produced from this document.

        These are stored per-point because the vector DB has no collection-level
        metadata store — all filterable/retrievable fields must live on each point.

        origin and name live inside json_content (the DoclingDocument), not at the
        outer ExportDocumentResponse level.
        """
        jc = document_data.get("json_content") or document_data.get("json") or {}
        if not isinstance(jc, dict):
            jc = {}
        origin = jc.get("origin") or {}
        meta: dict = {"Content-Type": "application/json"}
        if origin.get("filename"):
            meta["filename"] = origin["filename"]
        if origin.get("mimetype"):
            meta["source_mimetype"] = origin["mimetype"]
        if origin.get("uri"):
            meta["source_uri"] = origin["uri"]
        if origin.get("binary_hash") is not None:
            meta["binary_hash"] = str(origin["binary_hash"])
        name = jc.get("name") or document_data.get("name")
        if name:
            meta["document_name"] = name
        return meta

    def _format_item_mode(self, document_data: dict) -> list[Document]:
        """One Document per Docling content item."""
        items = self._extract_ordered_items(document_data)
        doc_meta = self._extract_doc_metadata(document_data)
        docs = []
        for it in items:
            text = it["text"]
            if not text:
                continue
            metadata = {
                **doc_meta,
                "docling_document": document_data,
                "page": it["page_no"],
                "label": it["label"],
                "chunk_mode": "item",
            }
            if it["bbox"]:
                metadata["bbox"] = it["bbox"]
            if it.get("level") is not None:
                metadata["level"] = it["level"]
            if it.get("marker"):
                metadata["marker"] = it["marker"]
            if it.get("enumerated") is not None:
                metadata["enumerated"] = it["enumerated"]
            if it.get("hyperlink"):
                metadata["hyperlink"] = it["hyperlink"]
            docs.append(Document(page_content=text, metadata=metadata))
        log.debug("Docling JSON item-mode produced %d documents", len(docs))
        return docs

    def _format_page_mode(self, document_data: dict) -> list[Document]:
        """One Document per page, all items on that page joined with newlines."""
        items = self._extract_ordered_items(document_data)
        doc_meta = self._extract_doc_metadata(document_data)
        jc = document_data.get("json_content") or document_data.get("json") or {}
        pages_map = (jc.get("pages") if isinstance(jc, dict) else None) or {}
        pages: dict[int, list[str]] = {}
        for it in items:
            text = it["text"]
            if not text:
                continue
            pages.setdefault(it["page_no"], []).append(text)
        docs = []
        for pno in sorted(pages.keys()):
            page_text = "\n".join(pages[pno])
            metadata = {
                **doc_meta,
                "docling_document": document_data,
                "page": pno,
                "chunk_mode": "page",
            }
            page_info = pages_map.get(str(pno)) or {}
            if isinstance(page_info, dict):
                # DoclingDocument serialises page size as {"size": {"width": …, "height": …}}
                size = page_info.get("size") or {}
                if size.get("width") is not None:
                    metadata["page_width"] = size["width"]
                if size.get("height") is not None:
                    metadata["page_height"] = size["height"]
            docs.append(Document(page_content=page_text, metadata=metadata))
        log.debug("Docling JSON page-mode produced %d documents", len(docs))
        return docs

    def _format_chunk_mode(self, document_data: dict) -> list[Document]:
        """Sliding-window chunks measured in the configured unit (characters or tiktoken tokens).

        Internally the text is tokenised word-by-word to preserve all whitespace and
        newlines. The window boundaries are determined by the cumulative size of those
        word-tokens in the requested unit so the assembled chunk text is lossless.

        chunk_size    – maximum chunk size in the configured unit.
        chunk_overlap – how many units of overlap to keep between consecutive chunks
                        (absolute, same unit as chunk_size).
        """
        items = self._extract_ordered_items(document_data)
        doc_meta = self._extract_doc_metadata(document_data)

        # Build a flat list of word-tokens and a parallel list of their page numbers.
        words: list[str] = []
        word_pages: list[int] = []
        for it in items:
            text = it["text"]
            if not text:
                continue
            text = text.replace('\r\n', '\n').replace('\r', '\n')
            # Each whitespace run and each non-whitespace run becomes one token so
            # that joining with "".join() reconstructs the original text exactly.
            item_words = re.findall(r'\n|[ \t]+|[^\s]+', text)
            words.extend(item_words)
            word_pages.extend([it["page_no"]] * len(item_words))

        if not words:
            document_text = document_data.get("md_content") or document_data.get("text") or ""
            return [
                Document(
                    page_content=document_text,
                    metadata={**doc_meta, "docling_document": document_data},
                )
            ]

        # Per-word sizes in the configured unit.
        sizes = self._word_sizes(words)

        chunk_size = max(1, self.chunk_size)
        chunk_overlap = max(0, min(self.chunk_overlap, chunk_size - 1))
        step_size = max(1, chunk_size - chunk_overlap)

        n = len(words)
        docs = []
        start_idx = 0

        while start_idx < n:
            # Collect words until chunk_size units are accumulated.
            acc = 0
            end_idx = start_idx
            while end_idx < n:
                acc += sizes[end_idx]
                end_idx += 1
                if acc >= chunk_size:
                    break

            chunk_text = "".join(words[start_idx:end_idx])
            page_no = word_pages[start_idx]
            metadata = {
                **doc_meta,
                "docling_document": document_data,
                "page": page_no,
                "chunk_mode": "chunk",
                "chunk_start_word": start_idx,
                "chunk_end_word": end_idx,
            }
            docs.append(Document(page_content=chunk_text, metadata=metadata))

            if end_idx >= n:
                break

            # Advance by step_size units from the current start.
            adv = 0
            new_start = start_idx
            while new_start < end_idx:
                adv += sizes[new_start]
                new_start += 1
                if adv >= step_size:
                    break
            # Guard: always advance by at least one word to prevent infinite loops.
            if new_start <= start_idx:
                new_start = start_idx + 1
            start_idx = new_start

        log.debug(
            "Docling JSON chunk-mode produced %d documents "
            "(chunk_size=%d %s, chunk_overlap=%d %s)",
            len(docs),
            chunk_size,
            self.chunk_content_type,
            chunk_overlap,
            self.chunk_content_type,
        )
        return docs

    def format_result(self, result_json: dict) -> list[Document]:
        document_data = result_json.get("document", {})
        json_content = document_data.get("json_content") or document_data.get("json")

        if isinstance(json_content, dict) and isinstance(json_content.get("pages"), dict):
            if self.chunk_mode == "item":
                return self._format_item_mode(document_data)
            elif self.chunk_mode == "page":
                return self._format_page_mode(document_data)
            else:  # "chunk" is the default
                return self._format_chunk_mode(document_data)

        # Fallback: structured JSON not available, return raw markdown/text.
        document_text = document_data.get("md_content") or document_data.get("text") or ""
        metadata = {"Content-Type": "application/json", "docling_document": document_data}
        log.debug("Docling JSON result extracted (fallback): keys=%s", list(document_data.keys()))
        return [Document(page_content=document_text, metadata=metadata)]

class Loader:
    def __init__(self, engine: str = '', **kwargs):
        self.engine = engine
        self.user = kwargs.get('user', None)
        self.metadata = kwargs.get('metadata', {})
        self.kwargs = kwargs

    def load(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        loader = self._get_loader(filename, file_content_type, file_path)
        docs = loader.load()
        return [Document(page_content=ftfy.fix_text(doc.page_content), metadata=doc.metadata) for doc in docs]

    async def aload(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """
        Async wrapper around `load`.

        Document loaders dispatched by `_get_loader` (PyMuPDF, Unstructured,
        python-docx, Tika, etc.) are uniformly synchronous and CPU/IO-bound.
        Calling `load` directly from an async handler would block the event
        loop for the entire parse — minutes for large PDFs. This offloads
        the work to a worker thread so the loop stays responsive.
        """
        return await asyncio.to_thread(self.load, filename, file_content_type, file_path)

    def _is_text_file(self, file_ext: str, file_content_type: str) -> bool:
        return file_ext in known_source_ext or (
            file_content_type
            and file_content_type.find('text/') >= 0
            # Avoid text/html files being detected as text
            and not file_content_type.find('html') >= 0
        )

    def _detect_text_encoding(self, file_path: str) -> str:
        """Detect the encoding of a text file with CJK-aware fallbacks.

        Langchain's ``TextLoader`` uses chardet internally when
        ``autodetect_encoding=True``, but chardet frequently misidentifies
        CJK encodings (e.g. GB18030 detected as GB2312 or even Cyrillic).
        This method replaces that by:

        1. Trying UTF-8 first (fast path for the vast majority of files).
        2. Using chardet as a *hint* to prioritise the right CJK codec
           family, but mapping subset names to their superset
           (e.g. GB2312 → gb18030).
        3. Validating that decoded text actually contains CJK characters,
           guarding against codecs that "succeed" but produce garbage.
        4. Falling back to latin-1 (always valid, ftfy fixes mojibake later).
        """
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
        except OSError:
            return 'utf-8'

        if not raw:
            return 'utf-8'

        # Fast path: most files are UTF-8
        try:
            raw.decode('utf-8')
            return 'utf-8'
        except UnicodeDecodeError:
            pass

        # Use chardet as a hint, not as ground truth
        import chardet

        detected = chardet.detect(raw)
        detected_enc = (detected.get('encoding') or '').lower().replace('-', '').replace('_', '')

        # Map chardet's detected encoding to the correct superset codec.
        # chardet often reports GB2312 for content that is actually GB18030;
        # GB18030 is a strict superset of both GB2312 and GBK.
        _ENC_FAMILY = {
            'gb2312': 'gb18030',
            'gb18030': 'gb18030',
            'gbk': 'gb18030',
            'big5': 'big5',
            'euckr': 'euc-kr',
            'eucjp': 'euc-jp',
            'iso2022jp': 'euc-jp',
            'shiftjis': 'shift_jis',
        }

        # Build priority list: chardet-hinted codec first, then remaining CJK
        base_order = ['gb18030', 'big5', 'euc-kr', 'euc-jp']
        hinted = _ENC_FAMILY.get(detected_enc)
        if hinted and hinted in base_order:
            ordered = [hinted] + [e for e in base_order if e != hinted]
        else:
            ordered = base_order

        for enc in ordered:
            try:
                text = raw.decode(enc)
                if text.strip() and self._has_cjk_characters(text):
                    log.info(
                        'Detected encoding %s for %s (chardet guessed %s)',
                        enc,
                        file_path,
                        detected.get('encoding'),
                    )
                    return enc
            except (UnicodeDecodeError, LookupError):
                continue

        # If chardet gave a non-CJK answer that isn't in our family map,
        # try it directly — it might be a valid Western encoding.
        chardet_encoding = detected.get('encoding')
        if chardet_encoding:
            try:
                raw.decode(chardet_encoding)
                log.info(
                    'Using chardet-detected encoding %s for %s',
                    chardet_encoding,
                    file_path,
                )
                return chardet_encoding
            except (UnicodeDecodeError, LookupError):
                pass

        # latin-1 is the ultimate fallback: every byte 0x00–0xFF is valid.
        # ftfy.fix_text() (applied downstream) repairs most mojibake that
        # results from treating Windows-1252 content as Latin-1.
        log.info('Falling back to latin-1 encoding for %s', file_path)
        return 'latin-1'

    @staticmethod
    def _has_cjk_characters(text: str, threshold: float = 0.05) -> bool:
        """Check if decoded text contains a meaningful proportion of CJK characters.

        This guards against codecs that technically "succeed" but decode the
        bytes into wrong Unicode codepoints (e.g. PUA chars, random symbols).
        A genuine CJK document should have at least ``threshold`` fraction of
        its non-whitespace characters in CJK Unicode blocks.
        """
        if not text:
            return False

        cjk_count = 0
        total = 0
        for ch in text:
            if ch.isspace():
                continue
            total += 1
            cp = ord(ch)
            if (
                0x4E00 <= cp <= 0x9FFF  # CJK Unified Ideographs
                or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
                or 0x20000 <= cp <= 0x2A6DF  # CJK Extension B
                or 0x2A700 <= cp <= 0x2B73F  # CJK Extension C
                or 0x2B740 <= cp <= 0x2B81F  # CJK Extension D
                or 0xF900 <= cp <= 0xFAFF  # CJK Compatibility Ideographs
                or 0x3000 <= cp <= 0x303F  # CJK Symbols and Punctuation
                or 0x3040 <= cp <= 0x309F  # Hiragana
                or 0x30A0 <= cp <= 0x30FF  # Katakana
                or 0xAC00 <= cp <= 0xD7AF  # Hangul Syllables
                or 0xFF00 <= cp <= 0xFFEF  # Halfwidth and Fullwidth Forms
            ):
                cjk_count += 1

        if total == 0:
            return False

        return (cjk_count / total) >= threshold

    def _get_loader(self, filename: str, file_content_type: str, file_path: str):
        file_ext = filename.split('.')[-1].lower()

        if (
            self.engine == 'external'
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL')
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY')
        ):
            loader = ExternalDocumentLoader(
                file_path=file_path,
                url=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL'),
                api_key=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY'),
                mime_type=file_content_type,
                user=self.user,
                headers=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_HEADERS'),
                metadata={
                    **self.metadata,
                    'file_name': filename,
                    'file_content_type': file_content_type,
                },
            )
        elif self.engine == 'tika' and self.kwargs.get('TIKA_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TikaLoader(
                    url=self.kwargs.get('TIKA_SERVER_URL'),
                    file_path=file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                )
        elif (
            self.engine == 'datalab_marker'
            and self.kwargs.get('DATALAB_MARKER_API_KEY')
            and file_ext
            in [
                'pdf',
                'xls',
                'xlsx',
                'ods',
                'doc',
                'docx',
                'odt',
                'ppt',
                'pptx',
                'odp',
                'html',
                'epub',
                'png',
                'jpeg',
                'jpg',
                'webp',
                'gif',
                'tiff',
            ]
        ):
            api_base_url = self.kwargs.get('DATALAB_MARKER_API_BASE_URL', '')
            if not api_base_url or api_base_url.strip() == '':
                api_base_url = 'https://www.datalab.to/api/v1/marker'  # https://github.com/open-webui/open-webui/pull/16867#issuecomment-3218424349

            loader = DatalabMarkerLoader(
                file_path=file_path,
                api_key=self.kwargs['DATALAB_MARKER_API_KEY'],
                api_base_url=api_base_url,
                additional_config=self.kwargs.get('DATALAB_MARKER_ADDITIONAL_CONFIG'),
                use_llm=self.kwargs.get('DATALAB_MARKER_USE_LLM', False),
                skip_cache=self.kwargs.get('DATALAB_MARKER_SKIP_CACHE', False),
                force_ocr=self.kwargs.get('DATALAB_MARKER_FORCE_OCR', False),
                paginate=self.kwargs.get('DATALAB_MARKER_PAGINATE', False),
                strip_existing_ocr=self.kwargs.get('DATALAB_MARKER_STRIP_EXISTING_OCR', False),
                disable_image_extraction=self.kwargs.get('DATALAB_MARKER_DISABLE_IMAGE_EXTRACTION', False),
                format_lines=self.kwargs.get('DATALAB_MARKER_FORMAT_LINES', False),
                output_format=self.kwargs.get('DATALAB_MARKER_OUTPUT_FORMAT', 'markdown'),
            )
        elif self.engine in ("docling", "docling_json") and self.kwargs.get("DOCLING_SERVER_URL"):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                # Build params for DoclingLoader
                params = self.kwargs.get('DOCLING_PARAMS', {})
                if not isinstance(params, dict):
                    try:
                        params = json.loads(params)
                    except json.JSONDecodeError:
                        log.error('Invalid DOCLING_PARAMS format, expected JSON object')
                        params = {}

                docling_timeout = self.kwargs.get("DOCLING_SERVE_TIMEOUT")
                if docling_timeout is not None:
                    try:
                        docling_timeout = int(docling_timeout)
                    except (ValueError, TypeError):
                        docling_timeout = None

                if self.engine == "docling_json":
                    try:
                        _chunk_size = int(self.kwargs.get("CHUNK_SIZE", 1000))
                    except (ValueError, TypeError):
                        _chunk_size = 1000
                    try:
                        _chunk_overlap = int(self.kwargs.get("CHUNK_OVERLAP", 100))
                    except (ValueError, TypeError):
                        _chunk_overlap = 100
                    _text_splitter = self.kwargs.get("TEXT_SPLITTER") or ""
                    _chunk_content_type = "token" if _text_splitter == "token" else "character"
                    loader = DoclingLoaderJson(
                        url=self.kwargs.get("DOCLING_SERVER_URL"),
                        api_key=self.kwargs.get("DOCLING_API_KEY", None),
                        file_path=file_path,
                        mime_type=file_content_type,
                        params=params,
                        timeout=docling_timeout,
                        status_callback=self.kwargs.get("DOCLING_STATUS_CALLBACK"),
                        chunk_mode=self.kwargs.get("DOCLING_JSON_CHUNK_MODE", "chunk"),
                        chunk_size=_chunk_size,
                        chunk_overlap=_chunk_overlap,
                        chunk_content_type=_chunk_content_type,
                        tiktoken_encoding_name=str(
                            self.kwargs.get("TIKTOKEN_ENCODING_NAME") or "cl100k_base"
                        ),
                    )
                else:
                    loader = DoclingLoader(
                        url=self.kwargs.get("DOCLING_SERVER_URL"),
                        api_key=self.kwargs.get("DOCLING_API_KEY", None),
                        file_path=file_path,
                        mime_type=file_content_type,
                        params=params,
                        timeout=docling_timeout,
                        status_callback=self.kwargs.get("DOCLING_STATUS_CALLBACK"),
                    )
        elif (
            self.engine == 'document_intelligence'
            and self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT') != ''
            and (
                file_ext in ['pdf', 'docx', 'ppt', 'pptx']
                or file_content_type
                in [
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.ms-powerpoint',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                ]
            )
        ):
            if self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY') != '':
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    api_key=self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY'),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
            else:
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    azure_credential=DefaultAzureCredential(),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
        elif self.engine == 'mineru' and file_ext in self.kwargs.get('MINERU_FILE_EXTENSIONS', ['pdf']):
            mineru_timeout = self.kwargs.get('MINERU_API_TIMEOUT', 300)
            if mineru_timeout:
                try:
                    mineru_timeout = int(mineru_timeout)
                except ValueError:
                    mineru_timeout = 300
            loader = MinerULoader(
                file_path=file_path,
                api_mode=self.kwargs.get('MINERU_API_MODE', 'local'),
                api_url=self.kwargs.get('MINERU_API_URL', 'http://localhost:8000'),
                api_key=self.kwargs.get('MINERU_API_KEY', ''),
                params=self.kwargs.get('MINERU_PARAMS', {}),
                timeout=mineru_timeout,
                max_markdown_bytes=MINERU_MAX_MARKDOWN_BYTES,
            )
        elif (
            self.engine == 'mistral_ocr'
            and self.kwargs.get('MISTRAL_OCR_API_KEY') != ''
            and file_ext in ['pdf']  # Mistral OCR currently only supports PDF and images
        ):
            loader = MistralLoader(
                base_url=self.kwargs.get('MISTRAL_OCR_API_BASE_URL'),
                api_key=self.kwargs.get('MISTRAL_OCR_API_KEY'),
                file_path=file_path,
                use_base64=self.kwargs.get('MISTRAL_OCR_USE_BASE64', False),
            )
        elif self.engine == 'paddleocr_vl' and self.kwargs.get('PADDLEOCR_VL_TOKEN') != '':
            loader = PaddleOCRVLLoader(
                api_url=self.kwargs.get('PADDLEOCR_VL_BASE_URL'),
                token=self.kwargs.get('PADDLEOCR_VL_TOKEN'),
                file_path=file_path,
            )
        else:
            if file_ext == 'pdf':
                loader = PyPDFLoader(
                    file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                    mode=self.kwargs.get('PDF_LOADER_MODE', 'page'),
                )
            elif file_ext == 'csv':
                loader = CSVLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'rst':
                try:
                    from langchain_community.document_loaders import UnstructuredRSTLoader

                    loader = UnstructuredRSTLoader(file_path, mode='elements')
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .rst file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'xml':
                try:
                    from langchain_community.document_loaders import UnstructuredXMLLoader

                    loader = UnstructuredXMLLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .xml file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext in ['htm', 'html']:
                loader = BSHTMLLoader(file_path, open_encoding='unicode_escape')
            elif file_ext == 'md':
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_content_type == 'application/epub+zip':
                try:
                    from langchain_community.document_loaders import UnstructuredEPubLoader

                    loader = UnstructuredEPubLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .epub files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif (
                file_content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                or file_ext == 'docx'
            ):
                loader = Docx2txtLoader(file_path)
            elif file_ext == 'doc' or file_content_type == 'application/msword':
                try:
                    from langchain_community.document_loaders import UnstructuredWordDocumentLoader

                    loader = UnstructuredWordDocumentLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .doc files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif file_content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            ] or file_ext in ['xls', 'xlsx']:
                try:
                    from langchain_community.document_loaders import UnstructuredExcelLoader

                    loader = UnstructuredExcelLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to pandas for Excel file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = ExcelLoader(file_path)
            elif file_content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ] or file_ext in ['ppt', 'pptx']:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader

                    loader = UnstructuredPowerPointLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to python-pptx for PowerPoint file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = PptxLoader(file_path)
            elif file_ext == 'msg':
                loader = OutlookMessageLoader(file_path)
            elif file_ext == 'odt':
                try:
                    from langchain_community.document_loaders import UnstructuredODTLoader

                    loader = UnstructuredODTLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .odt files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))

        return loader
