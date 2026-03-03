"""
title: O365 Search & Mail Summary Tool
description: |
  Vier Tools in einem für Microsoft 365:
  1. search_and_read_m365: Sucht in SharePoint/OneDrive nach Dokumenten (PDF, Word, Excel, PPT
     etc.) mit intelligenter mehrstufiger Copilot-Style Fallback-Suche.
  2. summarize_my_emails: Holt E-Mails aus Outlook, bewertet sie nach Relevanz
     (Copilot-Style Scoring: Direktanfragen, Action-Keywords, Themen-Match, System-Noise)
     und gruppiert sie in Handlungsrelevant / Informativ / Automatisch.
  3. search_email_detail: Sucht nach einer spezifischen Mail und zeigt den vollständigen
     Inhalt direkt im Chat an.
  4. generate_reply_mailto: Erstellt einen mailto:-Link zum Öffnen einer vorausgefüllten
     Antwort direkt in Standard Outlook (Empfänger, Betreff und Entwurf vorbefüllt).
author: open-webui
version: 5.1.0
required_open_webui_version: 0.3.9
"""

import httpx
import asyncio
import io
import re
import time
import datetime
from urllib.parse import quote
from pydantic import BaseModel, Field
from typing import Optional, Callable, Any


# ────────────────────────────────────────────────────────────────────────────
# Unterstützte Dokument-Endungen
# ────────────────────────────────────────────────────────────────────────────
DOCUMENT_EXTENSIONS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".doc": "Word (alt)",
    ".xlsx": "Excel",
    ".xls": "Excel (alt)",
    ".pptx": "PowerPoint",
    ".ppt": "PowerPoint (alt)",
    ".txt": "Text",
    ".csv": "CSV",
    ".md": "Markdown",
    ".odt": "OpenDocument Text",
    ".ods": "OpenDocument Tabelle",
    ".rtf": "RTF",
}

# KQL Dateifilter (Stufen 1–3)
KQL_FILE_FILTER = (
    "(FileExtension:pdf OR FileExtension:docx OR FileExtension:doc "
    "OR FileExtension:xlsx OR FileExtension:xls "
    "OR FileExtension:pptx OR FileExtension:ppt "
    "OR FileExtension:txt OR FileExtension:csv "
    "OR FileExtension:md OR FileExtension:rtf)"
)

# Deutsche + englische Stoppwörter die wir rausfiltern
# (häufige Wörter die die Suche nicht verbessern)
STOP_WORDS = {
    "der", "die", "das", "dem", "den", "ein", "eine", "einen", "einem", "einer",
    "und", "oder", "aber", "mit", "für", "von", "zu", "auf", "in", "an", "im",
    "ist", "sind", "war", "waren", "wird", "werden", "hat", "haben", "hatte",
    "ich", "du", "er", "sie", "wir", "ihr", "sie", "es", "se", "wie", "wo",
    "was", "wer", "welche", "welcher", "welches", "durch", "nach", "über",
    "unter", "aus", "bei", "seit", "bis", "gegen", "ohne", "um", "als",
    "the", "a", "an", "and", "or", "but", "with", "for", "from", "to", "of",
    "in", "on", "at", "by", "is", "are", "was", "were", "has", "have", "had",
    "it", "this", "that", "these", "those", "be", "do", "can", "will", "how",
}

# Häufige deutsche/englische Inhaltswörter, die zu generisch für eine Einzelwortsuche sind.
# Wörter die NICHT hier drin sind → wahrscheinlich Eigenname/Marke (z.B. "Eurotherme", "SCOM")
# → werden bei Einzelwort-Fallbacks bevorzugt zuerst gesucht.
EXTENDED_COMMON_WORDS = {
    # Allgemeine deutsche Substantive
    "rabatt", "preis", "kosten", "info", "information", "thema", "bereich",
    "punkt", "schritt", "regel", "recht", "fall", "art", "weg", "form",
    "teil", "ziel", "gruppe", "liste", "plan", "daten", "bericht", "datei",
    "dokument", "version", "stand", "status", "datum", "zeit", "inhalt",
    "mitarbeiter", "kunde", "firma", "partner", "team", "person", "user",
    "produkt", "loesung", "lösung", "system", "service", "prozess", "projekt",
    "ergebnis", "text", "beispiel", "vorgang", "aufgabe", "leistung",
    "funktion", "anleitung", "beschreibung", "uebersicht", "übersicht",
    "zusammenfassung", "vergünstigung", "vergünstigungen", "vorteil", "vorteile",
    "rabatte", "aktion", "angebot", "angebote", "gutschein",
    # Häufige Adjektive/Adverbien
    "neu", "alt", "intern", "extern", "allgemein", "aktuell", "wichtig",
    # Englische häufige Inhaltswörter
    "discount", "price", "cost", "guide", "manual", "report",
    "employee", "customer", "product", "solution", "process", "project",
    "benefit", "offer",    "voucher",
}


# ────────────────────────────────────────────────────────────────────────────
# E-Mail Scoring Signale (Copilot-Style)
# ────────────────────────────────────────────────────────────────────────────

# Direkte Handlungsaufforderungen → hoher Score
EMAIL_ACTION_SIGNALS = [
    # Deutsch
    "bitte", "könntest du", "kannst du", "können sie", "kurze frage", "kurze anfrage",
    "wäre es möglich", "bitte um", "brauchst du", "benötigst du", "benötigen sie",
    "freigeben", "freigabe", "genehmigen", "genehmigung", "bestätigen", "bestätigung",
    "klären", "klärung", "prüfen", "prüfe", "rückmeldung", "feedback", "dringend",
    "bis wann", "deadline", "frist", "termin", "asap", "sofort", "heute noch",
    "absicherung", "zuständig", "aktion erforderlich", "action required",
    # Englisch
    "please", "can you", "could you", "would you", "do you", "are you",
    "quick question", "approve", "approval", "confirm", "confirmation",
    "clarify", "clarification", "review", "check", "feedback needed",
    "urgent", "important", "action needed", "action required", "response needed",
]

# System-/Automatisierungsmails → negativer Score
EMAIL_SYSTEM_SIGNALS = [
    # Absender-Patterns (lowercase)
    "noreply", "no-reply", "donotreply", "do-not-reply", "notifications@",
    "mailer-daemon", "postmaster", "alerts@", "monitoring@", "automated@",
    "azure devops", "github", "gitlab", "jenkins", "jira", "confluence",
    "microsoft teams", "teams notification", "sharepoint",
    # Betreff-Patterns
    "pull request", "pr –", "pr:", "[pr]", "build succeeded", "build failed",
    "pipeline", "deployment", "release", "abandoned", "resolved", "closed",
    "message center", "service health", "azure notification",
    "unsubscribe", "newsletter", "marketing",
]


# ────────────────────────────────────────────────────────────────────────────
# Synonyme für häufige DE/EN Business-Begriffe (KQL OR-Expansion)
# ────────────────────────────────────────────────────────────────────────────
SYNONYM_MAP: dict[str, list[str]] = {
    # Deutsch → Synonyme (DE + EN)
    "anleitung":       ["tutorial", "guide", "howto", "handbuch"],
    "handbuch":        ["manual", "guide", "anleitung"],
    "dokumentation":   ["documentation", "doku", "handbuch"],
    "vergünstigung":   ["rabatt", "discount", "vorteil", "benefit", "preisnachlass"],
    "vergünstigungen": ["rabatte", "discounts", "vorteile", "benefits"],
    "rabatt":          ["vergünstigung", "discount", "preisnachlass"],
    "rabatte":         ["vergünstigungen", "discounts"],
    "vorteil":         ["benefit", "vergünstigung"],
    "vorteile":        ["benefits", "vergünstigungen"],
    "mitarbeiter":     ["employee", "personal"],
    "kunde":           ["customer", "client"],
    "kunden":          ["customers", "clients"],
    "bericht":         ["report", "auswertung", "analyse"],
    "berichte":        ["reports", "auswertungen"],
    "übersicht":       ["uebersicht", "overview", "überblick"],
    "uebersicht":      ["übersicht", "overview"],
    "angebot":         ["offer", "proposal"],
    "angebote":        ["offers", "proposals"],
    "prozess":         ["process", "ablauf", "workflow"],
    "projekt":         ["project", "vorhaben"],
    "projekte":        ["projects"],
    "schulung":        ["training", "kurs", "seminar", "weiterbildung"],
    "schulungen":      ["trainings", "kurse", "seminare"],
    "einführung":      ["einfuehrung", "introduction", "intro", "onboarding"],
    "einfuehrung":     ["einführung", "introduction", "intro"],
    "richtlinie":      ["policy", "regel", "guideline"],
    "richtlinien":     ["policies", "regeln", "guidelines"],
    "vorlage":         ["template", "muster"],
    "vorlagen":        ["templates"],
    "protokoll":       ["minutes", "log", "aufzeichnung"],
    "protokolle":      ["minutes", "logs"],
    "besprechung":     ["meeting", "sitzung"],
    "sitzung":         ["meeting", "besprechung"],
    "vertrag":         ["contract", "vereinbarung"],
    "verträge":        ["contracts"],
    "preisliste":      ["pricelist", "preise", "tarif"],
    "preise":          ["prices", "kosten"],
    "kosten":          ["costs", "preise", "aufwand"],
    "budget":          ["kosten", "cost", "haushalt"],
    "urlaubsantrag":   ["vacation request", "leave request", "urlaubsgesuch"],
    "gehalt":          ["salary", "lohn", "entlohnung"],
    "reisekosten":     ["travel costs", "reisespesen", "spesen"],
    # Englisch → Synonyme (EN + DE)
    "guide":           ["anleitung", "howto", "tutorial", "handbuch"],
    "tutorial":        ["anleitung", "guide", "howto"],
    "manual":          ["handbuch", "anleitung", "guide"],
    "documentation":   ["dokumentation", "doku"],
    "discount":        ["rabatt", "vergünstigung", "preisnachlass"],
    "benefit":         ["vorteil", "vergünstigung"],
    "benefits":        ["vorteile", "vergünstigungen"],
    "report":          ["bericht", "auswertung"],
    "reports":         ["berichte", "auswertungen"],
    "overview":        ["übersicht", "uebersicht"],
    "template":        ["vorlage", "muster"],
    "meeting":         ["besprechung", "sitzung"],
    "training":        ["schulung", "kurs"],
    "trainings":       ["schulungen", "kurse"],
    "policy":          ["richtlinie", "regel"],
    "policies":        ["richtlinien", "regeln"],
    "contract":        ["vertrag", "vereinbarung"],
    "contracts":       ["verträge"],
    "process":         ["prozess", "ablauf"],
    "project":         ["projekt", "vorhaben"],
    "offer":           ["angebot"],
    "offers":           ["angebote"],
    "introduction":    ["einführung", "einfuehrung", "intro"],
    "salary":          ["gehalt", "lohn"],
    "travel":          ["reise", "dienstreise"],
}


def _expand_term_with_synonyms(word: str, max_synonyms: int = 3) -> str:
    """
    Gibt einen KQL-Ausdruck zurück: Wort allein oder (Wort OR Syn1 OR Syn2 ...).
    Maximal max_synonyms Synonyme werden hinzugefügt.
    Beispiel: "Anleitung" → "(Anleitung OR tutorial OR guide OR howto)"
    """
    synonyms = SYNONYM_MAP.get(word.lower(), [])
    if not synonyms:
        return word
    all_terms = [word] + synonyms[:max_synonyms]
    return f"({' OR '.join(all_terms)})"


def _score_hit_by_keywords(text: str, words: list[str]) -> int:
    """
    Zählt wie oft die Suchbegriffe (case-insensitiv) im extrahierten Text vorkommen.
    Wird für das Keyword-Density-Ranking der Treffer verwendet.
    Ein höherer Wert = relevanterer Treffer.
    """
    if not text:
        return 0
    text_lower = text.lower()
    return sum(text_lower.count(w.lower()) for w in words)


def _compute_combined_score(
    doc: dict,
    search_words: list[str],
    api_rank: int = 0,
) -> float:
    """
    Kombiniertes Ranking wie Microsoft Copilot:
    1. Graph-API-Rank (höherer Rang = besser, invertiert: Platz 1 = höchster Score)
    2. Keyword-Density im extrahierten Text
    3. Aktualität (neuere Dateien werden bevorzugt – max. +20 Punkte)

    Alle drei Faktoren werden gewichtet summiert.
    """
    # --- Keyword-Dichte ---
    keyword_score = _score_hit_by_keywords(
        doc.get("text") or "", search_words
    )

    # --- API-Rank: Platz 1 → 15 Punkte, Platz 10 → ~6 Punkte ---
    # Graph liefert rank 1-basiert, niedrigerer Rank = relevanter
    rank_score = max(0, 15 - api_rank) if api_rank > 0 else 0

    # --- Aktualität: bis zu 20 Bonuspunkte für neuere Dateien ---
    recency_score = 0
    modified_str = doc.get("modified", "")
    if modified_str and modified_str != "?":
        try:
            import datetime
            mod_date = datetime.date.fromisoformat(modified_str[:10])
            today = datetime.date.today()
            days_old = (today - mod_date).days
            # 0 Tage alt → 20 Punkte, 365 Tage → 10 Punkte, 730+ Tage → 0 Punkte
            recency_score = max(0, 20 - int(days_old / 36.5))
        except Exception:
            pass

    return keyword_score + rank_score + recency_score


def _is_document(filename: str) -> bool:
    lower = filename.lower()
    return any(lower.endswith(ext) for ext in DOCUMENT_EXTENSIONS)


def _get_file_type(filename: str) -> str:
    lower = filename.lower()
    for ext, name in DOCUMENT_EXTENSIONS.items():
        if lower.endswith(ext):
            return name
    return "Unbekannt"


def _deduplicate_preserve_order(words: list[str]) -> list[str]:
    """Entfernt Duplikate, behält Reihenfolge bei (case-insensitive)."""
    seen = set()
    result = []
    for w in words:
        low = w.lower()
        if low not in seen:
            seen.add(low)
            result.append(w)
    return result


def _build_progressive_queries(raw_query: str) -> list[dict]:
    """
    Baut aus einem Suchbegriff automatisch progressive Suchanfragen
    – analog zur Suchlogik von Microsoft Copilot:

    Stufe 0: Phrasensuche ("alle Begriffe als Phrase") + Dateifilter
    Stufe 1: Alle Begriffe AND (mit Synonym-OR-Expansion) + Dateifilter
    Stufe 1b: Alle Begriffe AND (ohne Synonyme) + Dateifilter  [nur wenn Synonyme vorhanden]
    Stufe 2: OR-Breitensuche (alle Begriffe OR) + Dateifilter  ← Copilot-Style: breit & fehlertolerant
    Stufe 3: ~60% der Begriffe AND + Dateifilter
    Stufe 4: Top 3 AND (mit Synonymen) + Dateifilter
    Stufe 5: Top 2 AND (mit Synonymen) – KEIN Dateifilter
    Stufe 6+: Jeder Begriff einzeln (mit Synonymen) – KEIN Dateifilter

    Beispiel für "Eurothermen Mitarbeiter Vergünstigung Rabatt":
      Stufe 0: "Eurothermen Mitarbeiter Vergünstigung Rabatt" + FileFilter
      Stufe 1: Eurothermen AND Mitarbeiter AND (Vergünstigung OR rabatt OR discount OR vorteil ...) AND (Rabatt OR ..) + FileFilter
      Stufe 2: Eurothermen AND Mitarbeiter AND Vergünstigung AND Rabatt + FileFilter (ohne Synonyme)
      Stufe 3: Eurothermen OR Mitarbeiter OR (Vergünstigung OR rabatt ...) OR (Rabatt OR ..) + FileFilter  ← breite OR-Suche
      Stufe 4: Eurothermen AND Mitarbeiter AND Vergünstigung + FileFilter
      Stufe 5: Eurothermen AND Mitarbeiter – alle Typen
      Stufe 6: Eurothermen – alle Typen
    """
    # 1. Tokenisieren & bereinigen
    raw_words = re.sub(r"[^\w\s]", " ", raw_query).split()

    # 2. Stoppwörter rausfiltern
    filtered = [w for w in raw_words if w.lower() not in STOP_WORDS and len(w) > 1]
    if not filtered:
        filtered = raw_words

    # 3. Duplikate entfernen (Reihenfolge beibehalten)
    unique = _deduplicate_preserve_order(filtered)
    if not unique:
        unique = [raw_query.strip()]

    queries = []
    added_combos: set[tuple] = set()

    def _kql_and(words: list[str], with_file_filter: bool, use_synonyms: bool = False) -> str:
        parts = [_expand_term_with_synonyms(w) if use_synonyms else w for w in words]
        expr = " AND ".join(parts)
        if with_file_filter:
            return f"({expr}) AND {KQL_FILE_FILTER}"
        return expr

    def _add_query(
        words: list[str],
        with_file_filter: bool,
        label_prefix: str,
        use_synonyms: bool = False,
        is_phrase: bool = False,
    ):
        if not words:
            return
        combo_key = (tuple(w.lower() for w in words), with_file_filter, use_synonyms, is_phrase)
        if combo_key in added_combos:
            return
        added_combos.add(combo_key)

        word_count = len(words)
        filter_note = "+ Dokumentfilter" if with_file_filter else "alle Dateitypen"
        extras = []
        if is_phrase:
            extras.append("Phrase")
        if use_synonyms and not is_phrase:
            extras.append("+ Synonyme")
        extras_str = f" [{', '.join(extras)}]" if extras else ""
        label = f"{label_prefix}{extras_str} ({word_count} Begriff{'e' if word_count != 1 else ''}, {filter_note})"

        if is_phrase:
            phrase = " ".join(words)
            kql = f'"{phrase}" AND {KQL_FILE_FILTER}' if with_file_filter else f'"{phrase}"'
        else:
            kql = _kql_and(words, with_file_filter, use_synonyms)

        queries.append({
            "index": len(queries) + 1,
            "label": label,
            "words": words[:],
            "kql": kql,
            "with_file_filter": with_file_filter,
            "use_synonyms": use_synonyms,
        })

    total = len(unique)
    has_any_synonym = any(w.lower() in SYNONYM_MAP for w in unique)

    # Stufe 0: Phrasensuche (nur bei mehreren Begriffen)
    if total > 1:
        _add_query(unique, True, "Phrasensuche", is_phrase=True)

    # Stufe 1: Alle Begriffe AND + Synonym-Expansion + Dateifilter
    _add_query(unique, True, "Alle Begriffe AND", use_synonyms=True)

    # Stufe 1b: Alle Begriffe AND ohne Synonyme (Fallback wenn Synonyme vorhanden)
    if has_any_synonym:
        _add_query(unique, True, "Alle Begriffe AND")

    # Stufe 2 (NEU): OR-Breitensuche – Copilot-Style: alle Begriffe mit OR verknüpft
    # Findet Dokumente auch wenn nur EINIGE der Suchbegriffe enthalten sind
    # (z.B. "EUROTHERMEN Bad Schallerbach" wenn Suche nach "Eurothermen Vergünstigung Rabatt")
    if total > 1:
        def _kql_or(words: list[str], with_file_filter: bool, use_synonyms: bool = True) -> str:
            parts = [_expand_term_with_synonyms(w) if use_synonyms else w for w in words]
            expr = " OR ".join(parts)
            if with_file_filter:
                return f"({expr}) AND {KQL_FILE_FILTER}"
            return f"({expr})"

        combo_key_or = (tuple(w.lower() for w in unique), True, "or_broad", False)
        if combo_key_or not in added_combos:
            added_combos.add(combo_key_or)
            word_count = len(unique)
            kql_or = _kql_or(unique, True, use_synonyms=True)
            queries.append({
                "index": len(queries) + 1,
                "label": f"OR Breitensuche [+ Synonyme] ({word_count} Begriffe, + Dokumentfilter)",
                "words": unique[:],
                "kql": kql_or,
                "with_file_filter": True,
                "use_synonyms": True,
            })

    # Stufe 3: ~60% (mind. 4 Begriffe) + Dateifilter
    count_60 = max(4, round(total * 0.6))
    if count_60 < total:
        _add_query(unique[:count_60], True, "Hauptbegriffe AND")

    # Stufe 4: Top 3 Begriffe AND + Synonyme + Dateifilter
    if total > 3:
        _add_query(unique[:3], True, "Top 3 AND", use_synonyms=True)
    elif total == 3:
        _add_query(unique[:2], True, "Top 2 AND", use_synonyms=True)

    # Stufe 5: Top 2 Begriffe + Synonyme – KEIN Dateifilter
    if total >= 2:
        _add_query(unique[:2], False, "Top 2 AND", use_synonyms=True)

    # Stufe 6+: Jedes Wort einzeln + Synonyme – KEIN Dateifilter
    # Spezifische Begriffe (Eigennamen, seltene Wörter) werden zuerst gesucht
    words_by_specificity = sorted(
        unique,
        key=lambda w: (1 if w.lower() in EXTENDED_COMMON_WORDS else 0, len(w))
    )
    for word in words_by_specificity:
        _add_query([word], False, "Einzelbegriff", use_synonyms=True)

    return queries


async def _extract_text_from_bytes(filename: str, raw_bytes: bytes) -> str:
    """Extrahiert lesbaren Text aus Binärdaten je nach Dateityp."""
    lower = filename.lower()

    try:
        if lower.endswith((".txt", ".md", ".csv", ".rtf")):
            for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
                try:
                    return raw_bytes.decode(enc)[:8000]
                except UnicodeDecodeError:
                    continue
            return raw_bytes.decode("utf-8", errors="replace")[:8000]

        elif lower.endswith((".docx", ".odt")):
            try:
                from docx import Document
                doc = Document(io.BytesIO(raw_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n".join(paragraphs)[:8000]
            except ImportError:
                return "[python-docx nicht installiert – pip install python-docx]"
            except Exception as e:
                return f"[Word-Extraktion fehlgeschlagen: {e}]"

        elif lower.endswith((".xlsx", ".ods")):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(
                    io.BytesIO(raw_bytes), read_only=True, data_only=True
                )
                lines = []
                for sheet_name in wb.sheetnames[:3]:
                    ws = wb[sheet_name]
                    lines.append(f"=== Sheet: {sheet_name} ===")
                    for row in ws.iter_rows(max_row=100, values_only=True):
                        row_vals = [str(c) if c is not None else "" for c in row]
                        if any(v.strip() for v in row_vals):
                            lines.append("\t".join(row_vals))
                    if len("\n".join(lines)) > 7000:
                        break
                wb.close()
                return "\n".join(lines)[:8000]
            except ImportError:
                return "[openpyxl nicht installiert – pip install openpyxl]"
            except Exception as e:
                return f"[Excel-Extraktion fehlgeschlagen: {e}]"

        elif lower.endswith(".pdf"):
            try:
                from pdfminer.high_level import extract_text_to_fp
                from pdfminer.layout import LAParams
                output = io.StringIO()
                extract_text_to_fp(
                    io.BytesIO(raw_bytes), output,
                    laparams=LAParams(), output_type="text", codec="utf-8"
                )
                text = output.getvalue().strip()
                if text:
                    return text[:8000]
            except Exception:
                pass
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(raw_bytes))
                text = "\n".join(
                    p.extract_text() or "" for p in reader.pages[:15]
                ).strip()
                if text:
                    return text[:8000]
            except Exception:
                pass
            return "[PDF-Text nicht extrahierbar – pip install pdfminer.six]"

        elif lower.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw_bytes))
                lines = []
                for i, slide in enumerate(prs.slides[:20]):
                    lines.append(f"--- Folie {i+1} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
                    if len("\n".join(lines)) > 7000:
                        break
                return "\n".join(lines)[:8000]
            except ImportError:
                return "[python-pptx nicht installiert – pip install python-pptx]"
            except Exception as e:
                return f"[PowerPoint-Extraktion fehlgeschlagen: {e}]"

        else:
            return "[Dateityp nicht unterstützt]"

    except Exception as e:
        return f"[Extraktionsfehler: {e}]"


# ────────────────────────────────────────────────────────────────────────────
# E-Mail Scoring Funktion
# ────────────────────────────────────────────────────────────────────────────

def _score_email(
    msg: dict,
    user_email: str,
    user_topics: list[str],
) -> tuple[int, list[str], bool]:
    """
    Bewertet eine E-Mail nach Copilot-Logik.
    Gibt (score, reasons, is_system) zurück.

    is_system=True → Mail kommt IMMER in ⧁ Automatisch, egal wie hoch score!

    Scoring:
    +5  Direkt an User adressiert (to:, nicht nur cc:)
    +4  Action-Keywords im Betreff oder Vorschau (Handlungsaufforderung)
    +3  Betreff enthält Fragezeichen (direkte Frage)
    +2  User-Topics treffen im Betreff/Vorschau
    +1  Ungelesen
    -1  Nur CC (kein direkter Empfänger)
    SYSTEM-FLAG: noreply, DevOps, Azure, Monitoring → immer ⧁
    """
    score = 0
    reasons = []

    subject   = (msg.get("subject") or "").lower()
    preview   = (msg.get("bodyPreview") or "").lower()
    sender    = (msg.get("from", {}).get("emailAddress", {}).get("address") or "").lower()
    is_read   = msg.get("isRead", True)
    to_list   = [r.get("emailAddress", {}).get("address", "").lower()
                 for r in (msg.get("toRecipients") or [])]
    text      = subject + " " + preview + " " + sender

    # ── Direkt-Check ─────────────────────────────────────────────
    user_email_lower = user_email.lower()
    if user_email_lower in to_list:
        score += 5
        reasons.append("Direkt adressiert")
    else:
        score -= 1
        reasons.append("Nur CC")

    # ── System-Noise ─────────────────────────────────────────────
    is_system = any(sig in text for sig in EMAIL_SYSTEM_SIGNALS)
    if is_system:
        score -= 4
        reasons.append("System/Automatisch")

    # ── Action-Keywords ───────────────────────────────────────────
    matched_actions = [sig for sig in EMAIL_ACTION_SIGNALS if sig in subject or sig in preview]
    if matched_actions:
        score += 4
        reasons.append(f"Action: '{matched_actions[0]}'")

    # ── Fragezeichen im Betreff = direkte Frage ───────────────────
    if "?" in (msg.get("subject") or ""):
        score += 3
        reasons.append("Frage im Betreff")

    # ── User-Topics ───────────────────────────────────────────────
    matched_topics = [t for t in user_topics if t.lower() in subject or t.lower() in preview]
    if matched_topics:
        score += 2 * len(matched_topics)
        reasons.append(f"Themen: {', '.join(matched_topics[:3])}")

    # ── Ungelesen ─────────────────────────────────────────────────
    if not is_read:
        score += 1
        reasons.append("Ungelesen")

    return score, reasons, is_system


class Tools:
    class Valves(BaseModel):
        debug_mode: bool = Field(
            default=False,
            description="Debug-Modus: Zeigt alle Suchstufen mit vollem KQL-Query-String in der Ausgabe an.",
        )
        max_downloads: int = Field(
            default=5,
            description=(
                "Maximale Anzahl Dokumente die heruntergeladen & extrahiert werden (1–15). "
                "Niedrigerer Wert = weniger Graph-API-Calls, höherer Wert = besseres Ranking. "
                "Standard: 5 (entspricht ~7 API-Calls gesamt statt 17)."
            ),
        )
        oauth_client_secret: str = Field(
            default="",
            description=(
                "Azure App Client Secret – nur für Confidential Client Apps nötig. "
                "Bei PKCE / Public Client Apps leer lassen."
            ),
        )

    class UserValves(BaseModel):
        topics_of_interest: str = Field(
            default="",
            description=(
                "Deine Arbeitsthemen als kommagetrennte Liste – werden beim Mail-Scoring bevorzugt. "
                "Beispiel: PowerAutomate, KI, ISMS, SharePoint, Azure, Monitoring"
            ),
        )
        email_days_back: int = Field(
            default=1,
            description=(
                "Wie viele Tage zurück sollen Mails geladen werden? "
                "1 = nur heute, 7 = letzte Woche. Standard: 1"
            ),
        )
        email_max_fetch: int = Field(
            default=50,
            description=(
                "Maximale Anzahl Mails die von der Graph API geladen werden (10–100). "
                "Höherer Wert = bessere Abdeckung, mehr API-Last. Standard: 50"
            ),
        )

    def __init__(self):
        self.base_url = "https://graph.microsoft.com/v1.0"
        self.valves = self.Valves()
        self.user_valves = self.UserValves()

    async def search_and_read_m365(
        self,
        query: str,
        __oauth_token__: Optional[dict] = None,
        __request__: Optional[Any] = None,
        __event_emitter__: Optional[Callable[[Any], Any]] = None,
    ) -> str:
        """
        Sucht in SharePoint/OneDrive nach Dokumenten und liest deren Inhalt.

        Aus dem Suchbegriff werden automatisch bis zu 5 progressive
        AND-Suchanfragen generiert – von strikt (alle Wörter) bis breit
        (ein Wort, kein Dateifilter). Die erste Stufe mit Treffern gewinnt.

        Beispiel: query = "SCOM Windows Service Monitor Anleitung"
          Stufe 1: SCOM AND Windows AND Service AND Monitor AND Anleitung (+ Dateifilter)
          Stufe 2: SCOM AND Windows AND Service AND Monitor (+ Dateifilter)
          Stufe 3: SCOM AND Windows AND Service (+ Dateifilter)
          Stufe 4: SCOM AND Windows (alle Dateitypen)
          Stufe 5: SCOM (alle Dateitypen)

        :param query: Suchbegriff(e) – kann mehrere Wörter enthalten.
                     Längere Eingaben werden automatisch progressiv gestutzt.
        """
        # ── Status-Emitter ─────────────────────────────────────────────
        async def emit_status(msg: str, done: bool = False):
            if __event_emitter__:
                await __event_emitter__(
                    {"type": "status", "data": {"description": msg, "done": done}}
                )

        # ── Auth ───────────────────────────────────────────────────────
        # Strategie:
        # 1. __oauth_token__ vorhanden + gültig  → direkt verwenden
        # 2. __oauth_token__ vorhanden + abgelaufen + refresh_token da
        #    → automatisch via Microsoft Token-Endpoint erneuern
        # 3. __oauth_token__ fehlt komplett (bekannter Open WebUI Bug)
        #    → Fallback auf Request-Cookie "oauth_id_token"
        # 4. Alles fehlt → Fehlermeldung
        token = None
        _token_source = "keiner"

        if __oauth_token__:
            _access_token  = __oauth_token__.get("access_token")
            _refresh_token = __oauth_token__.get("refresh_token")
            _expires_at    = __oauth_token__.get("expires_at", 0)
            _userinfo      = __oauth_token__.get("userinfo", {})

            _now = time.time()
            # Token noch gültig? (60 Sekunden Puffer)
            _is_expired = bool(_expires_at) and _now >= (_expires_at - 60)

            if _access_token:
                # ── Pfad 1: Token gültig ───────────────────────────
                token = _access_token
                _token_source = "__oauth_token__ (gültig)"

            elif _refresh_token:
                # ── Pfad 2: Token abgelaufen, erneuern ────────────────
                await emit_status("🔄 Access Token abgelaufen – erneuere automatisch...")
                try:
                    # client_id + tenant_id direkt aus dem Token lesen –
                    # kein manuelles Konfigurieren nötig!
                    _client_id = _userinfo.get("aud", "")
                    _tenant_id = _userinfo.get("tid", "common")
                    _scope     = __oauth_token__.get(
                        "scope", "https://graph.microsoft.com/.default offline_access"
                    )

                    _payload = {
                        "grant_type":    "refresh_token",
                        "refresh_token": _refresh_token,
                        "client_id":     _client_id,
                        "scope":         _scope,
                    }
                    if self.valves.oauth_client_secret:
                        _payload["client_secret"] = self.valves.oauth_client_secret

                    async with httpx.AsyncClient(timeout=10.0) as _rc:
                        _rr = await _rc.post(
                            f"https://login.microsoftonline.com/{_tenant_id}/oauth2/v2.0/token",
                            data=_payload,
                        )

                    if _rr.status_code == 200:
                        token = _rr.json().get("access_token")
                        _token_source = "refresh_token (auto-erneuert ✅)"
                        await emit_status("✅ Token erfolgreich erneuert.")
                    else:
                        # Refresh schlug fehl – alten Token trotzdem versuchen
                        token = _access_token
                        _token_source = f"__oauth_token__ (abgelaufen, Refresh HTTP {_rr.status_code})"

                except Exception as _e_ref:
                    token = _access_token  # letzter Ausweg
                    _token_source = f"__oauth_token__ (Refresh-Fehler: {str(_e_ref)[:60]})"

            elif _access_token:
                # ── Pfad 2b: abgelaufen, kein refresh_token ────────────
                token = _access_token
                _token_source = "__oauth_token__ (abgelaufen, kein refresh_token)"

        # ── Pfad 3: Cookie-Fallback (wenn __oauth_token__ komplett fehlt)
        if not token and __request__ is not None:
            try:
                _cookies = dict(__request__.cookies)
                for _cname in ("oauth_id_token", "id_token", "access_token", "token"):
                    _cval = _cookies.get(_cname)
                    if _cval:
                        token = _cval
                        _token_source = f"Cookie '{_cname}'"
                        break
            except Exception:
                pass

        # ── Pfad 4: Gar nichts verfügbar ────────────────────────────
        if not token:
            return "❌ Nicht über Microsoft angemeldet. Bitte via ENTRA ID einloggen."

        auth_headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        }

        # ── Queries aufbauen ───────────────────────────────────────────
        fallback_queries = _build_progressive_queries(query)
        total_stages = len(fallback_queries)

        # Zeige dem User was passiert (im Status)
        stage_preview = " → ".join(
            f"[{q['index']}] {' AND '.join(q['words'][:4])}{'...' if len(q['words']) > 4 else ''}"
            for q in fallback_queries
        )
        await emit_status(f"🧠 {total_stages} Suchstufen vorbereitet: {stage_preview}")

        async with httpx.AsyncClient(
            timeout=30.0, follow_redirects=True
        ) as client:

            # ── 5-stufige Fallback-Suche ───────────────────────────────
            hits = []
            used_query = None
            search_log = []

            for q in fallback_queries:
                words_display = " AND ".join(q["words"])
                filter_note = "+ Dateifilter" if q["with_file_filter"] else "alle Typen"
                await emit_status(
                    f"🔍 Stufe {q['index']}/{total_stages}: "
                    f"{words_display} [{filter_note}]"
                )

                payload = {
                    "requests": [
                        {
                            "entityTypes": ["driveItem"],
                            "query": {"queryString": q["kql"]},
                            "fields": [
                                "name", "id", "parentReference", "webUrl",
                                "size", "lastModifiedDateTime", "createdBy",
                                "file", "folder", "summary",
                            ],
                            # 15 Ergebnisse pro Stage (Copilot nutzt auch breiteres Fetch)
                            "size": 15,
                            "from": 0,
                        }
                    ]
                }

                try:
                    r = await client.post(
                        f"{self.base_url}/search/query",
                        headers=auth_headers,
                        json=payload,
                        timeout=15.0,
                    )

                    if r.status_code == 401:
                        return "❌ OAuth-Token abgelaufen (Files.Read.All benötigt)."
                    if r.status_code == 403:
                        return "❌ Zugriff verweigert – 'Files.Read.All' Delegated Permission fehlt."

                    r.raise_for_status()

                    raw_hits = (
                        r.json()
                        .get("value", [{}])[0]
                        .get("hitsContainers", [{}])[0]
                        .get("hits", [])
                    )

                    # Client-seitig Ordner rausfiltern & Dokumentcheck
                    # Wir behalten den Graph-API-Rank für das spätere Ranking
                    doc_hits = [
                        h for h in raw_hits
                        if not h.get("resource", {}).get("folder")
                        and _is_document(h.get("resource", {}).get("name", ""))
                    ]

                    count = len(doc_hits)
                    if count > 0:
                        search_log.append({
                            "compact": f"  ✅ Stufe {q['index']} [{q['label']}] → {count} Dokument(e)",
                            "debug":   f"  ✅ Stufe {q['index']} [{q['label']}] → {count} Dokument(e)\n"
                                       f"     📝 Query: {q['kql']}",
                        })
                        hits = doc_hits  # Treffer – Top-N werden nach API-Rank vorselektiert
                        used_query = q
                        break
                    else:
                        search_log.append({
                            "compact": f"  ❌ Stufe {q['index']} [{q['label']}] → 0 Treffer",
                            "debug":   f"  ❌ Stufe {q['index']} [{q['label']}] → 0 Treffer\n"
                                       f"     📝 Query: {q['kql']}",
                        })

                except httpx.HTTPStatusError as e:
                    search_log.append({
                        "compact": f"  ⚠️ Stufe {q['index']} [{q['label']}] → API-Fehler {e.response.status_code}",
                        "debug":   f"  ⚠️ Stufe {q['index']} [{q['label']}] → API-Fehler {e.response.status_code}\n"
                                   f"     📝 Query: {q['kql']}",
                    })
                    continue
                except Exception as e:
                    search_log.append({
                        "compact": f"  ⚠️ Stufe {q['index']} [{q['label']}] → {str(e)[:60]}",
                        "debug":   f"  ⚠️ Stufe {q['index']} [{q['label']}] → {str(e)[:60]}\n"
                                   f"     📝 Query: {q['kql']}",
                    })
                    continue

            # ── Kein Ergebnis ──────────────────────────────────────────
            debug = self.valves.debug_mode
            log_key = "debug" if debug else "compact"
            if not hits:
                log = "\n".join(entry[log_key] for entry in search_log)
                return (
                    f"ℹ️ Keine Dokumente für '{query}' gefunden ({total_stages} Stufen).\n\n"
                    f"**Suchprotokoll:**\n```\n{log}\n```\n\n"
                    f"**Tipps:**\n"
                    f"- Kürzere, spezifischere Begriffe verwenden\n"
                    f"- Prüfen ob die Dokumente in SharePoint/OneDrive indiziert sind\n"
                    f"- Neue Dokumente können bis zu 24h brauchen"
                )

            # ── Vorselektion nach API-Rank (reduziert Download-Calls) ───
            # Graph liefert hits bereits nach Relevanz sortiert (rank=1 = bester Treffer).
            # Wir laden nur die Top-N herunter – spart API-Calls ohne Qualitätsverlust.
            _max_dl = max(1, min(15, self.valves.max_downloads))
            hits_total = len(hits)
            hits_to_download = sorted(
                hits,
                key=lambda h: h.get("rank", 999)
            )[:_max_dl]

            await emit_status(
                f"📄 {hits_total} Dokument(e) via Stufe {used_query['index']} gefunden – "
                f"lade Top {len(hits_to_download)} nach API-Rank (von {hits_total}) herunter..."
            )

            # ── Download & Extraktion (parallel) ──────────────────────
            async def download_and_extract(hit: dict) -> dict:
                res = hit.get("resource", {})
                name = res.get("name", "Unbekannt")
                item_id = res.get("id")
                drive_id = res.get("parentReference", {}).get("driveId")
                web_url = res.get("webUrl", "")
                modified = res.get("lastModifiedDateTime", "")
                file_size = res.get("size", 0)
                author = (
                    res.get("createdBy", {})
                    .get("user", {})
                    .get("displayName", "Unbekannt")
                )
                # summary = kurzer Textausschnitt von Graph API (Keyword-Kontext)
                api_summary = hit.get("summary", "")
                # Graph-API-Rank: Platz 1 = relevantester Treffer laut Microsoft
                api_rank = hit.get("rank", 0)

                result = {
                    "name": name,
                    "url": web_url,
                    "type": _get_file_type(name),
                    "modified": modified[:10] if modified else "?",
                    "size_kb": round(file_size / 1024, 1) if file_size else 0,
                    "author": author,
                    "text": None,
                    "summary": api_summary,
                    "api_rank": api_rank,
                    "score": 0,
                    "error": None,
                }

                if not drive_id or not item_id:
                    result["error"] = "Keine Drive-ID/Item-ID verfügbar"
                    return result

                if file_size and file_size > 25 * 1024 * 1024:
                    result["error"] = f"Datei zu groß ({result['size_kb']} KB) – Max: 25 MB"
                    return result

                try:
                    dl_url = (
                        f"{self.base_url}/drives/{drive_id}/items/{item_id}/content"
                    )
                    dr = await client.get(dl_url, headers=auth_headers, timeout=20.0)

                    if dr.status_code == 200:
                        result["text"] = await _extract_text_from_bytes(name, dr.content)
                    elif dr.status_code == 302:
                        loc = dr.headers.get("Location")
                        if loc:
                            dr2 = await client.get(loc, timeout=20.0)
                            result["text"] = (
                                await _extract_text_from_bytes(name, dr2.content)
                                if dr2.status_code == 200
                                else None
                            )
                            if dr2.status_code != 200:
                                result["error"] = f"Redirect HTTP {dr2.status_code}"
                        else:
                            result["error"] = "Kein Redirect-Ziel"
                    else:
                        result["error"] = f"HTTP {dr.status_code}"

                except asyncio.TimeoutError:
                    result["error"] = "Timeout (>20s)"
                except Exception as e:
                    result["error"] = str(e)[:150]

                return result

            docs_all = await asyncio.gather(*[download_and_extract(h) for h in hits_to_download])

            # ── Kombiniertes Ranking (Copilot-Style) ───────────────────
            # Kombiniert drei Faktoren wie Microsoft Copilot:
            # 1. Graph-API-Rank (Microsoft's eigene Relevanz-Bewertung)
            # 2. Keyword-Density im extrahierten Text
            # 3. Aktualität (neuere Dateien werden bevorzugt – bis +20 Punkte)
            for doc in docs_all:
                doc["score"] = _compute_combined_score(
                    doc, used_query["words"], api_rank=doc.get("api_rank", 0)
                )
            docs_ranked = sorted(docs_all, key=lambda d: d["score"], reverse=True)
            top_docs = docs_ranked[:3]

            # ── Ausgabe ────────────────────────────────────────────────
            await emit_status(
                f"✅ Top {len(top_docs)} von {len(docs_all)} geladenen Dokument(en) (aus {hits_total} Treffern) nach kombiniertem Score ausgewählt.",
                done=True,
            )

            used_words = " AND ".join(used_query["words"])

            output = [
                f"## 📁 SharePoint/OneDrive: '{query}'\n",
                f"**Erfolgreich mit Stufe {used_query['index']} von {total_stages}:** `{used_words}`",
                f"*{hits_total} Treffer gefunden, {len(docs_all)} heruntergeladen (Top {_max_dl} nach API-Rank), Top {len(top_docs)} nach kombiniertem Score angezeigt.*",
            ]

            if debug:
                log_text = "\n".join(entry[log_key] for entry in search_log)
                output += [
                    "\n**🔍 Debug – Suchprotokoll (alle Stufen + KQL):**",
                    "```",
                    log_text,
                    "```",
                    "**📊 Ranking aller Treffer (Graph-API-Rank + Keyword-Dichte + Aktualität):**",
                    "```",
                    *[
                        f"  [Score:{d['score']:>5.1f} | API-Rank:{d.get('api_rank', '?'):>2} | Geändert:{d['modified']}] {d['name']}"
                        for d in docs_ranked
                    ],
                    "```\n",
                ]

            output += [
                f"\n**Top {len(top_docs)} Dokument(e):**\n",
                "---\n",
            ]

            for i, doc in enumerate(top_docs, 1):
                score_badge = f" *(Score: {doc['score']})*" if debug else ""
                output.append(f"### Dokument {i}: {doc['name']}{score_badge}")
                output.append(f"- **Typ:** {doc['type']}")
                output.append(f"- **Größe:** {doc['size_kb']} KB")
                output.append(f"- **Geändert:** {doc['modified']}")
                output.append(f"- **Autor:** {doc['author']}")
                output.append(f"- **Link:** [{doc['name']}]({doc['url']})")
                output.append("")

                if doc.get("error"):
                    output.append(f"> ⚠️ {doc['error']}")
                    # Zeige API-Summary als Fallback wenn vorhanden
                    if doc.get("summary"):
                        output.append(f"> 💡 **Vorschau:** {doc['summary']}")
                elif doc.get("text") and doc["text"].strip():
                    output.append("**Inhalt:**")
                    output.append("```")
                    output.append(doc["text"].strip())
                    output.append("```")
                elif doc.get("summary"):
                    # Kein Text extrahierbar → API-Summary als Vorschau
                    output.append("> 💡 **Vorschau (API Summary):**")
                    output.append(f"> {doc['summary']}")
                else:
                    output.append("> ℹ️ Kein Text extrahierbar.")

                output.append("")
                output.append("---")
                output.append("")

            return "\n".join(output)

    async def summarize_my_emails(
        self,
        days_back: int = 0,
        topic_filter: str = "",
        __user__: Optional[dict] = None,
        __oauth_token__: Optional[dict] = None,
        __request__: Optional[Any] = None,
        __event_emitter__: Optional[Callable[[Any], Any]] = None,
    ) -> str:
        """
        Holt E-Mails aus Outlook und erstellt eine priorisierte Zusammenfassung –
        analog zu Microsoft Copilot's täglichem Mail-Briefing.

        E-Mails werden automatisch in drei Kategorien eingeteilt:
        🔴 Handlungsrelevant – direkte Anfragen, Action-Keywords, Fragen an dich
        🟡 Informativ – relevante Informationen, passt zu deinen Themen
        ⚪ Automatisch/System – DevOps, Monitoring, Newsletter (zusammengefasst)

        Die User-E-Mail wird automatisch aus deinem Login gelesen.
        Deine Arbeitsthemen kannst du in den UserValves einstellen.

        :param days_back: Wie viele Tage zurück (0 = UserValve-Einstellung verwenden,
                         1 = nur heute, 7 = letzte Woche)
        :param topic_filter: Optionale Themen-Ergänzung für diese Anfrage
                             (z.B. "Azure, SCOM") – wird mit UserValve-Themen kombiniert
        """
        # ── Status-Emitter ──────────────────────────────────────────────
        async def emit_status(msg: str, done: bool = False):
            if __event_emitter__:
                await __event_emitter__(
                    {"type": "status", "data": {"description": msg, "done": done}}
                )

        # ── Auth ────────────────────────────────────────────────────────
        # Priorität: ① gültiger access_token  ② sofort refresh (auch wenn leer!)
        #            ③ abgelaufener token als Fallback  ④ Cookie  ⑤ Fehler
        token = None
        _refresh_error = ""

        if __oauth_token__:
            _access_token  = (__oauth_token__.get("access_token") or "").strip()
            _refresh_token = (__oauth_token__.get("refresh_token") or "").strip()
            _expires_at    = __oauth_token__.get("expires_at", 0)
            _userinfo      = __oauth_token__.get("userinfo", {})
            _now           = time.time()
            _is_expired    = bool(_expires_at) and _now >= (_expires_at - 60)

            # ① Gültiger access_token vorhanden
            if _access_token and not _is_expired:
                token = _access_token

            # ② access_token leer ODER abgelaufen → sofort Refresh, keine weiteren Bedingungen
            elif _refresh_token:
                await emit_status("🔄 Token fehlt/abgelaufen – erneuere automatisch...")
                try:
                    _client_id = _userinfo.get("aud", "")
                    _tenant_id = _userinfo.get("tid", "common")
                    _scope = (__oauth_token__.get("scope") or
                              "https://graph.microsoft.com/.default offline_access")
                    _payload = {
                        "grant_type":    "refresh_token",
                        "refresh_token": _refresh_token,
                        "client_id":     _client_id,
                        "scope":         _scope,
                    }
                    if self.valves.oauth_client_secret:
                        _payload["client_secret"] = self.valves.oauth_client_secret
                    async with httpx.AsyncClient(timeout=15.0) as _rc:
                        _rr = await _rc.post(
                            f"https://login.microsoftonline.com/{_tenant_id}/oauth2/v2.0/token",
                            data=_payload,
                        )
                    if _rr.status_code == 200:
                        token = _rr.json().get("access_token")
                        await emit_status("✅ Token erfolgreich erneuert.")
                    else:
                        _refresh_error = f"HTTP {_rr.status_code}: {_rr.text[:150]}"
                        await emit_status(f"⚠️ Refresh fehlgeschlagen ({_refresh_error}) – versuche alten Token...")
                        token = _access_token or None
                except Exception as _ex:
                    _refresh_error = str(_ex)[:150]
                    await emit_status(f"⚠️ Refresh-Fehler: {_refresh_error}")
                    token = _access_token or None

            # ③ Kein refresh_token, aber (abgelaufener) access_token vorhanden → trotzdem probieren
            elif _access_token:
                token = _access_token

        # ④ Cookie-Fallback
        if not token and __request__ is not None:
            try:
                _cookies = dict(__request__.cookies)
                for _cname in ("oauth_id_token", "id_token", "access_token", "token"):
                    _cval = _cookies.get(_cname)
                    if _cval:
                        token = _cval
                        break
            except Exception:
                pass

        if not token:
            _hint = f"\n\nRefresh-Fehler: {_refresh_error}" if _refresh_error else ""
            return (
                "❌ Kein gültiger Token verfügbar. Bitte Open WebUI neu laden "
                f"und via ENTRA ID einloggen.{_hint}"
            )

        auth_headers = {
            "Authorization": f"Bearer {token}",
            "Accept": "application/json",
        }

        # ── User-Infos aus Open WebUI ────────────────────────────────────
        # __user__ liefert automatisch: id, email, name, role
        user_email = ""
        user_name  = ""
        if __user__:
            user_email = __user__.get("email", "")
            user_name  = __user__.get("name", "")

        # Fallback: aus dem OAuth-Token (userinfo)
        if not user_email and __oauth_token__:
            _ui = __oauth_token__.get("userinfo", {})
            user_email = _ui.get("email") or _ui.get("upn") or _ui.get("preferred_username", "")
            user_name  = user_name or _ui.get("name", "")

        # ── Themen zusammenführen ────────────────────────────────────────
        _uv_topics = self.user_valves.topics_of_interest
        _all_topics_raw = ", ".join(filter(None, [_uv_topics, topic_filter]))
        user_topics = [
            t.strip() for t in _all_topics_raw.split(",") if t.strip()
        ] if _all_topics_raw else []

        # ── Zeitraum berechnen ───────────────────────────────────────────
        _days = days_back if days_back > 0 else self.user_valves.email_days_back
        _days = max(1, min(30, _days))
        since_dt = datetime.datetime.now(datetime.timezone.utc) - datetime.timedelta(days=_days)
        since_iso = since_dt.strftime("%Y-%m-%dT%H:%M:%SZ")
        since_label = "heute" if _days == 1 else f"letzten {_days} Tage"

        await emit_status(
            f"📬 Hole E-Mails der {since_label}"
            f"{f' für {user_name}' if user_name else ''}"
            f"{f' (Themen: {_all_topics_raw})' if user_topics else ''}..."
        )

        # ── Hilfsfunktion: Token per Refresh erneuern ────────────────────
        async def _try_refresh_token() -> str | None:
            """Versucht den Access Token via Refresh Token zu erneuern. Gibt neuen Token zurück oder None."""
            if not __oauth_token__:
                return None
            _rt = (__oauth_token__.get("refresh_token") or "").strip()
            if not _rt:
                return None
            try:
                _ui  = __oauth_token__.get("userinfo", {})
                _cid = _ui.get("aud", "")
                _tid = _ui.get("tid", "common")
                _sc  = (__oauth_token__.get("scope") or
                        "https://graph.microsoft.com/.default offline_access")
                _pl  = {"grant_type": "refresh_token", "refresh_token": _rt,
                        "client_id": _cid, "scope": _sc}
                if self.valves.oauth_client_secret:
                    _pl["client_secret"] = self.valves.oauth_client_secret
                async with httpx.AsyncClient(timeout=15.0) as _rc:
                    _rr = await _rc.post(
                        f"https://login.microsoftonline.com/{_tid}/oauth2/v2.0/token",
                        data=_pl,
                    )
                if _rr.status_code == 200:
                    return _rr.json().get("access_token")
            except Exception:
                pass
            return None

        # ── Graph API: Mails abrufen (mit Retry-on-401) ──────────────────
        # GET /me/messages mit Datumsfilter und relevanten Feldern
        # Mail.Read Permission wird benötigt
        _max_fetch = max(10, min(100, self.user_valves.email_max_fetch))
        params = {
            "$filter": f"receivedDateTime ge {since_iso}",
            "$orderby": "receivedDateTime desc",
            "$top": str(_max_fetch),
            "$select": (
                "id,subject,from,toRecipients,ccRecipients,"
                "receivedDateTime,isRead,importance,bodyPreview,webLink"
            ),
        }

        try:
            async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
                r = await client.get(
                    f"{self.base_url}/me/messages",
                    headers=auth_headers,
                    params=params,
                    timeout=20.0,
                )

                # ── Retry-on-401: Token lokal gültig, aber Microsoft sagt "invalid" ──
                # Tritt auf wenn expires_at nicht gesetzt (=0) und Token tatsächlich
                # abgelaufen ist (z.B. "Signing key is invalid")
                if r.status_code == 401:
                    await emit_status("🔄 Token von Microsoft abgelehnt – erneuere und versuche erneut...")
                    _new_token = await _try_refresh_token()
                    if _new_token:
                        await emit_status("✅ Token erneuert – wiederhole Anfrage...")
                        # Retry mit neuem Token
                        r = await client.get(
                            f"{self.base_url}/me/messages",
                            headers={**auth_headers, "Authorization": f"Bearer {_new_token}"},
                            params=params,
                            timeout=20.0,
                        )

                # Nach eventuellem Retry: Fehler prüfen
                if r.status_code == 401:
                    return (
                        "❌ Graph API: Zugriff verweigert (HTTP 401) – auch nach Token-Refresh.\n"
                        "Bitte Open WebUI neu laden und via Microsoft neu anmelden.\n"
                        f"API-Details: {r.text[:300]}"
                    )
                if r.status_code == 403:
                    return (
                        "❌ Zugriff verweigert – 'Mail.Read' Delegated Permission fehlt.\n"
                        "Bitte in der Azure App Registration unter API Permissions hinzufügen."
                    )
                r.raise_for_status()

                messages = r.json().get("value", [])

        except httpx.HTTPStatusError as e:
            return f"❌ Graph API Fehler: HTTP {e.response.status_code}\n{e.response.text[:300]}"
        except Exception as e:
            return f"❌ Verbindungsfehler: {str(e)[:200]}"

        if not messages:
            return (
                f"ℹ️ Keine E-Mails für die {since_label} gefunden.\n"
                f"Filter: receivedDateTime >= {since_iso}"
            )

        await emit_status(
            f"🧠 {len(messages)} E-Mails geladen – bewerte nach Copilot-Logik..."
        )

        # ── Scoring ───────────────────────────────────────────────
        scored = []
        user_email_lower = user_email.lower()
        for msg in messages:
            sender_addr = (msg.get("from", {}).get("emailAddress", {}).get("address") or "").lower()
            # Eigene versendete Mails (übergelaufene Replies im Inbox) herausfiltern
            if sender_addr and sender_addr == user_email_lower:
                continue
            score, reasons, is_system = _score_email(msg, user_email, user_topics)
            scored.append({
                "msg": msg,
                "score": score,
                "reasons": reasons,
                "is_system": is_system,
            })

        # Sortieren nach Score (absteigend), dann nach Datum (neueste zuerst)
        scored.sort(key=lambda x: (x["score"], x["msg"].get("receivedDateTime", "")), reverse=True)

        # ── Thread-Grouping ─────────────────────────────────────────
        # Betreff-Präfixe entfernen: AW, WG, Re, Fwd, EW etc.
        _prefix_re = re.compile(
            r'^\s*(?:(?:AW|WG|Re|Fwd?|EW|SV|TR|FWD)\s*(?:\[\d+\])?:\s*)+',
            re.IGNORECASE,
        )
        def _norm_subject(subj: str) -> str:
            return _prefix_re.sub("", subj).strip().lower()

        # Pro normalisiertem Betreff: höchst-scorendes (= erstes nach sort) Element beibehalten,
        # thread_count mitzählen
        seen_threads: dict[str, dict] = {}
        for entry in scored:
            ns = _norm_subject(entry["msg"].get("subject") or "")
            if ns not in seen_threads:
                seen_threads[ns] = {**entry, "thread_count": 1}
            else:
                seen_threads[ns]["thread_count"] += 1
                # Neuestes Datum für Link-Anzeige merken (scored ist bereits score-sortiert,
                # aber wir wollen bei gleichen Threads das späteste Datum)
                existing_dt = seen_threads[ns]["msg"].get("receivedDateTime", "")
                new_dt = entry["msg"].get("receivedDateTime", "")
                if new_dt > existing_dt:
                    seen_threads[ns]["msg"] = entry["msg"]

        scored = list(seen_threads.values())
        # Nochmal sortieren nach Score desc, dann Datum desc
        scored.sort(key=lambda x: (x["score"], x["msg"].get("receivedDateTime", "")), reverse=True)

        # ── Kategorisierung (Copilot-Style) ───────────────────────────
        # WICHTIG: System-Mails IMMER in ⧁, egal wie hoch der Score ist!
        # (z.B. Azure DevOps PR: direkt+5, action+4, ungelesen+1, system-4 = 6
        # → Wäre ohne is_system-Flag fälschlicherweise in 🔴 gelandet!)
        action_items  = [e for e in scored if e["score"] >= 5 and not e["is_system"]]
        info_items    = [e for e in scored if 1 <= e["score"] < 5 and not e["is_system"]]
        system_items  = [e for e in scored if e["is_system"] or e["score"] < 1]

        debug = self.valves.debug_mode

        await emit_status(
            f"✅ {len(action_items)} handlungsrelevant, {len(info_items)} informativ, "
            f"{len(system_items)} automatisch/System.",
            done=True,
        )

        # ── Ausgabe aufbauen ─────────────────────────────────────────────
        header_name = f" für **{user_name}**" if user_name else ""
        header_period = since_label
        header_topics = f" | Themen: *{_all_topics_raw}*" if user_topics else ""

        output = [
            # ── LLM-Anweisung ──────────────────────────────────────────────
            "> **📌 Anweisung an den Assistenten:** "
            "Gib diese Tool-Ausgabe **vollständig und unverändert** weiter – kürze oder paraphrasiere **nicht**. "
            "Zeige alle `[↗ In Outlook öffnen]`-Links als **klickbare Markdown-Links** an. "
            "Falls der User mehr Details zu einer Mail möchte, nutze das Tool `search_email_detail`.\n",
            f"## 📬 E-Mail-Briefing{header_name} – {header_period}{header_topics}\n",
            f"*{len(messages)} Mails geladen · {_days} Tag(e) · Scoring: Direkt+5, Action+4, Frage+3, Thema+2, Ungelesen+1, CC-1, System-4*\n",
            "---\n",
        ]

        # ── 🔴 HANDLUNGSRELEVANT ─────────────────────────────────────────
        output.append("### 🔴 Handlungsrelevant")
        if action_items:
            for e in action_items:
                msg        = e["msg"]
                subj       = msg.get("subject") or "(kein Betreff)"
                sender     = msg.get("from", {}).get("emailAddress", {})
                from_str   = f"{sender.get('name', '')} <{sender.get('address', '')}>"
                received   = (msg.get("receivedDateTime") or "")[:16].replace("T", " ")
                link       = msg.get("webLink", "")
                unread     = "🔵 " if not msg.get("isRead", True) else ""
                score_str  = f" *(Score: {e['score']}, {', '.join(e['reasons'])})*" if debug else ""
                tc         = e.get("thread_count", 1)
                thread_str = f" *({tc} Mails im Thread)*" if tc > 1 else ""

                # Kurze Vorschau: 1. Satz, max 100 Zeichen
                raw_preview = (msg.get("bodyPreview") or "").replace("\r\n", " ").replace("\n", " ").strip()
                first_dot   = raw_preview.find(". ")
                short_prev  = raw_preview[:first_dot + 1] if 0 < first_dot < 100 else raw_preview[:100]
                if len(raw_preview) > len(short_prev):
                    short_prev += " …"

                output.append(f"\n#### {unread}{subj}{score_str}{thread_str}")
                output.append(f"- **Von:** {from_str} · {received}")
                if short_prev:
                    output.append(f"- **Vorschau:** {short_prev}")
                if link:
                    output.append(f"- **[↗ In Outlook öffnen]({link})**")
        else:
            output.append("*Keine direkt handlungsrelevanten Mails.*")

        output.append("\n---\n")

        # ── 🟡 INFORMATIV ────────────────────────────────────────────────
        output.append("### 🟡 Informativ")
        if info_items:
            for e in info_items[:10]:
                msg        = e["msg"]
                subj       = msg.get("subject") or "(kein Betreff)"
                sender     = msg.get("from", {}).get("emailAddress", {})
                from_name  = sender.get("name") or sender.get("address", "?")
                received   = (msg.get("receivedDateTime") or "")[:16].replace("T", " ")
                unread     = "🔵 " if not msg.get("isRead", True) else ""
                score_str  = f" *(Score: {e['score']}, {', '.join(e['reasons'])})*" if debug else ""
                link       = msg.get("webLink", "")
                tc         = e.get("thread_count", 1)
                thread_str = f" *({tc}×)*" if tc > 1 else ""

                output.append(f"- {unread}**{subj}**{score_str}{thread_str}  ")
                output.append(f"  *{from_name}*, {received} – [↗ In Outlook öffnen]({link})" if link else f"  *{from_name}*, {received}")

            if len(info_items) > 10:
                output.append(f"\n*… und {len(info_items) - 10} weitere informative Threads.*")
        else:
            output.append("*Keine informativen Mails.*")

        output.append("\n---\n")

        # ── ⚪ AUTOMATISCH ────────────────────────────────────────────────
        output.append("### ⚪ Automatisch / System")
        if system_items:
            for e in system_items[:5]:
                msg        = e["msg"]
                subj       = msg.get("subject") or "(kein Betreff)"
                sender     = msg.get("from", {}).get("emailAddress", {})
                from_name  = sender.get("name") or sender.get("address", "?")
                link       = msg.get("webLink", "")
                link_str   = f" [↗]({link})" if link else ""
                score_str  = f" *(Score: {e['score']})*" if debug else ""
                tc         = e.get("thread_count", 1)
                thread_str = f" *({tc}×)*" if tc > 1 else ""
                output.append(f"- {subj}{score_str}{thread_str} – *{from_name}*{link_str}")
            if len(system_items) > 5:
                output.append(f"- *… und {len(system_items) - 5} weitere System-Threads (nicht einzeln aufgeführt).*")
        else:
            output.append("*Keine System-Mails.*")

        output.append("\n---\n")
        output.append(
            f"> 💡 **Tipp:** Stelle in den UserValves deine Arbeitsthemen ein "
            f"(z.B. `PowerAutomate, KI, ISMS`) für besseres Scoring."
        )

        return "\n".join(output)


    async def search_email_detail(
        self,
        query: str,
        max_results: int = 3,
        __user__: Optional[dict] = None,
        __oauth_token__: Optional[dict] = None,
        __request__: Optional[Any] = None,
        __event_emitter__: Optional[Callable[[Any], Any]] = None,
    ) -> str:
        """
        Sucht nach E-Mails und zeigt den vollständigen Inhalt direkt im Chat an –
        ohne dass der User den Outlook-Link öffnen muss.

        Durchsucht Betreff, Body und Absender gleichzeitig (Microsoft Graph Volltext-Suche).
        Für jede gefundene Mail wird der komplette Inhalt als lesbarer Text ausgegeben.

        Beispielaufrufe:
        - "Zeig mir die Mail von Bauernberger über Anlagen"
        - "Was steht in der Mail zum GMP-Statusupdate?"
        - "Detail zur Mail von Sonja über Zielvereinbarungen"

        :param query: Suchbegriff – Betreff, Absender oder Inhalt (alles durchsuchbar)
        :param max_results: Wie viele Mails anzeigen (1–5, Standard: 3)
        """
        async def emit_status(msg: str, done: bool = False):
            if __event_emitter__:
                await __event_emitter__(
                    {"type": "status", "data": {"description": msg, "done": done}}
                )

        # ── Auth (identisch zu summarize_my_emails) ──────────────────────
        token = None
        _refresh_error = ""

        if __oauth_token__:
            _access_token  = (__oauth_token__.get("access_token") or "").strip()
            _refresh_token = (__oauth_token__.get("refresh_token") or "").strip()
            _expires_at    = __oauth_token__.get("expires_at", 0)
            _userinfo      = __oauth_token__.get("userinfo", {})
            _now           = time.time()
            _is_expired    = bool(_expires_at) and _now >= (_expires_at - 60)

            if _access_token and not _is_expired:
                token = _access_token
            elif _refresh_token:
                await emit_status("🔄 Token erneuern...")
                try:
                    _pl = {
                        "grant_type":    "refresh_token",
                        "refresh_token": _refresh_token,
                        "client_id":     _userinfo.get("aud", ""),
                        "scope":         (__oauth_token__.get("scope") or
                                         "https://graph.microsoft.com/.default offline_access"),
                    }
                    if self.valves.oauth_client_secret:
                        _pl["client_secret"] = self.valves.oauth_client_secret
                    async with httpx.AsyncClient(timeout=15.0) as _rc:
                        _rr = await _rc.post(
                            f"https://login.microsoftonline.com/{_userinfo.get('tid', 'common')}/oauth2/v2.0/token",
                            data=_pl,
                        )
                    if _rr.status_code == 200:
                        token = _rr.json().get("access_token")
                        await emit_status("✅ Token erneuert.")
                    else:
                        _refresh_error = f"HTTP {_rr.status_code}"
                        token = _access_token or None
                except Exception as _ex:
                    _refresh_error = str(_ex)[:100]
                    token = _access_token or None
            elif _access_token:
                token = _access_token

        if not token and __request__ is not None:
            try:
                for _cn in ("oauth_id_token", "id_token", "access_token", "token"):
                    _cv = dict(__request__.cookies).get(_cn)
                    if _cv:
                        token = _cv
                        break
            except Exception:
                pass

        if not token:
            return "❌ Nicht über Microsoft angemeldet. Bitte via ENTRA ID einloggen."

        auth_headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
        _max = max(1, min(5, max_results))

        # ── Hilfsfunktion: HTML → lesbarer Text ─────────────────────────
        def _html_to_text(html_str: str) -> str:
            """Entfernt HTML-Tags, normalisiert Whitespace."""
            # Block-Elemente → Zeilenumbruch
            text = re.sub(r'<(?:br|p|div|tr|li|h[1-6])[^>]*>', '\n', html_str, flags=re.IGNORECASE)
            # Alle anderen Tags entfernen
            text = re.sub(r'<[^>]+>', '', text)
            # HTML-Entities
            text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')\
                       .replace('&amp;', '&').replace('&quot;', '"').replace('&#39;', "'")
            # Mehrfache Leerzeilen reduzieren
            text = re.sub(r'\n{3,}', '\n\n', text)
            return text.strip()

        # ── Retry-on-401 Hilfsfunktion ───────────────────────────────────
        async def _refresh_and_get_token() -> str | None:
            if not __oauth_token__:
                return None
            _rt = (__oauth_token__.get("refresh_token") or "").strip()
            if not _rt:
                return None
            try:
                _ui = __oauth_token__.get("userinfo", {})
                _pl = {
                    "grant_type": "refresh_token", "refresh_token": _rt,
                    "client_id": _ui.get("aud", ""),
                    "scope": (__oauth_token__.get("scope") or
                              "https://graph.microsoft.com/.default offline_access"),
                }
                if self.valves.oauth_client_secret:
                    _pl["client_secret"] = self.valves.oauth_client_secret
                async with httpx.AsyncClient(timeout=15.0) as _rc:
                    _rr = await _rc.post(
                        f"https://login.microsoftonline.com/{_ui.get('tid', 'common')}/oauth2/v2.0/token",
                        data=_pl,
                    )
                if _rr.status_code == 200:
                    return _rr.json().get("access_token")
            except Exception:
                pass
            return None

        await emit_status(f"🔍 Suche Mails: \"{query}\"...")

        # ── Schritt 1: Volltextsuche via $search ─────────────────────────
        # Graph API $search durchsucht Betreff, Body & Absender gleichzeitig
        # HINWEIS: $orderby ist mit $search nicht kombinierbar → client-seitig sortieren
        search_params = {
            "$search": f'"{query}"',
            "$top": str(_max * 3),   # Mehr holen damit client-seitig sortieren sinnvoll ist
            "$select": "id,subject,from,toRecipients,receivedDateTime,isRead,webLink,importance",
        }

        try:
            async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:

                r = await client.get(
                    f"{self.base_url}/me/messages",
                    headers=auth_headers,
                    params=search_params,
                    timeout=20.0,
                )

                # Retry-on-401
                if r.status_code == 401:
                    await emit_status("🔄 Token abgelaufen – erneuere...")
                    _nt = await _refresh_and_get_token()
                    if _nt:
                        auth_headers["Authorization"] = f"Bearer {_nt}"
                        r = await client.get(
                            f"{self.base_url}/me/messages",
                            headers=auth_headers,
                            params=search_params,
                            timeout=20.0,
                        )

                if r.status_code == 401:
                    return "❌ Zugriff verweigert. Bitte Open WebUI neu laden und erneut anmelden."
                if r.status_code == 403:
                    return "❌ Mail.Read Permission fehlt. Bitte in Azure App Registration hinzufügen."
                r.raise_for_status()

                candidates = r.json().get("value", [])

                if not candidates:
                    # Fallback: $filter auf Betreff
                    await emit_status(f"ℹ️ Keine $search-Treffer – versuche Betreff-Filter...")
                    fallback_params = {
                        "$filter": f"contains(subject, '{query}')",
                        "$top": str(_max),
                        "$select": "id,subject,from,toRecipients,receivedDateTime,isRead,webLink",
                        "$orderby": "receivedDateTime desc",
                    }
                    r2 = await client.get(
                        f"{self.base_url}/me/messages",
                        headers=auth_headers,
                        params=fallback_params,
                        timeout=20.0,
                    )
                    if r2.status_code == 200:
                        candidates = r2.json().get("value", [])

                if not candidates:
                    return (
                        f"ℹ️ Keine Mails gefunden für: **\"{query}\"**\n\n"
                        "Tipps:\n"
                        "- Verwende Schlüsselwörter aus dem Betreff oder Absendernamen\n"
                        "- Für aktuelle Mails: `summarize_my_emails` verwenden"
                    )

                # Client-seitig nach Datum sortieren (neueste zuerst) und auf _max begrenzen
                candidates.sort(
                    key=lambda m: m.get("receivedDateTime", ""), reverse=True
                )
                candidates = candidates[:_max]

                await emit_status(
                    f"📨 {len(candidates)} Treffer – lade vollständigen Inhalt..."
                )

                # ── Schritt 2: Vollständigen Body für jeden Treffer holen ──
                output = [
                    # ── LLM-Anweisung ──────────────────────────────────────────────
                    "> **📌 Anweisung an den Assistenten:** "
                    "Gib diesen vollständigen Mail-Inhalt **unverändert** aus. "
                    "Zeige den `[↗ In Outlook öffnen]`-Link als **klickbaren Markdown-Link** an. "
                    "Wenn der User eine Antwort verfassen möchte: schreibe den Entwurf, dann rufe "
                    "`generate_reply_mailto` auf mit: "
                    "to_email=Absender-E-Mail, subject=Betreff, body=dein Entwurf, "
                    "original_from='Name <email>' (Von-Feld der Mail), "
                    "original_date=Empfangsdatum, original_to=An-Empfänger, "
                    "original_body=Mail-Text (aus dem Inhalt-Block oben). "
                    "Für eine neue Mail: generate_reply_mailto ohne original_*-Parameter aufrufen.\n",
                    f"## 🔍 E-Mail Detail: \"{query}\"\n",
                    f"*{len(candidates)} Treffer gefunden – vollständiger Inhalt:*\n",
                    "---\n",
                ]

                for i, msg in enumerate(candidates, 1):
                    msg_id   = msg.get("id", "")
                    subject  = msg.get("subject") or "(kein Betreff)"
                    sender   = msg.get("from", {}).get("emailAddress", {})
                    from_str = f"{sender.get('name', '')} <{sender.get('address', '')}>"
                    received = (msg.get("receivedDateTime") or "")[:16].replace("T", " ")
                    link     = msg.get("webLink", "")
                    unread   = "🔵 " if not msg.get("isRead", True) else ""

                    output.append(f"### {unread}Mail {i}: {subject}")
                    output.append(f"- **Von:** {from_str}")
                    output.append(f"- **Erhalten:** {received}")
                    if link:
                        output.append(f"- **[↗ In Outlook öffnen]({link})**")
                    # Metadaten für LLM (generate_reply_mailto Parameter)
                    output.append(
                        f"<!-- reply_meta: to={sender.get('address','')}, "
                        f"web_link={link}, subject={subject}, "
                        f"from_str={from_str}, received={received} -->"
                    )
                    output.append("")

                    # Vollständigen Body abrufen
                    if msg_id:
                        try:
                            rb = await client.get(
                                f"{self.base_url}/me/messages/{msg_id}",
                                headers=auth_headers,
                                params={"$select": "body"},
                                timeout=15.0,
                            )
                            if rb.status_code == 200:
                                body_obj     = rb.json().get("body", {})
                                content_type = body_obj.get("contentType", "text")
                                content      = body_obj.get("content", "")

                                if content_type == "html":
                                    body_text = _html_to_text(content)
                                else:
                                    body_text = content.strip()

                                # Zitat-Trenner erkennen und kürzen
                                # (meistens beginnt das Zitat mit "Von: ..." oder "-----")
                                _quote_markers = [
                                    "\nVon: ", "\nFrom: ", "\n-----Original",
                                    "\n________________________________",
                                    "\n------ Weitergeleitete Nachricht",
                                ]
                                for _qm in _quote_markers:
                                    _qi = body_text.find(_qm)
                                    if 0 < _qi < len(body_text) - 50:
                                        body_text = body_text[:_qi].strip() + "\n\n*[… zitierter Vorgänger-Thread gekürzt]*"
                                        break

                                output.append("**Inhalt:**")
                                output.append("```")
                                output.append(body_text[:4000])  # Max 4000 Zeichen
                                if len(body_text) > 4000:
                                    output.append(f"\n[… {len(body_text) - 4000} weitere Zeichen – öffne den Link für den vollen Text]")
                                output.append("```")
                            else:
                                output.append(f"> ⚠️ Body konnte nicht geladen werden (HTTP {rb.status_code})")
                        except Exception as _be:
                            output.append(f"> ⚠️ Body-Fehler: {str(_be)[:100]}")
                    else:
                        output.append("> ⚠️ Keine Message-ID für Body-Abruf verfügbar.")

                    output.append("\n---\n")

                await emit_status(
                    f"✅ {len(candidates)} Mail(s) vollständig geladen.", done=True
                )
                return "\n".join(output)

        except httpx.HTTPStatusError as e:
            return f"❌ Graph API Fehler: HTTP {e.response.status_code}\n{e.response.text[:300]}"
        except Exception as e:
            return f"❌ Verbindungsfehler: {str(e)[:200]}"

    def generate_reply_mailto(
        self,
        to_email: str,
        subject: str,
        body: str,
        cc_email: str = "",
        original_from: str = "",
        original_date: str = "",
        original_to: str = "",
        original_body: str = "",
    ) -> str:
        """
        Erstellt einen mailto:-Link der Outlook öffnet – als Antwort (mit Original-Mail
        darunter wie bei Outlook "Antworten") oder als neue Mail.

        Dieses Tool nach JEDEM verfassten Mail-Entwurf aufrufen!

        Für eine ANTWORT: original_from, original_date, original_to und original_body
        aus den Mail-Metadaten befüllen → Outlook öffnet sich mit Entwurf und
        Original-Mail darunter, genau wie "Antworten" in Outlook.

        Für eine NEUE MAIL: original_* Parameter weglassen.

        :param to_email:      Empfänger-E-Mail-Adresse
        :param subject:       Betreff – ohne AW: übergeben, wird bei Antwort automatisch gesetzt
        :param body:          Verfasster Entwurf-Text
        :param cc_email:      Optional: CC (kommagetrennt)
        :param original_from: Absender der Original-Mail, z.B. "Max Muster <max@firma.at>"
        :param original_date: Datum der Original-Mail, z.B. "2026-02-24 13:57"
        :param original_to:   Empfänger der Original-Mail
        :param original_body: Text der Original-Mail (wird als Zitat darunter eingefügt)
        """
        is_reply = bool(original_from or original_body)

        # ── Markdown → Plain Text für Outlook ───────────────────────────
        # Outlook rendert kein Markdown – **fett** würde literal erscheinen.
        # Chat-Anzeige bleibt weiterhin Markdown, nur der mailto-Body wird konvertiert.
        def _md_to_plain(md: str) -> str:
            t = md
            # Code-Blöcke zuerst (vor einzelnen Backticks)
            t = re.sub(r'```[^\n]*\n(.*?)```', r'\1', t, flags=re.DOTALL)
            # Überschriften: ### Titel → TITEL (Großbuchstaben als Ersatz)
            t = re.sub(r'^#{1,6}\s+(.+)$', lambda m: m.group(1).upper(), t, flags=re.MULTILINE)
            # Fett+Kursiv: ***text*** → text
            t = re.sub(r'\*{3}(.+?)\*{3}', r'\1', t)
            # Fett: **text** → text
            t = re.sub(r'\*{2}(.+?)\*{2}', r'\1', t)
            # Kursiv: *text* → text
            t = re.sub(r'\*(.+?)\*', r'\1', t)
            # Unterstrichen/Bold: __text__, _text_
            t = re.sub(r'__(.+?)__', r'\1', t)
            t = re.sub(r'_(.+?)_',   r'\1', t)
            # Inline-Code: `code` → code
            t = re.sub(r'`(.+?)`', r'\1', t)
            # Links: [Text](URL) → Text (URL)
            t = re.sub(r'\[(.+?)\]\((.+?)\)', r'\1 (\2)', t)
            # Horizontale Linien entfernen
            t = re.sub(r'^[-*_]{3,}\s*$', '', t, flags=re.MULTILINE)
            # Mehrfache Leerzeilen auf max. 2 reduzieren
            t = re.sub(r'\n{3,}', '\n\n', t)
            return t.strip()

        # ── Betreff ─────────────────────────────────────────────────────
        _subj = subject.strip()
        if is_reply and not re.match(r'^(?:AW|RE|WG|FW)\s*:', _subj, re.IGNORECASE):
            _subj = f"AW: {_subj}"
        if not is_reply:
            _subj = re.sub(r'^(?:AW|RE|WG|FW)\s*:\s*', '', _subj, flags=re.IGNORECASE).strip()
        _orig_subj = re.sub(r'^(?:AW|RE|WG|FW)\s*:\s*', '', _subj, flags=re.IGNORECASE).strip()

        # Entwurf: Markdown für Chat-Anzeige, Plain für mailto
        _reply_md   = body.strip()          # bleibt Markdown (für Chat-Blockquote)
        _reply_plain = _md_to_plain(body)   # Plain Text für mailto-Body

        # ── Mailto-Body: Entwurf + Outlook-Style Zitat ───────────────────
        _reply_text = _reply_plain
        _mailto_body = _reply_text

        if is_reply:
            # Budget: Ziel <1500 Zeichen RAW (nach URL-Encoding max ~4500)
            _budget = max(200, 1400 - len(_reply_text))
            _orig_preview = (original_body or "").strip()
            if len(_orig_preview) > _budget:
                _orig_preview = _orig_preview[:_budget] + "\r\n[...]"

            _quote_parts = ["", "\r\n________________________________"]
            if original_from:  _quote_parts.append(f"Von: {original_from}")
            if original_date:  _quote_parts.append(f"Gesendet: {original_date}")
            _quote_parts.append(f"An: {original_to}" if original_to else f"An: {to_email}")
            _quote_parts.append(f"Betreff: {_orig_subj}")
            _quote_parts.append("")
            _quote_parts.append(_orig_preview)

            _mailto_body = _reply_text + "\r\n".join(_quote_parts)

        # ── Mailto-URL ───────────────────────────────────────────────────
        _params = [f"subject={quote(_subj)}"]
        if cc_email.strip():
            _params.append(f"cc={quote(cc_email.strip())}")
        if _mailto_body:
            _params.append(f"body={quote(_mailto_body)}")
        mailto_url = f"mailto:{to_email}?{'&'.join(_params)}"

        # ── Chat-Ausgabe ─────────────────────────────────────────────────
        _icon = "↩️ Antwort" if is_reply else "✉️ Neue Mail"

        lines = [
            "",
            "---",
            f"### {_icon}",
            "",
            "| | |",
            "|---|---|",
            f"| **An** | {to_email} |",
        ]
        if cc_email.strip():
            lines.append(f"| **CC** | {cc_email.strip()} |")
        lines.append(f"| **Betreff** | {_subj} |")
        lines.append("")

        # Entwurf als Blockquote im Chat – Markdown bleibt erhalten (wird gerendert)
        for _line in _reply_md.splitlines():
            lines.append(f"> {_line}" if _line.strip() else ">")

        # Original darunter – wie in Outlook
        if is_reply:
            lines += [
                ">",
                "> ________________________________",
                f"> Von: {original_from}"    if original_from else "",
                f"> Gesendet: {original_date}" if original_date else "",
                f"> An: {original_to}"       if original_to else f"> An: {to_email}",
                f"> Betreff: {_orig_subj}",
                ">",
            ]
            for _ol in (original_body or "").strip()[:600].splitlines():
                lines.append(f"> {_ol}" if _ol.strip() else ">")
            if len(original_body or "") > 600:
                lines.append("> [...]")

        # Link – kein ** drum herum (bricht Markdown-Parser in Open WebUI)
        _link_label = "In Outlook antworten" if is_reply else "In Outlook schreiben"
        _hint = (
            "Öffnet Outlook mit Entwurf und Original-Mail vorausgefüllt."
            if is_reply else
            "Öffnet Outlook mit vorausgefülltem Empfänger und Betreff."
        )

        lines += [
            "",
            "---",
            f"[↗ {_link_label}]({mailto_url})",
            "",
            f"*{_hint}*",
        ]

        return "\n".join(l for l in lines if l is not None)

    async def get_calendar_events(
        self,
        date_input: str = "morgen",
        days: int = 1,
        __user__: Optional[dict] = None,
        __oauth_token__: Optional[dict] = None,
        __request__: Optional[Any] = None,
        __event_emitter__: Optional[Callable[[Any], Any]] = None,
    ) -> str:
        """
        Ruft Kalendertermine über Microsoft Graph ab – wie Microsoft Copilot.

        Perfekt für:
        - Tagesvorbereitung: "Was habe ich morgen?"
        - Rückblick: "Was war gestern los? Fasse meinen Tag zusammen."
        - Wochenplanung: "Zeig mir diese Woche."

        Gibt eine strukturierte Übersicht zurück:
        - Alle Termine mit Uhrzeit, Ort, Teilnehmern
        - Für vergangene Tage: Anweisung zur Zusammenfassung im Copilot-Stil

        :param date_input: "heute", "morgen", "gestern", "übermorgen",
                           "diese woche", "nächste woche", "letzte woche",
                           oder ISO-Datum wie "2026-02-25"
        :param days:       Anzahl Tage (1–7). Nur wenn date_input ein konkretes Datum ist.
        """
        from datetime import datetime as _dt, timedelta as _td, date as _date_t

        _DE_DAYS   = ["Montag","Dienstag","Mittwoch","Donnerstag","Freitag","Samstag","Sonntag"]
        _DE_MONTHS = ["","Januar","Februar","März","April","Mai","Juni",
                      "Juli","August","September","Oktober","November","Dezember"]

        async def emit_status(msg: str, done: bool = False):
            if __event_emitter__:
                await __event_emitter__({"type": "status", "data": {"description": msg, "done": done}})

        # ── Datum parsen ──────────────────────────────────────────────────
        today = _dt.now().date()
        _raw  = date_input.strip().lower()
        _days = max(1, min(days, 7))

        if _raw in ("heute", "today", "h"):
            start_date = today;            _days = 1; label = "heutigen Tag"
        elif _raw in ("morgen", "tomorrow", "m"):
            start_date = today + _td(1);   _days = 1; label = "morgigen Tag"
        elif _raw in ("gestern", "yesterday", "g"):
            start_date = today - _td(1);   _days = 1; label = "gestrigen Tag"
        elif _raw in ("übermorgen", "uberm"):
            start_date = today + _td(2);   _days = 1; label = "übermorgigen Tag"
        elif "nächste woche" in _raw or "next week" in _raw:
            start_date = today - _td(today.weekday()) + _td(7)
            _days = 7;                                  label = "nächste Woche"
        elif "letzte woche" in _raw or "last week" in _raw:
            start_date = today - _td(today.weekday()) - _td(7)
            _days = 7;                                  label = "letzte Woche"
        elif "woche" in _raw or "week" in _raw:
            start_date = today - _td(today.weekday())
            _days = 7;                                  label = "diese Woche"
        else:
            try:
                start_date = _dt.strptime(date_input.strip(), "%Y-%m-%d").date()
                label = start_date.strftime("%d.%m.%Y")
            except ValueError:
                return (
                    f"❌ Ungültiges Datum: '{date_input}'.\n"
                    "Verwende 'heute', 'morgen', 'gestern', 'diese woche' oder 'YYYY-MM-DD'."
                )

        end_date = start_date + _td(_days)
        is_past  = end_date <= today

        await emit_status(f"📅 Lade Kalender für {label}...")

        # ── Auth: gleiche Logik wie summarize_my_emails ───────────────────
        _access_token  = ""
        _refresh_token = ""
        _userinfo      = {}

        if __oauth_token__:
            _access_token  = (__oauth_token__.get("access_token")  or "").strip()
            _refresh_token = (__oauth_token__.get("refresh_token") or "").strip()
            _userinfo      = __oauth_token__.get("userinfo", {}) or {}

        if not _access_token and not _refresh_token:
            return "❌ Kein OAuth-Token – bitte in Open WebUI mit Microsoft anmelden."

        # client_id + tenant_id aus dem Token (wie in den anderen Funktionen)
        _client_id     = _userinfo.get("aud", "")
        _tenant_id     = _userinfo.get("tid", "common")
        _client_secret = (self.valves.oauth_client_secret or "").strip()
        _scope         = (
            __oauth_token__.get("scope")
            if __oauth_token__ else None
        ) or "https://graph.microsoft.com/.default offline_access"

        async def _try_refresh() -> str | None:
            if not (_refresh_token and _client_id and _tenant_id):
                return None
            try:
                _payload = {
                    "grant_type":    "refresh_token",
                    "refresh_token": _refresh_token,
                    "client_id":     _client_id,
                    "scope":         _scope,
                }
                if _client_secret:
                    _payload["client_secret"] = _client_secret
                async with httpx.AsyncClient() as _rc:
                    _r = await _rc.post(
                        f"https://login.microsoftonline.com/{_tenant_id}/oauth2/v2.0/token",
                        data=_payload,
                        timeout=15.0,
                    )
                    if _r.status_code == 200:
                        return _r.json().get("access_token", "")
                    await emit_status(f"⚠️ Token-Refresh fehlgeschlagen (HTTP {_r.status_code})")
                    return None
            except Exception as _e:
                await emit_status(f"⚠️ Refresh-Fehler: {str(_e)[:80]}")
                return None

        # ── Graph API: calendarView ──────────────────────────────────────
        _start_str = f"{start_date}T00:00:00"
        _end_str   = f"{end_date}T00:00:00"

        _params = {
            "startDateTime": _start_str,
            "endDateTime":   _end_str,
            "$select": (
                "subject,start,end,location,organizer,attendees,"
                "bodyPreview,isAllDay,isCancelled,showAs,webLink,"
                "isOrganizer,onlineMeetingUrl"
            ),
            "$orderby": "start/dateTime",
            "$top": "50",
        }

        try:
            async with httpx.AsyncClient() as client:

                async def _fetch(token: str):
                    return await client.get(
                        f"{self.base_url}/me/calendarView",
                        headers={
                            "Authorization": f"Bearer {token}",
                            "Prefer": 'outlook.timezone="W. Europe Standard Time"',
                        },
                        params=_params,
                        timeout=20.0,
                    )

                resp = await _fetch(_access_token)

                # Retry-on-401
                if resp.status_code == 401:
                    await emit_status("🔄 Token abgelaufen – erneuere...")
                    _new = await _try_refresh()
                    if _new:
                        _access_token = _new
                        resp = await _fetch(_access_token)

                if resp.status_code != 200:
                    return (
                        f"❌ Graph API: HTTP {resp.status_code}\n"
                        f"{resp.text[:300]}"
                    )

                events = resp.json().get("value", [])

            # ── Termine verarbeiten ───────────────────────────────────────
            await emit_status(f"✅ {len(events)} Termine geladen.", done=True)

            # LLM-Instruktion oben
            is_week = _days > 1
            if is_past and not is_week:
                _summary_task = (
                    "Fasse den Tag danach kompakt zusammen "
                    "(Vormittag/Mittag/Nachmittag, Highlights, Gesamteindruck)."
                )
            elif not is_week:
                _summary_task = (
                    "Füge danach eine kurze Tagesvorbereitung hinzu: "
                    "wichtigste Termine, Empfehlungen, Pufferzeiten."
                )
            else:
                _summary_task = (
                    "Fasse die Woche danach kompakt zusammen: "
                    "Verteilung der Termine pro Tag, wichtigste Events."
                )

            output = [
                "> **📌 Anweisung an den Assistenten:** "
                "Gib die folgende Kalenderübersicht **zuerst vollständig und wörtlich aus** "
                "(Markdown, alle ↗-Links müssen klickbar bleiben – NICHT umschreiben, NICHT kürzen). "
                f"{_summary_task}\n",
            ]

            # Datum-Label Wochentag
            def _fmt_date(d: _date_t) -> str:
                return f"{_DE_DAYS[d.weekday()]}, {d.day}. {_DE_MONTHS[d.month]} {d.year}"

            if not is_week:
                output.append(f"## 📅 {_fmt_date(start_date)} — Kalenderübersicht\n")
            else:
                output.append(
                    f"## 📅 {label}: {_fmt_date(start_date)} – {_fmt_date(end_date - _td(1))}\n"
                )

            if not events:
                output.append("*Keine Termine gefunden.*")
                return "\n".join(output)

            output.append(f"*{len(events)} Termin(e) insgesamt*\n")
            output.append("---")

            # Termine nach Tag gruppieren (relevant bei Wochenansicht)
            from collections import defaultdict as _dd
            by_day: dict = _dd(list)
            for ev in events:
                _start_raw = ev.get("start", {}).get("dateTime", "")
                _day_key   = _start_raw[:10] if _start_raw else str(start_date)
                by_day[_day_key].append(ev)

            for _day_str in sorted(by_day.keys()):
                _day_events = by_day[_day_str]
                if is_week:
                    _d = _dt.strptime(_day_str, "%Y-%m-%d").date()
                    output.append(f"\n### {_fmt_date(_d)}")

                for ev in _day_events:
                    _subj      = ev.get("subject") or "(kein Titel)"
                    _is_allday = ev.get("isAllDay", False)
                    _cancelled = ev.get("isCancelled", False)
                    _show_as   = ev.get("showAs", "")      # free / busy / tentative / oof
                    _is_org    = ev.get("isOrganizer", False)
                    _web       = ev.get("webLink", "")

                    # Zeitformatierung
                    _s_raw = ev.get("start", {}).get("dateTime", "")
                    _e_raw = ev.get("end",   {}).get("dateTime", "")
                    if _is_allday:
                        _time_str = "Ganztag"
                    else:
                        _s_t = _s_raw[11:16] if len(_s_raw) >= 16 else "?"
                        _e_t = _e_raw[11:16] if len(_e_raw) >= 16 else "?"
                        _time_str = f"{_s_t} – {_e_t}"

                    # Titel (plain) + kleines ↗-Icon zum Öffnen in Outlook
                    _open_icon = f" [↗]({_web})" if _web else ""
                    _cancel_fmt = f"~~{_subj}~~" if _cancelled else _subj
                    _free_fmt   = f"*(Frei)* {_cancel_fmt}" if _show_as == "free" else _cancel_fmt

                    output.append(f"\n#### 🕐 {_time_str}  {_free_fmt}{_open_icon}")

                    # Fokus-Zeit erkennen
                    _is_focus = "focus" in _subj.lower() or "fokus" in _subj.lower()
                    if _is_focus:
                        output.append("- 🔕 *Fokus-/Sperrzeit – keine Meeting-Einladungen*")

                    # Organisator
                    _org = ev.get("organizer", {}).get("emailAddress", {})
                    if _org.get("name") and not _is_org:
                        output.append(f"- 👤 Organisiert von: **{_org['name']}**")
                    elif _is_org:
                        output.append("- 👤 Du bist Organisator")

                    # Ort / Online
                    _loc = (ev.get("location") or {}).get("displayName", "").strip()
                    _om  = ev.get("onlineMeetingUrl") or ""
                    if _loc:
                        output.append(f"- 📍 {_loc}")
                    if _om:
                        output.append(f"- 💻 [Teams-Beitritt]({_om})")

                    # Teilnehmer (max. 5 anzeigen)
                    _attendees = ev.get("attendees") or []
                    _att_names = [
                        a.get("emailAddress", {}).get("name", "")
                        for a in _attendees
                        if a.get("emailAddress", {}).get("name")
                           and a.get("type") != "resource"
                    ]
                    if _att_names:
                        _shown = _att_names[:5]
                        _rest  = len(_att_names) - len(_shown)
                        _att_str = ", ".join(_shown)
                        if _rest > 0:
                            _att_str += f" (+{_rest} weitere)"
                        output.append(f"- 👥 {_att_str}")

            output.append("\n---")
            return "\n".join(output)

        except httpx.HTTPStatusError as e:
            return f"❌ Graph API Fehler: HTTP {e.response.status_code}\n{e.response.text[:300]}"
        except Exception as e:
            return f"❌ Verbindungsfehler: {str(e)[:200]}"
