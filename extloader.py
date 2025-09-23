import fitz  # PyMuPDF
import base64
import logging
import os
import tempfile
import time
import json
import re
from typing import Dict, List, Tuple, Any, Optional
from fastapi import FastAPI, Request, HTTPException, status
from fastapi.responses import JSONResponse

# Fast native document processors
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

# Image processing
try:
    from PIL import Image
    import io
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

# Azure Document Intelligence compatible output formats
AZURE_DOC_INTEL_COMPATIBLE = True
DEFAULT_OUTPUT_FORMAT = "json"  # or "markdown"

# Image processing efficiency settings
AUTO_COMPRESS_IMAGES = os.getenv("AUTO_COMPRESS_IMAGES", "true").lower() == "true"
DEFAULT_COMPRESSION_WIDTH = int(os.getenv("FILE_IMAGE_COMPRESSION_WIDTH", "1024"))
DEFAULT_COMPRESSION_HEIGHT = int(os.getenv("FILE_IMAGE_COMPRESSION_HEIGHT", "1024"))
COMPRESSION_QUALITY = int(os.getenv("IMAGE_COMPRESSION_QUALITY", "85"))

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize the FastAPI application
app = FastAPI(
    title="External Content Processing Engine",
    description="An API to extract text and images from documents for OpenWebUI.",
    version="1.0.0"
)

def extract_docx_structure_azure_format(doc, filename: str, output_format: str = "json") -> str:
    """Extract structured content from DOCX in Azure Document Intelligence format."""
    
    # Build content string for span calculation
    full_content = ""
    content_elements = []
    
    # Process document elements in reading order
    element_id = 0
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            start_offset = len(full_content)
            text_content = paragraph.text.strip()
            
            # Determine element type
            style_name = paragraph.style.name if paragraph.style else ""
            if "Heading" in style_name:
                level = 1
                if "Heading 2" in style_name:
                    level = 2
                elif "Heading 3" in style_name:
                    level = 3
                elif "Heading 4" in style_name:
                    level = 4
                elif "Heading 5" in style_name:
                    level = 5
                elif "Heading 6" in style_name:
                    level = 6
                
                element_type = "title"
                role = f"heading{level}"
            else:
                element_type = "paragraph"
                role = "paragraph"
            
            content_elements.append({
                "id": f"element_{element_id}",
                "kind": element_type,
                "role": role,
                "content": text_content,
                "boundingRegions": [{"pageNumber": 1}],  # Simplified for DOCX
                "spans": [{"offset": start_offset, "length": len(text_content)}]
            })
            
            full_content += text_content + "\n\n"
            element_id += 1
    
    # Process tables
    table_id = 0
    for table in doc.tables:
        start_offset = len(full_content)
        table_content = ""
        
        # Extract table data
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            table_data.append(row_data)
        
        # Create table content for spans
        for row in table_data:
            table_content += " | ".join(row) + "\n"
        
        content_elements.append({
            "id": f"table_{table_id}",
            "kind": "table",
            "rowCount": len(table_data),
            "columnCount": len(table_data[0]) if table_data else 0,
            "cells": [
                {
                    "kind": "content",
                    "rowIndex": row_idx,
                    "columnIndex": col_idx,
                    "content": cell_content,
                    "spans": [{"offset": start_offset + table_content.find(cell_content), "length": len(cell_content)}]
                }
                for row_idx, row in enumerate(table_data)
                for col_idx, cell_content in enumerate(row)
                if cell_content.strip()
            ],
            "boundingRegions": [{"pageNumber": 1}],
            "spans": [{"offset": start_offset, "length": len(table_content)}]
        })
        
        full_content += table_content + "\n\n"
        table_id += 1
    
    # Create Azure Document Intelligence compatible response
    if output_format.lower() == "markdown":
        # Return markdown format like Azure
        markdown_content = ""
        for element in content_elements:
            if element["kind"] == "title":
                level = int(element["role"].replace("heading", ""))
                markdown_content += "#" * level + " " + element["content"] + "\n\n"
            elif element["kind"] == "paragraph":
                markdown_content += element["content"] + "\n\n"
            elif element["kind"] == "table":
                # Convert to markdown table
                table_rows = []
                max_cols = element.get("columnCount", 0)
                current_row = [""] * max_cols
                
                for cell in element["cells"]:
                    row_idx = cell["rowIndex"]
                    col_idx = cell["columnIndex"]
                    if row_idx >= len(table_rows):
                        table_rows.extend([[""] * max_cols for _ in range(row_idx - len(table_rows) + 1)])
                    table_rows[row_idx][col_idx] = cell["content"]
                
                if table_rows:
                    # Add table header
                    markdown_content += "| " + " | ".join(table_rows[0]) + " |\n"
                    markdown_content += "|" + "---|" * len(table_rows[0]) + "\n"
                    
                    # Add table rows
                    for row in table_rows[1:]:
                        markdown_content += "| " + " | ".join(row) + " |\n"
                    markdown_content += "\n"
        
        return markdown_content.strip()
    
    else:
        # Return Azure Document Intelligence JSON format
        azure_response = {
            "apiVersion": "2024-11-30",
            "modelId": "prebuilt-layout",
            "stringIndexType": "textElements",
            "content": full_content.strip(),
            "pages": [
                {
                    "pageNumber": 1,
                    "angle": 0,
                    "width": 8.5,
                    "height": 11,
                    "unit": "inch",
                    "spans": [{"offset": 0, "length": len(full_content.strip())}]
                }
            ],
            "paragraphs": [
                elem for elem in content_elements 
                if elem["kind"] in ["paragraph", "title"]
            ],
            "tables": [
                elem for elem in content_elements 
                if elem["kind"] == "table"
            ],
            "styles": [],
            "documents": [
                {
                    "docType": f"document:{filename}",
                    "boundingRegions": [{"pageNumber": 1}],
                    "fields": {},
                    "spans": [{"offset": 0, "length": len(full_content.strip())}],
                    "confidence": 0.99
                }
            ]
        }
        
        return json.dumps(azure_response, indent=2)

def process_docx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process DOCX files using python-docx with Azure Document Intelligence compatibility."""
    if not DOCX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-docx not available")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        doc = DocxDocument(temp_file_path)
        
        # Extract structured content in Azure Document Intelligence format
        structured_content = extract_docx_structure_azure_format(doc, filename, "json")
        
        os.unlink(temp_file_path)
        
        # Estimate page count (rough approximation)
        page_count = max(1, len(structured_content) // 3000)
        
        return structured_content, page_count
        
    except Exception as e:
        if 'temp_file_path' in locals():
            try:
                os.unlink(temp_file_path)
            except:
                pass
        raise HTTPException(status_code=500, detail=f"Failed to process DOCX: {e}")

def process_xlsx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process XLSX files using openpyxl."""
    if not OPENPYXL_AVAILABLE:
        raise HTTPException(status_code=500, detail="openpyxl not available")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        workbook = openpyxl.load_workbook(temp_file_path, read_only=True)
        full_text = ""
        sheet_count = 0
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_count += 1
            full_text += f"\n=== Sheet: {sheet_name} ===\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    full_text += "\t".join(row_text) + "\n"
        
        workbook.close()
        os.unlink(temp_file_path)
        
        return full_text.strip(), sheet_count
        
    except Exception as e:
        if 'temp_file_path' in locals():
            try:
                os.unlink(temp_file_path)
            except:
                pass
        raise HTTPException(status_code=500, detail=f"Failed to process XLSX: {e}")

def process_pptx(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process PPTX files using python-pptx."""
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-pptx not available")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        prs = Presentation(temp_file_path)
        full_text = ""
        slide_count = 0
        
        for slide in prs.slides:
            slide_count += 1
            full_text += f"\n=== Slide {slide_count} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + "\n"
                
                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            full_text += cell.text + " "
                        full_text += "\n"
        
        os.unlink(temp_file_path)
        
        return full_text.strip(), slide_count
        
    except Exception as e:
        if 'temp_file_path' in locals():
            try:
                os.unlink(temp_file_path)
            except:
                pass
        raise HTTPException(status_code=500, detail=f"Failed to process PPTX: {e}")

def process_xls_legacy(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process legacy XLS files using openpyxl (limited support)."""
    # For legacy XLS files, we'll attempt basic processing
    # Note: openpyxl has limited XLS support, this is a fallback
    try:
        return process_xlsx(file_bytes, filename)
    except Exception:
        # If openpyxl fails on XLS, return basic error info
        raise HTTPException(status_code=500, detail="Legacy XLS format not fully supported. Please convert to XLSX.")

def process_ppt_legacy(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process legacy PPT files - limited support."""
    # Legacy PPT format is not supported by python-pptx
    raise HTTPException(status_code=500, detail="Legacy PPT format not supported. Please convert to PPTX.")

def process_rtf(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """Process RTF files using striprtf."""
    if not RTF_AVAILABLE:
        raise HTTPException(status_code=500, detail="striprtf not available")
    
    try:
        rtf_content = file_bytes.decode('utf-8', errors='ignore')
        plain_text = rtf_to_text(rtf_content)
        
        # Estimate page count
        page_count = max(1, len(plain_text) // 3000)
        
        return plain_text.strip(), page_count
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process RTF: {e}")

def compress_image_for_efficiency(image_bytes: bytes, image_ext: str) -> tuple[bytes, bool]:
    """
    Compress image for processing efficiency if auto-compression is enabled.
    Returns: (processed_bytes, was_compressed)
    """
    if not PILLOW_AVAILABLE or not AUTO_COMPRESS_IMAGES:
        return image_bytes, False
    
    try:
        image = Image.open(io.BytesIO(image_bytes))
        
        # Convert to RGB if necessary (for JPEG compatibility)
        if image.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = background
        
        # Resize if image is larger than compression settings for efficiency
        original_size = (image.width, image.height)
        if image.width > DEFAULT_COMPRESSION_WIDTH or image.height > DEFAULT_COMPRESSION_HEIGHT:
            image.thumbnail((DEFAULT_COMPRESSION_WIDTH, DEFAULT_COMPRESSION_HEIGHT), Image.Resampling.LANCZOS)
        
        # Compress with quality setting for efficiency
        output = io.BytesIO()
        image.save(output, format='JPEG', quality=COMPRESSION_QUALITY, optimize=True)
        compressed_bytes = output.getvalue()
        
        # Only use compressed version if it's actually smaller
        if len(compressed_bytes) < len(image_bytes):
            original_mb = len(image_bytes) / (1024 * 1024)
            compressed_mb = len(compressed_bytes) / (1024 * 1024)
            logger.info(f"Image compressed for efficiency: {original_mb:.1f}MB -> {compressed_mb:.1f}MB")
            return compressed_bytes, True
        else:
            return image_bytes, False
            
    except Exception as e:
        logger.warning(f"Image compression failed, using original: {e}")
        return image_bytes, False

def process_images_efficiently(doc, extract_images_flag: str, filename: str) -> list:
    """
    Process PDF images with efficiency optimizations.
    Returns: base64_images list
    """
    base64_images = []
    images_processed = 0
    images_compressed = 0
    
    if extract_images_flag != 'true':
        logger.info("No image extraction requested (X-Extract-Images != 'true')")
        return base64_images
    
    logger.info("Starting efficient image extraction...")
    
    try:
        # Process images efficiently
        for page_num in range(len(doc)):
            image_list = doc.get_page_images(page_num, full=True)
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Apply efficiency compression if enabled
                    processed_bytes, was_compressed = compress_image_for_efficiency(image_bytes, image_ext)
                    
                    if was_compressed:
                        images_compressed += 1
                    
                    # Encode to base64
                    encoded_string = base64.b64encode(processed_bytes).decode("utf-8")
                    
                    # Determine the correct MIME type for data URI
                    if was_compressed:
                        data_uri = f"data:image/jpeg;base64,{encoded_string}"
                    else:
                        data_uri = f"data:image/{image_ext};base64,{encoded_string}"
                    
                    base64_images.append(data_uri)
                    images_processed += 1
                    
                except Exception as e:
                    logger.warning(f"Failed to process image {images_processed + 1}: {e}")
                    continue
        
        logger.info(f"Image processing complete: {images_processed} processed, {images_compressed} compressed for efficiency")
        
    except Exception as e:
        logger.error(f"Error during image processing: {e}")
    
    return base64_images

def process_non_pdf_fast(file_bytes: bytes, filename: str, file_ext: str) -> tuple[str, int]:
    """Process non-PDF documents using fast native Python processors."""
    logger.info(f"Processing {file_ext} file with fast native processor: {filename}")
    
    # Route to appropriate processor based on file extension
    if file_ext == ".docx":
        return process_docx(file_bytes, filename)
    elif file_ext == ".xlsx":
        return process_xlsx(file_bytes, filename)
    elif file_ext == ".pptx":
        return process_pptx(file_bytes, filename)
    elif file_ext == ".xls":
        return process_xls_legacy(file_bytes, filename)
    elif file_ext == ".ppt":
        return process_ppt_legacy(file_bytes, filename)
    elif file_ext == ".rtf":
        return process_rtf(file_bytes, filename)
    else:
        # Fallback for unsupported file types
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: {file_ext}. Supported types: .docx, .xlsx, .pptx, .rtf"
        )

@app.put("/process")
async def process_document(request: Request):
    """
    Processes an uploaded document to extract text and, optionally, images.
    """
    start_time = time.time()
    
    logger.info("Received request to /process endpoint.")
    logger.info("=== DOCUMENT PROCESSING REQUEST RECEIVED ===")
    logger.info(f"All headers: {dict(request.headers)}")
    
    file_bytes = await request.body()
    if not file_bytes:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No file content received."
        )

    filename = request.headers.get("x-filename", "unknown_file")
    file_ext = os.path.splitext(filename)[1].lower()
    
    # Extract image flag - EXACTLY like the original working version
    extract_images_flag = request.headers.get("x-extract-images", "false").lower()
    
    # Check for Azure Document Intelligence compatible output format
    output_format = request.headers.get("outputContentFormat", DEFAULT_OUTPUT_FORMAT).lower()
    
    # Debug logging
    logger.info(f"=== DEBUG: filename='{filename}', file_ext='{file_ext}'")
    logger.info(f"Raw X-Extract-Images header: {request.headers.get('x-extract-images')}")
    logger.info(f"Processed extract_images_flag: {extract_images_flag}")
    logger.info(f"Output format: {output_format}")
    logger.info(f"Will extract images: {extract_images_flag == 'true'}")
    logger.info(f"Will take PDF path: {file_ext == '.pdf'}")
    logger.info(f"File size: {len(file_bytes)} bytes")

    base64_images = []
    full_text = ""
    page_count = 0

    # PDF processing - use original working logic with efficiency improvements
    if file_ext == ".pdf":
        logger.info("Taking PDF processing path with PyMuPDF")
        try:
            # Open the PDF from the in-memory byte stream - EXACTLY like original
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            page_count = doc.page_count
            
            # 1. Extract all text content from the document - EXACTLY like original
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                full_text += page.get_text() + "\n\n"

            # 2. Extract images with efficiency optimizations
            base64_images = process_images_efficiently(doc, extract_images_flag, filename)

            doc.close()  # Close the document

        except Exception as e:
            logger.error(f"Failed to process PDF '{filename}': {e}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"An error occurred while processing the PDF: {e}",
            )
            
    # Non-PDF processing - use Azure-compatible processing for DOCX
    else:
        try:
            if file_ext == ".docx":
                # Use Azure-compatible processing for DOCX
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                    temp_file.write(file_bytes)
                    temp_file_path = temp_file.name
                
                doc = DocxDocument(temp_file_path)
                full_text = extract_docx_structure_azure_format(doc, filename, output_format)
                page_count = max(1, len(full_text) // 3000)
                
                os.unlink(temp_file_path)
            else:
                # Use existing fast native processors for other formats
                full_text, page_count = process_non_pdf_fast(file_bytes, filename, file_ext)
                
            logger.info(f"Successfully processed {file_ext} document. Pages/Sheets: {page_count}")
            
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Error processing {file_ext} document: {e}")
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to process document: {e}"
            )

    # Prepare the final JSON response payload with processing status
    response_payload = {
        "page_content": full_text.strip(),
        "metadata": {
            "source": filename,
            "page_count": page_count,
            "images_extracted": len(base64_images),
            "processing_status": "completed",
            "processing_time_ms": int((time.time() - start_time) * 1000) if 'start_time' in locals() else None,
            "output_format": output_format,
            "azure_compatible": AZURE_DOC_INTEL_COMPATIBLE
        },
        "images": base64_images,
    }

    logger.info(f"Returning response with {len(full_text)} chars of text and {len(base64_images)} images")
    
    return JSONResponse(content=response_payload)

@app.get("/")
def read_root():
    """A simple root endpoint to confirm the server is running."""
    processors_available = []
    if DOCX_AVAILABLE:
        processors_available.append("DOCX")
    if OPENPYXL_AVAILABLE:
        processors_available.append("XLSX/XLS")
    if PPTX_AVAILABLE:
        processors_available.append("PPTX")
    if RTF_AVAILABLE:
        processors_available.append("RTF")
    
    return {
        "status": "ok", 
        "message": "Content Processing Engine is running.",
        "pdf_processing": "PyMuPDF",
        "document_processing": "Azure Document Intelligence Compatible",
        "processors_available": processors_available,
        "supported_formats": ["PDF", "DOCX", "XLSX", "PPTX", "RTF"] if all([DOCX_AVAILABLE, OPENPYXL_AVAILABLE, PPTX_AVAILABLE, RTF_AVAILABLE]) else ["PDF"] + processors_available,
        "azure_compatible": AZURE_DOC_INTEL_COMPATIBLE,
        "output_formats": ["json", "markdown"],
        "image_processing": {
            "auto_compress": AUTO_COMPRESS_IMAGES,
            "compression_width": DEFAULT_COMPRESSION_WIDTH,
            "compression_height": DEFAULT_COMPRESSION_HEIGHT,
            "compression_quality": COMPRESSION_QUALITY,
            "compression_available": PILLOW_AVAILABLE
        }
    }

@app.get("/health")
def health_check():
    """Health check endpoint."""
    return {
        "status": "healthy",
        "pdf_processor": "PyMuPDF - Available",
        "docx_processor": "python-docx - Available (Azure Compatible)" if DOCX_AVAILABLE else "python-docx - Not Available",
        "xlsx_processor": "openpyxl - Available" if OPENPYXL_AVAILABLE else "openpyxl - Not Available", 
        "pptx_processor": "python-pptx - Available" if PPTX_AVAILABLE else "python-pptx - Not Available",
        "rtf_processor": "striprtf - Available" if RTF_AVAILABLE else "striprtf - Not Available",
        "image_compression": "Pillow - Available" if PILLOW_AVAILABLE else "Pillow - Not Available",
        "azure_compatibility": "Enabled",
        "compression_settings": {
            "auto_compress": AUTO_COMPRESS_IMAGES,
            "max_width": DEFAULT_COMPRESSION_WIDTH,
            "max_height": DEFAULT_COMPRESSION_HEIGHT,
            "quality": COMPRESSION_QUALITY
        }
    }

@app.get("/test-image-extraction")
def test_image_extraction():
    """Test endpoint to verify image extraction capability."""
    return {
        "status": "ready",
        "message": "Send a PDF with 'X-Extract-Images: true' header to /process to test image extraction",
        "pdf_processing": "PyMuPDF",
        "azure_compatibility": "Azure Document Intelligence Compatible Output",
        "expected_headers": {
            "X-Filename": "your-file.pdf",
            "X-Extract-Images": "true",
            "outputContentFormat": "json or markdown (optional)"
        }
    }